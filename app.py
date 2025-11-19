from flask import Flask, render_template, request, jsonify, send_file, session, redirect, url_for
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
from datetime import datetime, date, time, timedelta
import os
import json
import logging
from openai import OpenAI
import markdown

# Configura logging per vedere output con Gunicorn
# IMPORTANTE: Configura il logging PRIMA di creare l'app per vedere i log con Gunicorn
import sys
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(sys.stdout),  # Forza output su stdout per Gunicorn
        logging.StreamHandler(sys.stderr)   # Anche su stderr
    ]
)
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from io import BytesIO
import requests
import base64

app = Flask(__name__)
# Database in directory instance per persistenza con Docker
instance_path = os.path.join(os.getcwd(), 'instance')
os.makedirs(instance_path, exist_ok=True)
app.config['SQLALCHEMY_DATABASE_URI'] = f'sqlite:///{os.path.join(instance_path, "courses.db")}'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['COURSES_DIR'] = 'courses'
app.config['MD_SOURCE_DIR'] = 'MD'
app.config['OPENAI_API_KEY'] = os.environ.get('OPENAI_API_KEY', '')
app.config['CANVA_CLIENT_ID'] = os.environ.get('CANVA_CLIENT_ID', '')
app.config['CANVA_CLIENT_SECRET'] = os.environ.get('CANVA_CLIENT_SECRET', '')
app.config['CANVA_REDIRECT_URI'] = os.environ.get('CANVA_REDIRECT_URI', 'http://localhost:5001/api/canva/callback')
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'dev-secret-key-change-in-production')

# Abilita sessioni per OAuth
app.secret_key = app.config['SECRET_KEY']

db = SQLAlchemy(app)
CORS(app)

# Modelli Database
class Course(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    code = db.Column(db.String(50), unique=True, nullable=False)
    name = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text)
    total_hours = db.Column(db.Integer, nullable=False)
    theory_hours = db.Column(db.Integer, nullable=False)
    practice_hours = db.Column(db.Integer, nullable=False)
    num_lessons = db.Column(db.Integer, default=0)  # Numero di lezioni previste
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    export_count = db.Column(db.Integer, default=0)  # Contatore esportazioni
    lessons = db.relationship('Lesson', backref='course', lazy=True, cascade='all, delete-orphan')

class Lesson(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)
    title = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text)
    lesson_type = db.Column(db.String(20), nullable=False)  # 'theory' o 'practice'
    duration_hours = db.Column(db.Float, nullable=False)
    order = db.Column(db.Integer, nullable=False)
    content = db.Column(db.Text)  # Contenuto markdown
    objectives = db.Column(db.Text)  # JSON array di obiettivi
    materials = db.Column(db.Text)  # JSON array di materiali
    exercises = db.Column(db.Text)  # JSON array di esercizi
    lesson_date = db.Column(db.Date, nullable=True)  # Data della lezione
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    modification_count = db.Column(db.Integer, default=0)  # Contatore modifiche

class TrainingMaterial(db.Model):
    """Materiali di addestramento per l'AI (documenti, URL)"""
    id = db.Column(db.Integer, primary_key=True)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)
    material_type = db.Column(db.String(20), nullable=False)  # 'file' o 'url'
    filename = db.Column(db.String(500))  # Nome file se tipo 'file'
    url = db.Column(db.String(1000))  # URL se tipo 'url'
    content_extracted = db.Column(db.Text)  # Contenuto estratto dal documento/URL
    summary = db.Column(db.Text)  # Riassunto generato da ChatGPT
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class Question(db.Model):
    """Domande a risposta multipla per i corsi"""
    id = db.Column(db.Integer, primary_key=True)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)
    question_number = db.Column(db.Integer, nullable=False)  # Numero della domanda (1-30)
    question_text = db.Column(db.Text, nullable=False)  # Testo della domanda
    option_a = db.Column(db.Text, nullable=False)
    option_b = db.Column(db.Text, nullable=False)
    option_c = db.Column(db.Text, nullable=False)
    option_d = db.Column(db.Text, nullable=False)
    correct_answer = db.Column(db.String(1), nullable=False)  # 'A', 'B', 'C', o 'D'
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class Preference(db.Model):
    """Preferenze dell'applicazione"""
    id = db.Column(db.Integer, primary_key=True)
    key = db.Column(db.String(100), unique=True, nullable=False)
    value = db.Column(db.Text)  # Valore come stringa (JSON per valori complessi)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

class CourseTemplate(db.Model):
    """Template di corsi riutilizzabili"""
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text)
    code_prefix = db.Column(db.String(50))  # Prefisso per il codice corso (es: "TEMPLATE_")
    course_data = db.Column(db.Text, nullable=False)  # JSON con dati del corso
    lessons_data = db.Column(db.Text)  # JSON con dati delle lezioni
    is_predefined = db.Column(db.Boolean, default=False)  # True per template predefiniti
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

class LessonVersion(db.Model):
    """Versioni storiche delle lezioni per il versioning"""
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    version_number = db.Column(db.Integer, nullable=False)  # Numero versione incrementale
    title = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text)
    lesson_type = db.Column(db.String(20), nullable=False)
    duration_hours = db.Column(db.Float, nullable=False)
    order = db.Column(db.Integer, nullable=False)
    content = db.Column(db.Text)
    objectives = db.Column(db.Text)  # JSON array
    materials = db.Column(db.Text)  # JSON array
    exercises = db.Column(db.Text)  # JSON array
    lesson_date = db.Column(db.Date, nullable=True)
    comment = db.Column(db.Text)  # Commento opzionale per questa versione
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    created_by = db.Column(db.String(100))  # Utente che ha creato la versione (opzionale)
    
    lesson = db.relationship('Lesson', backref='versions')

class LessonNote(db.Model):
    """Note private per lezioni"""
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    content = db.Column(db.Text, nullable=False)
    is_private = db.Column(db.Boolean, default=True)  # True per note private
    created_by = db.Column(db.String(100))  # Utente che ha creato la nota
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    
    lesson = db.relationship('Lesson', backref='notes')

class LessonComment(db.Model):
    """Commenti collaborativi per lezioni"""
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    content = db.Column(db.Text, nullable=False)
    author = db.Column(db.String(100), nullable=False)  # Nome autore commento
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    parent_id = db.Column(db.Integer, db.ForeignKey('lesson_comment.id'), nullable=True)  # Per commenti annidati
    
    lesson = db.relationship('Lesson', backref='comments')
    parent = db.relationship('LessonComment', remote_side=[id], backref='replies')

class LessonReminder(db.Model):
    """Promemoria per revisioni delle lezioni"""
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=False)
    title = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text)
    reminder_date = db.Column(db.DateTime, nullable=False)  # Data/ora promemoria
    is_completed = db.Column(db.Boolean, default=False)
    created_by = db.Column(db.String(100))
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    
    lesson = db.relationship('Lesson', backref='reminders')

class Notification(db.Model):
    """Sistema di notifiche"""
    id = db.Column(db.Integer, primary_key=True)
    lesson_id = db.Column(db.Integer, db.ForeignKey('lesson.id'), nullable=True)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=True)
    type = db.Column(db.String(50), nullable=False)  # 'comment', 'reminder', 'note', 'version', etc.
    title = db.Column(db.String(200), nullable=False)
    message = db.Column(db.Text, nullable=False)
    is_read = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    
    lesson = db.relationship('Lesson', backref='notifications')
    course = db.relationship('Course', backref='notifications')

class AIAnalysis(db.Model):
    """Analisi AI salvate per i corsi"""
    id = db.Column(db.Integer, primary_key=True)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)
    analysis_type = db.Column(db.String(50), nullable=False)  # 'balance', 'suggestions', 'duplicates', 'coherence'
    title = db.Column(db.String(200), nullable=False)
    analysis_content = db.Column(db.Text, nullable=False)  # Contenuto dell'analisi in Markdown
    analysis_data = db.Column(db.Text)  # JSON con dati aggiuntivi (es. percentuali)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    
    course = db.relationship('Course', backref='ai_analyses')

# Creare directory per i corsi
os.makedirs(app.config['COURSES_DIR'], exist_ok=True)
os.makedirs(app.config['MD_SOURCE_DIR'], exist_ok=True)

# Funzione di inizializzazione database (funziona anche con Gunicorn)
def init_database():
    """Inizializza il database e applica migrazioni"""
    import sys
    try:
        # Assicurati che la directory instance esista
        instance_path = os.path.join(os.getcwd(), 'instance')
        os.makedirs(instance_path, exist_ok=True)
        
        db_path = os.path.join(instance_path, 'courses.db')
        logger.info(f"🔧 Inizializzazione database: {db_path}")
        logger.info(f"   Directory esiste: {os.path.exists(instance_path)}")
        logger.info(f"   Database esiste: {os.path.exists(db_path)}")
        
        # FORZA l'inizializzazione anche se siamo già in un app_context
        with app.app_context():
            # Crea tutte le tabelle - FORZA la creazione
            logger.info("   Creazione tabelle...")
            try:
                db.create_all()
                logger.info("   ✓ Tabelle create")
            except Exception as create_error:
                logger.error(f"   ❌ Errore in db.create_all(): {create_error}")
                # Prova a creare le tabelle manualmente
                try:
                    from sqlalchemy import text
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS course (
                            id INTEGER PRIMARY KEY,
                            code VARCHAR(50) UNIQUE NOT NULL,
                            name VARCHAR(200) NOT NULL,
                            description TEXT,
                            total_hours INTEGER NOT NULL,
                            theory_hours INTEGER NOT NULL,
                            practice_hours INTEGER NOT NULL,
                            num_lessons INTEGER DEFAULT 0,
                            created_at DATETIME
                        )
                    """))
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS lesson (
                            id INTEGER PRIMARY KEY,
                            course_id INTEGER NOT NULL,
                            title VARCHAR(200) NOT NULL,
                            description TEXT,
                            lesson_type VARCHAR(20),
                            duration_hours REAL,
                            "order" INTEGER,
                            content TEXT,
                            objectives TEXT,
                            materials TEXT,
                            exercises TEXT,
                            lesson_date DATE,
                            created_at DATETIME,
                            FOREIGN KEY (course_id) REFERENCES course(id)
                        )
                    """))
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS preference (
                            id INTEGER PRIMARY KEY,
                            key VARCHAR(100) UNIQUE NOT NULL,
                            value TEXT,
                            updated_at DATETIME
                        )
                    """))
                    db.session.commit()
                    logger.info("   ✓ Tabelle create manualmente")
                except Exception as manual_error:
                    logger.error(f"   ❌ Errore creazione manuale: {manual_error}")
                    raise
            
            # Verifica che le tabelle siano state create
            from sqlalchemy import inspect, text
            inspector = inspect(db.engine)
            tables = inspector.get_table_names()
            logger.info(f"   Tabelle esistenti: {tables}")
            
            # Migrazioni: aggiungi colonne se non esistono
            if 'course' in tables:
                try:
                    columns = [col['name'] for col in inspector.get_columns('course')]
                    if 'num_lessons' not in columns:
                        db.session.execute(text('ALTER TABLE course ADD COLUMN num_lessons INTEGER DEFAULT 0'))
                        db.session.commit()
                        logger.info("   ✓ Colonna num_lessons aggiunta")
                    if 'export_count' not in columns:
                        db.session.execute(text('ALTER TABLE course ADD COLUMN export_count INTEGER DEFAULT 0'))
                        db.session.commit()
                        logger.info("   ✓ Colonna export_count aggiunta")
                except Exception as e:
                    logger.info(f"   Info course migrations: {e}")
            
            if 'lesson' in tables:
                try:
                    lesson_columns = [col['name'] for col in inspector.get_columns('lesson')]
                    if 'lesson_date' not in lesson_columns:
                        db.session.execute(text('ALTER TABLE lesson ADD COLUMN lesson_date DATE'))
                        db.session.commit()
                        logger.info("   ✓ Colonna lesson_date aggiunta")
                    if 'updated_at' not in lesson_columns:
                        db.session.execute(text('ALTER TABLE lesson ADD COLUMN updated_at DATETIME'))
                        db.session.commit()
                        logger.info("   ✓ Colonna updated_at aggiunta")
                    if 'modification_count' not in lesson_columns:
                        db.session.execute(text('ALTER TABLE lesson ADD COLUMN modification_count INTEGER DEFAULT 0'))
                        db.session.commit()
                        logger.info("   ✓ Colonna modification_count aggiunta")
                except Exception as e:
                    logger.info(f"   Info lesson migrations: {e}")
            
            if 'preference' not in tables:
                try:
                    db.create_all()
                    logger.info("   ✓ Tabella Preference creata")
                except Exception as e:
                    logger.info(f"   Info Preference: {e}")
            
            # Migrazione per course_template
            if 'course_template' not in tables:
                try:
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS course_template (
                            id INTEGER PRIMARY KEY,
                            name VARCHAR(200) NOT NULL,
                            description TEXT,
                            code_prefix VARCHAR(50),
                            course_data TEXT NOT NULL,
                            lessons_data TEXT,
                            is_predefined BOOLEAN DEFAULT 0,
                            created_at DATETIME,
                            updated_at DATETIME
                        )
                    """))
                    db.session.commit()
                    logger.info("   ✓ Tabella course_template creata")
                except Exception as e:
                    logger.info(f"   Info course_template creation: {e}")
            
            # Migrazione per lesson_version
            if 'lesson_version' not in tables:
                try:
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS lesson_version (
                            id INTEGER PRIMARY KEY,
                            lesson_id INTEGER NOT NULL,
                            version_number INTEGER NOT NULL,
                            title VARCHAR(200) NOT NULL,
                            description TEXT,
                            lesson_type VARCHAR(20) NOT NULL,
                            duration_hours REAL NOT NULL,
                            "order" INTEGER NOT NULL,
                            content TEXT,
                            objectives TEXT,
                            materials TEXT,
                            exercises TEXT,
                            lesson_date DATE,
                            comment TEXT,
                            created_at DATETIME,
                            created_by VARCHAR(100),
                            FOREIGN KEY (lesson_id) REFERENCES lesson(id)
                        )
                    """))
                    db.session.commit()
                    logger.info("   ✓ Tabella lesson_version creata")
                except Exception as e:
                    logger.info(f"   Info lesson_version creation: {e}")
            
            # Migrazioni per tabelle collaborazione
            if 'lesson_note' not in tables:
                try:
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS lesson_note (
                            id INTEGER PRIMARY KEY,
                            lesson_id INTEGER NOT NULL,
                            content TEXT NOT NULL,
                            is_private BOOLEAN DEFAULT 1,
                            created_by VARCHAR(100),
                            created_at DATETIME,
                            updated_at DATETIME,
                            FOREIGN KEY (lesson_id) REFERENCES lesson(id)
                        )
                    """))
                    db.session.commit()
                    logger.info("   ✓ Tabella lesson_note creata")
                except Exception as e:
                    logger.info(f"   Info lesson_note creation: {e}")
            
            if 'lesson_comment' not in tables:
                try:
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS lesson_comment (
                            id INTEGER PRIMARY KEY,
                            lesson_id INTEGER NOT NULL,
                            content TEXT NOT NULL,
                            author VARCHAR(100) NOT NULL,
                            created_at DATETIME,
                            updated_at DATETIME,
                            parent_id INTEGER,
                            FOREIGN KEY (lesson_id) REFERENCES lesson(id),
                            FOREIGN KEY (parent_id) REFERENCES lesson_comment(id)
                        )
                    """))
                    db.session.commit()
                    logger.info("   ✓ Tabella lesson_comment creata")
                except Exception as e:
                    logger.info(f"   Info lesson_comment creation: {e}")
            
            if 'lesson_reminder' not in tables:
                try:
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS lesson_reminder (
                            id INTEGER PRIMARY KEY,
                            lesson_id INTEGER NOT NULL,
                            title VARCHAR(200) NOT NULL,
                            description TEXT,
                            reminder_date DATETIME NOT NULL,
                            is_completed BOOLEAN DEFAULT 0,
                            created_by VARCHAR(100),
                            created_at DATETIME,
                            FOREIGN KEY (lesson_id) REFERENCES lesson(id)
                        )
                    """))
                    db.session.commit()
                    logger.info("   ✓ Tabella lesson_reminder creata")
                except Exception as e:
                    logger.info(f"   Info lesson_reminder creation: {e}")
            
            if 'notification' not in tables:
                try:
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS notification (
                            id INTEGER PRIMARY KEY,
                            lesson_id INTEGER,
                            course_id INTEGER,
                            type VARCHAR(50) NOT NULL,
                            title VARCHAR(200) NOT NULL,
                            message TEXT NOT NULL,
                            is_read BOOLEAN DEFAULT 0,
                            created_at DATETIME,
                            FOREIGN KEY (lesson_id) REFERENCES lesson(id),
                            FOREIGN KEY (course_id) REFERENCES course(id)
                        )
                    """))
                    db.session.commit()
                    logger.info("   ✓ Tabella notification creata")
                except Exception as e:
                    logger.info(f"   Info notification creation: {e}")
            
            # Migrazione per ai_analysis
            if 'ai_analysis' not in tables:
                try:
                    db.session.execute(text("""
                        CREATE TABLE IF NOT EXISTS ai_analysis (
                            id INTEGER PRIMARY KEY,
                            course_id INTEGER NOT NULL,
                            analysis_type VARCHAR(50) NOT NULL,
                            title VARCHAR(200) NOT NULL,
                            analysis_content TEXT NOT NULL,
                            analysis_data TEXT,
                            created_at DATETIME,
                            FOREIGN KEY (course_id) REFERENCES course(id)
                        )
                    """))
                    db.session.commit()
                    logger.info("   ✓ Tabella ai_analysis creata")
                except Exception as e:
                    logger.info(f"   Info ai_analysis creation: {e}")
            
            # Test query per verificare che funzioni
            try:
                count = Course.query.count()
                logger.info(f"   ✓ Database funzionante (corsi: {count})")
            except Exception as e:
                logger.warning(f"   ⚠️  Errore test query: {e}")
                
            logger.info("✅ Database inizializzato correttamente")
    except Exception as e:
        error_msg = f"❌ ERRORE inizializzazione database: {e}"
        logger.error(error_msg)
        import traceback
        logger.error(traceback.format_exc())
        raise  # Rilancia l'eccezione per vedere l'errore

# Flag per tracciare se il database è stato inizializzato
_db_initialized = False

# Inizializza il database all'avvio (funziona anche con Gunicorn)
# Usa un try/except per non bloccare l'avvio se c'è un problema
logger.info("=" * 60)
logger.info("AVVIO APPLICAZIONE - Inizializzazione database")
logger.info("=" * 60)
try:
    init_database()
    _db_initialized = True
    logger.info("✅ Database inizializzato all'avvio")
except Exception as e:
    logger.warning(f"⚠️  Attenzione: errore durante inizializzazione database all'avvio: {e}")
    logger.warning("   L'applicazione continuerà, ma potrebbe non funzionare correttamente")
    logger.warning("   Il database verrà inizializzato alla prima richiesta")
    import traceback
    logger.warning(traceback.format_exc())
    _db_initialized = False
logger.info("=" * 60)

# Hook per inizializzare il database alla prima richiesta se necessario
@app.before_request
def ensure_db_initialized():
    """Assicura che il database sia inizializzato prima di ogni richiesta"""
    global _db_initialized
    if not _db_initialized:
        try:
            # Verifica se il database esiste provando una query
            try:
                from sqlalchemy import inspect
                inspector = inspect(db.engine)
                tables = inspector.get_table_names()
                if 'course' in tables:
                    _db_initialized = True
                    logger.info("Database già esistente, skip inizializzazione")
                    return
            except:
                pass
            
            # Se non esiste, inizializzalo
            logger.info("Database non trovato, inizializzazione in corso...")
            init_database()
            _db_initialized = True
            logger.info("Database inizializzato con successo in before_request")
        except Exception as e:
            logger.error(f"⚠️  Errore inizializzazione database in before_request: {e}")
            import traceback
            logger.error(traceback.format_exc())

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/health')
def health():
    """Health check endpoint - inizializza anche il database se necessario"""
    try:
        # Verifica che il database funzioni
        try:
            db.session.execute(db.text('SELECT 1'))
            db_status = 'ok'
        except Exception as db_err:
            # Se il database non esiste, prova a inizializzarlo
            error_str = str(db_err)
            if 'no such table' in error_str.lower() or 'operationalerror' in error_str.lower():
                logger.warning("Database non trovato in /health, inizializzazione in corso...")
                try:
                    init_database()
                    db_status = 'initialized'
                except Exception as init_err:
                    db_status = f'error: {str(init_err)}'
            else:
                db_status = f'error: {str(db_err)}'
    except Exception as e:
        db_status = f'error: {str(e)}'
    
    return jsonify({
        'status': 'ok',
        'message': 'Server is running',
        'database': db_status,
        'db_path': app.config['SQLALCHEMY_DATABASE_URI']
    }), 200

@app.route('/api/db-status', methods=['GET'])
def db_status():
    """Endpoint per verificare lo stato del database"""
    try:
        from sqlalchemy import inspect
        inspector = inspect(db.engine)
        tables = inspector.get_table_names()
        
        course_count = 0
        lesson_count = 0
        try:
            course_count = Course.query.count()
            lesson_count = Lesson.query.count()
        except:
            pass
        
        db_path = app.config['SQLALCHEMY_DATABASE_URI']
        db_file_exists = os.path.exists(os.path.join('instance', 'courses.db'))
        
        return jsonify({
            'status': 'ok',
            'database_exists': db_file_exists,
            'database_path': db_path,
            'tables': tables,
            'course_count': course_count,
            'lesson_count': lesson_count
        }), 200
    except Exception as e:
        import traceback
        return jsonify({
            'status': 'error',
            'error': str(e),
            'traceback': traceback.format_exc()
        }), 500

@app.route('/api/init-db', methods=['POST', 'GET'])
@app.route('/init-db', methods=['POST', 'GET'])  # Endpoint alternativo senza /api
def force_init_db():
    """Endpoint per forzare l'inizializzazione del database (utile per debug)"""
    try:
        logger.info("=== FORZATURA INIZIALIZZAZIONE DATABASE ===")
        init_database()
        logger.info("=== INIZIALIZZAZIONE COMPLETATA ===")
        
        # Verifica
        from sqlalchemy import inspect
        inspector = inspect(db.engine)
        tables = inspector.get_table_names()
        
        return jsonify({
            'status': 'success',
            'message': 'Database inizializzato con successo',
            'tables': tables,
            'tables_count': len(tables)
        }), 200
    except Exception as e:
        import traceback
        logger.error(f"Errore in force_init_db: {e}")
        logger.error(traceback.format_exc())
        return jsonify({
            'status': 'error',
            'error': str(e),
            'traceback': traceback.format_exc()
        }), 500

@app.route('/api/courses/search', methods=['GET'])
def search_courses():
    """Ricerca avanzata su corsi e lezioni"""
    try:
        from sqlalchemy.orm import joinedload
        from sqlalchemy import or_, and_
        
        # Parametri di ricerca
        query = request.args.get('q', '').strip()
        lesson_type = request.args.get('lesson_type', '')
        min_duration = request.args.get('min_duration', type=float)
        date_from = request.args.get('date_from', '')
        date_to = request.args.get('date_to', '')
        search_in_content = request.args.get('search_in_content', 'true').lower() == 'true'
        
        # Query base
        courses_query = Course.query.options(joinedload(Course.lessons))
        
        # Filtro per ricerca testuale
        if query:
            if search_in_content:
                # Cerca nei corsi e nelle lezioni (incluso contenuto Markdown)
                courses_query = courses_query.join(Lesson).filter(
                    or_(
                        Course.name.ilike(f'%{query}%'),
                        Course.code.ilike(f'%{query}%'),
                        Course.description.ilike(f'%{query}%'),
                        Lesson.title.ilike(f'%{query}%'),
                        Lesson.description.ilike(f'%{query}%'),
                        Lesson.content.ilike(f'%{query}%')
                    )
                ).distinct()
            else:
                # Cerca solo nei corsi
                courses_query = courses_query.filter(
                    or_(
                        Course.name.ilike(f'%{query}%'),
                        Course.code.ilike(f'%{query}%'),
                        Course.description.ilike(f'%{query}%')
                    )
                )
        
        # Filtro per tipo lezione
        if lesson_type:
            courses_query = courses_query.join(Lesson).filter(
                Lesson.lesson_type == lesson_type
            ).distinct()
        
        # Filtro per durata minima
        if min_duration is not None:
            courses_query = courses_query.join(Lesson).filter(
                Lesson.duration_hours >= min_duration
            ).distinct()
        
        # Filtro per data
        if date_from or date_to:
            courses_query = courses_query.join(Lesson)
            date_filters = []
            if date_from:
                try:
                    from_date = datetime.strptime(date_from, '%Y-%m-%d').date()
                    date_filters.append(Lesson.lesson_date >= from_date)
                except ValueError:
                    pass
            if date_to:
                try:
                    to_date = datetime.strptime(date_to, '%Y-%m-%d').date()
                    date_filters.append(Lesson.lesson_date <= to_date)
                except ValueError:
                    pass
            if date_filters:
                courses_query = courses_query.filter(and_(*date_filters)).distinct()
        
        courses = courses_query.all()
        
        # Serializza i risultati dei corsi
        result = []
        matched_lessons = []
        
        for c in courses:
            try:
                lessons_list = list(c.lessons)
                lessons_with_dates = [l for l in lessons_list if l.lesson_date is not None]
                start_date = None
                end_date = None
                if lessons_with_dates:
                    dates = [l.lesson_date for l in lessons_with_dates]
                    start_date = min(dates)
                    end_date = max(dates)
                
                # Trova lezioni che matchano i criteri di ricerca
                course_matched_lessons = []
                for lesson in lessons_list:
                    lesson_matches = False
                    match_reasons = []
                    
                    # Verifica se la lezione matcha i filtri
                    if lesson_type and lesson.lesson_type != lesson_type:
                        continue
                    if min_duration is not None and lesson.duration_hours < min_duration:
                        continue
                    if date_from or date_to:
                        if lesson.lesson_date:
                            if date_from:
                                try:
                                    from_date = datetime.strptime(date_from, '%Y-%m-%d').date()
                                    if lesson.lesson_date < from_date:
                                        continue
                                except ValueError:
                                    pass
                            if date_to:
                                try:
                                    to_date = datetime.strptime(date_to, '%Y-%m-%d').date()
                                    if lesson.lesson_date > to_date:
                                        continue
                                except ValueError:
                                    pass
                        else:
                            continue
                    
                    # Verifica se matcha la query testuale
                    if query:
                        query_lower = query.lower()
                        if lesson.title and query_lower in lesson.title.lower():
                            lesson_matches = True
                            match_reasons.append('titolo')
                        if lesson.description and query_lower in lesson.description.lower():
                            lesson_matches = True
                            match_reasons.append('descrizione')
                        if lesson.content and query_lower in lesson.content.lower():
                            lesson_matches = True
                            match_reasons.append('contenuto')
                    else:
                        # Se non c'è query ma matcha i filtri, includila
                        lesson_matches = True
                    
                    if lesson_matches:
                        # Estrai snippet del contenuto se matcha
                        content_snippet = None
                        if query and lesson.content:
                            content_lower = lesson.content.lower()
                            query_lower = query.lower()
                            idx = content_lower.find(query_lower)
                            if idx >= 0:
                                # Estrai 100 caratteri prima e dopo il match
                                start = max(0, idx - 100)
                                end = min(len(lesson.content), idx + len(query) + 100)
                                snippet = lesson.content[start:end]
                                if start > 0:
                                    snippet = '...' + snippet
                                if end < len(lesson.content):
                                    snippet = snippet + '...'
                                content_snippet = snippet
                        
                        course_matched_lessons.append({
                            'id': lesson.id,
                            'title': lesson.title,
                            'description': lesson.description,
                            'lesson_type': lesson.lesson_type,
                            'duration_hours': lesson.duration_hours,
                            'order': lesson.order,
                            'lesson_date': lesson.lesson_date.isoformat() if lesson.lesson_date else None,
                            'match_reasons': match_reasons,
                            'content_snippet': content_snippet
                        })
                
                result.append({
                    'id': c.id,
                    'code': c.code,
                    'name': c.name,
                    'description': c.description,
                    'total_hours': c.total_hours,
                    'theory_hours': c.theory_hours,
                    'practice_hours': c.practice_hours,
                    'num_lessons': c.num_lessons or 0,
                    'lessons_count': len(lessons_list),
                    'start_date': start_date.isoformat() if start_date else None,
                    'end_date': end_date.isoformat() if end_date else None,
                    'matched_lessons': course_matched_lessons
                })
                
                # Aggiungi lezioni matchate alla lista globale
                matched_lessons.extend([{
                    **lesson_data,
                    'course_id': c.id,
                    'course_name': c.name,
                    'course_code': c.code
                } for lesson_data in course_matched_lessons])
                
            except Exception as e:
                logger.error(f"Errore processando corso {c.id}: {e}")
                result.append({
                    'id': c.id,
                    'code': c.code,
                    'name': c.name,
                    'description': c.description or '',
                    'total_hours': c.total_hours or 0,
                    'theory_hours': c.theory_hours or 0,
                    'practice_hours': c.practice_hours or 0,
                    'num_lessons': getattr(c, 'num_lessons', 0) or 0,
                    'lessons_count': 0,
                    'start_date': None,
                    'end_date': None,
                    'matched_lessons': []
                })
        
        return jsonify({
            'results': result,
            'matched_lessons': matched_lessons,
            'count': len(result),
            'lessons_count': len(matched_lessons),
            'query': query,
            'filters': {
                'lesson_type': lesson_type,
                'min_duration': min_duration,
                'date_from': date_from,
                'date_to': date_to,
                'search_in_content': search_in_content
            }
        })
    except Exception as e:
        import traceback
        logger.error(f"Errore nella ricerca: {e}\n{traceback.format_exc()}")
        return jsonify({'error': f'Errore nella ricerca: {str(e)}'}), 500

@app.route('/api/courses', methods=['GET'])
def get_courses():
    try:
        from sqlalchemy.orm import joinedload
        
        # Prova a caricare i corsi - se fallisce, inizializza il database
        try:
            courses = Course.query.options(joinedload(Course.lessons)).all()
        except Exception as db_error:
            # Se il database non esiste, inizializzalo
            error_str = str(db_error)
            if 'no such table' in error_str.lower() or 'operationalerror' in error_str.lower():
                logger.warning(f"Database non inizializzato, inizializzazione in corso... Errore: {db_error}")
                try:
                    init_database()
                    logger.info("Database inizializzato con successo in get_courses")
                    # Riprova la query dopo l'inizializzazione
                    courses = Course.query.options(joinedload(Course.lessons)).all()
                except Exception as init_error:
                    logger.error(f"Errore durante inizializzazione database in get_courses: {init_error}")
                    import traceback
                    logger.error(traceback.format_exc())
                    return jsonify({'error': f'Database non inizializzato: {str(init_error)}'}), 500
            else:
                # Altro tipo di errore, rilancia
                raise
        result = []
        for c in courses:
            try:
                # Calcola la prima e ultima data delle lezioni
                # Forza il caricamento delle lezioni
                lessons_list = list(c.lessons)
                lessons_with_dates = [l for l in lessons_list if l.lesson_date is not None]
                start_date = None
                end_date = None
                if lessons_with_dates:
                    dates = [l.lesson_date for l in lessons_with_dates]
                    start_date = min(dates)
                    end_date = max(dates)
                
                result.append({
                    'id': c.id,
                    'code': c.code,
                    'name': c.name,
                    'description': c.description,
                    'total_hours': c.total_hours,
                    'theory_hours': c.theory_hours,
                    'practice_hours': c.practice_hours,
                    'num_lessons': c.num_lessons or 0,
                    'lessons_count': len(lessons_list),
                    'start_date': start_date.isoformat() if start_date else None,
                    'end_date': end_date.isoformat() if end_date else None
                })
            except Exception as e:
                # Se c'è un errore con un corso, logga ma continua
                print(f"Errore processando corso {c.id}: {e}")
                result.append({
                    'id': c.id,
                    'code': c.code,
                    'name': c.name,
                    'description': c.description or '',
                    'total_hours': c.total_hours or 0,
                    'theory_hours': c.theory_hours or 0,
                    'practice_hours': c.practice_hours or 0,
                    'num_lessons': getattr(c, 'num_lessons', 0) or 0,
                    'lessons_count': 0,
                    'start_date': None,
                    'end_date': None
                })
        return jsonify(result)
    except Exception as e:
        import traceback
        error_msg = f"Errore nel caricamento dei corsi: {str(e)}"
        print(error_msg)
        traceback.print_exc()
        return jsonify({'error': error_msg}), 500

@app.route('/api/md-files', methods=['GET'])
def list_md_files():
    """Lista i file MD disponibili nella cartella MD"""
    md_dir = app.config['MD_SOURCE_DIR']
    files = []
    if os.path.exists(md_dir):
        for filename in os.listdir(md_dir):
            if filename.endswith('.md'):
                filepath = os.path.join(md_dir, filename)
                size = os.path.getsize(filepath)
                files.append({
                    'filename': filename,
                    'size': size,
                    'size_kb': round(size / 1024, 2)
                })
    return jsonify(files)

@app.route('/api/courses/from-md', methods=['POST'])
def create_course_from_md():
    """Crea un corso completo da un file MD"""
    data = request.json
    filename = data.get('filename')
    
    if not filename:
        return jsonify({'error': 'Nome file non fornito'}), 400
    
    filepath = os.path.join(app.config['MD_SOURCE_DIR'], filename)
    if not os.path.exists(filepath):
        return jsonify({'error': f'File {filename} non trovato'}), 404
    
    try:
        # Leggi e converti il file
        with open(filepath, 'rb') as f:
            content = f.read()
        
        # Se è RTF, estrai il Markdown
        if content.startswith(b'{\\rtf'):
            content = extract_markdown_from_rtf(content)
        else:
            content = content.decode('utf-8')
        
        # Parsa il corso e le lezioni
        course_data, lessons_data = parse_course_from_markdown(content, filename)
        
        course_code = course_data.get('code', filename.replace('.md', '').upper())
        
        # Verifica se il corso esiste già
        existing_course = Course.query.filter_by(code=course_code).first()
        
        if existing_course and not data.get('overwrite', False):
            return jsonify({
                'error': f'Il corso con codice "{course_code}" esiste già. Vuoi sovrascriverlo?',
                'existing_course_id': existing_course.id,
                'code': course_code
            }), 409
        
        # Se esiste e dobbiamo sovrascrivere, elimina lezioni esistenti
        if existing_course and data.get('overwrite', False):
            Lesson.query.filter_by(course_id=existing_course.id).delete()
            course = existing_course
            course.name = course_data.get('name', course.name)
            course.description = course_data.get('description', course.description)
            course.total_hours = course_data.get('total_hours', course.total_hours)
            course.theory_hours = course_data.get('theory_hours', course.theory_hours)
            course.practice_hours = course_data.get('practice_hours', course.practice_hours)
        else:
            # Crea nuovo corso
            course = Course(
                code=course_code,
                name=course_data.get('name', 'Nuovo Corso'),
                description=course_data.get('description', ''),
                total_hours=course_data.get('total_hours', 80),
                theory_hours=course_data.get('theory_hours', 40),
                practice_hours=course_data.get('practice_hours', 40)
            )
            db.session.add(course)
        
        db.session.flush()
        
        # Crea le lezioni
        for lesson_data in lessons_data:
            lesson = Lesson(
                course_id=course.id,
                title=lesson_data['title'],
                description=lesson_data.get('description', ''),
                lesson_type=lesson_data.get('lesson_type', 'theory'),
                duration_hours=lesson_data.get('duration_hours', 2.0),
                order=lesson_data.get('order', len(lessons_data)),
                content=lesson_data.get('content', ''),
                objectives=json.dumps(lesson_data.get('objectives', [])),
                materials=json.dumps(lesson_data.get('materials', [])),
                exercises=json.dumps(lesson_data.get('exercises', []))
            )
            db.session.add(lesson)
        
        db.session.commit()
        return jsonify({
            'id': course.id,
            'message': f'Corso {"aggiornato" if existing_course and data.get("overwrite") else "creato"} con successo da {filename}',
            'lessons_count': len(lessons_data)
        }), 201
    except Exception as e:
        db.session.rollback()
        import traceback
        return jsonify({'error': f'Errore durante la creazione: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses', methods=['POST'])
def create_course():
    try:
        # Verifica che il database sia inizializzato
        try:
            # Prova una query semplice per verificare che il database funzioni
            Course.query.limit(1).all()
        except Exception as db_error:
            # Se il database non esiste, inizializzalo
            error_str = str(db_error)
            if 'no such table' in error_str.lower() or 'operationalerror' in error_str.lower():
                logger.warning(f"Database non inizializzato in create_course, inizializzazione in corso... Errore: {db_error}")
                try:
                    init_database()
                    logger.info("Database inizializzato con successo in create_course")
                except Exception as init_error:
                    logger.error(f"Errore durante inizializzazione database in create_course: {init_error}")
                    import traceback
                    logger.error(traceback.format_exc())
                    return jsonify({'error': f'Database non inizializzato: {str(init_error)}'}), 500
            else:
                # Altro tipo di errore, rilancia
                raise
        
        data = request.json
        
        # Validazione dei dati obbligatori
        if not data:
            return jsonify({'error': 'Nessun dato fornito'}), 400
        
        # Verifica campi obbligatori
        required_fields = ['code', 'name', 'total_hours', 'theory_hours', 'practice_hours']
        missing_fields = [field for field in required_fields if field not in data or data[field] is None or data[field] == '']
        
        if missing_fields:
            return jsonify({
                'error': f'Campi obbligatori mancanti: {", ".join(missing_fields)}'
            }), 400
        
        # Validazione valori numerici
        try:
            total_hours = int(data['total_hours'])
            theory_hours = int(data['theory_hours'])
            practice_hours = int(data['practice_hours'])
            
            if total_hours <= 0 or theory_hours < 0 or practice_hours < 0:
                return jsonify({
                    'error': 'Le ore devono essere numeri positivi'
                }), 400
            
            if theory_hours + practice_hours > total_hours:
                return jsonify({
                    'error': 'La somma delle ore teoriche e pratiche non può superare le ore totali'
                }), 400
        except (ValueError, TypeError):
            return jsonify({
                'error': 'Le ore devono essere numeri validi'
            }), 400
        
        # Validazione code e name
        code = data['code'].strip()
        name = data['name'].strip()
        
        if not code:
            return jsonify({'error': 'Il codice corso non può essere vuoto'}), 400
        if not name:
            return jsonify({'error': 'Il nome corso non può essere vuoto'}), 400
        
        # Verifica se il codice esiste già
        existing_course = Course.query.filter_by(code=code).first()
        if existing_course:
            return jsonify({
                'error': f'Un corso con il codice "{code}" esiste già'
            }), 409
        
        # Crea il corso
        course = Course(
            code=code,
            name=name,
            description=data.get('description', '').strip(),
            total_hours=total_hours,
            theory_hours=theory_hours,
            practice_hours=practice_hours,
            num_lessons=data.get('num_lessons', 0)
        )
        db.session.add(course)
        db.session.commit()
        return jsonify({'id': course.id, 'message': 'Corso creato con successo'}), 201
    except Exception as e:
        db.session.rollback()
        import traceback
        error_msg = f"Errore durante la creazione del corso: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return jsonify({'error': error_msg}), 500

@app.route('/api/courses/<int:course_id>', methods=['GET'])
def get_course(course_id):
    course = Course.query.get_or_404(course_id)
    lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
    return jsonify({
        'id': course.id,
        'code': course.code,
        'name': course.name,
        'description': course.description,
        'total_hours': course.total_hours,
        'theory_hours': course.theory_hours,
        'practice_hours': course.practice_hours,
        'num_lessons': course.num_lessons or 0,
        'lessons': [{
            'id': l.id,
            'title': l.title,
            'description': l.description,
            'lesson_type': l.lesson_type,
            'duration_hours': l.duration_hours,
            'order': l.order,
            'content': l.content,
            'objectives': json.loads(l.objectives) if l.objectives else [],
            'materials': json.loads(l.materials) if l.materials else [],
            'exercises': json.loads(l.exercises) if l.exercises else [],
            'lesson_date': l.lesson_date.isoformat() if l.lesson_date else None
        } for l in lessons]
    })

@app.route('/api/courses/<int:course_id>', methods=['PUT'])
def update_course(course_id):
    course = Course.query.get_or_404(course_id)
    data = request.json
    course.code = data.get('code', course.code)
    course.name = data.get('name', course.name)
    course.description = data.get('description', course.description)
    course.total_hours = data.get('total_hours', course.total_hours)
    course.theory_hours = data.get('theory_hours', course.theory_hours)
    course.practice_hours = data.get('practice_hours', course.practice_hours)
    course.num_lessons = data.get('num_lessons', course.num_lessons)
    db.session.commit()
    return jsonify({'message': 'Corso aggiornato con successo'})

@app.route('/api/courses/<int:course_id>', methods=['DELETE'])
def delete_course(course_id):
    course = Course.query.get_or_404(course_id)
    
    # Elimina tutte le lezioni associate
    Lesson.query.filter_by(course_id=course_id).delete()
    
    # Elimina il corso
    db.session.delete(course)
    db.session.commit()
    
    return jsonify({'message': 'Corso eliminato con successo'}), 200

@app.route('/api/courses/<int:course_id>/lessons', methods=['POST'])
def create_lesson(course_id):
    course = Course.query.get_or_404(course_id)
    data = request.json
    lesson_date = None
    if data.get('lesson_date'):
        try:
            lesson_date = datetime.strptime(data['lesson_date'], '%Y-%m-%d').date()
        except (ValueError, TypeError):
            pass
    
    lesson = Lesson(
        course_id=course_id,
        title=data['title'],
        description=data.get('description', ''),
        lesson_type=data['lesson_type'],
        duration_hours=data['duration_hours'],
        order=data['order'],
        content=data.get('content', ''),
        objectives=json.dumps(data.get('objectives', [])),
        materials=json.dumps(data.get('materials', [])),
        exercises=json.dumps(data.get('exercises', [])),
        lesson_date=lesson_date
    )
    db.session.add(lesson)
    db.session.commit()
    return jsonify({'id': lesson.id, 'message': 'Lezione creata con successo'}), 201

@app.route('/api/courses/<int:course_id>/lessons/reorder', methods=['PUT'])
def reorder_lessons(course_id):
    """Riordina le lezioni di un corso"""
    try:
        data = request.json
        order_data = data.get('order', [])
        
        if not order_data:
            return jsonify({'error': 'Nessun ordine fornito'}), 400
        
        # Verifica che il corso esista
        course = Course.query.get_or_404(course_id)
        
        # Aggiorna l'ordine di ogni lezione
        for item in order_data:
            lesson_id = item.get('lesson_id')
            new_order = item.get('order')
            
            if lesson_id and new_order:
                lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first()
                if lesson:
                    lesson.order = new_order
        
        db.session.commit()
        return jsonify({'message': 'Ordine lezioni aggiornato con successo'}), 200
        
    except Exception as e:
        db.session.rollback()
        logger.error(f"Errore nel riordino lezioni: {e}")
        return jsonify({'error': f'Errore nel riordino lezioni: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>', methods=['PUT'])
def update_lesson(course_id, lesson_id):
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    data = request.json
    
    # Salva versione corrente prima di aggiornare (versioning)
    try:
        # Calcola il prossimo numero di versione
        max_version = db.session.query(db.func.max(LessonVersion.version_number)).filter_by(lesson_id=lesson_id).scalar() or 0
        next_version = max_version + 1
        
        # Crea una versione della lezione corrente
        version = LessonVersion(
            lesson_id=lesson_id,
            version_number=next_version,
            title=lesson.title,
            description=lesson.description,
            lesson_type=lesson.lesson_type,
            duration_hours=lesson.duration_hours,
            order=lesson.order,
            content=lesson.content,
            objectives=lesson.objectives,
            materials=lesson.materials,
            exercises=lesson.exercises,
            lesson_date=lesson.lesson_date,
            comment=data.get('version_comment', '')  # Commento opzionale dalla richiesta
        )
        db.session.add(version)
    except Exception as e:
        logger.warning(f"Errore nel salvataggio versione: {e}")
        # Continua anche se il salvataggio della versione fallisce
    
    # Aggiorna la lezione
    lesson.title = data.get('title', lesson.title)
    lesson.description = data.get('description', lesson.description)
    lesson.lesson_type = data.get('lesson_type', lesson.lesson_type)
    lesson.duration_hours = data.get('duration_hours', lesson.duration_hours)
    lesson.order = data.get('order', lesson.order)
    lesson.content = data.get('content', lesson.content)
    lesson.objectives = json.dumps(data.get('objectives', json.loads(lesson.objectives) if lesson.objectives else []))
    lesson.materials = json.dumps(data.get('materials', json.loads(lesson.materials) if lesson.materials else []))
    lesson.exercises = json.dumps(data.get('exercises', json.loads(lesson.exercises) if lesson.exercises else []))
    
    # Aggiorna la data della lezione
    if 'lesson_date' in data:
        if data['lesson_date']:
            try:
                lesson.lesson_date = datetime.strptime(data['lesson_date'], '%Y-%m-%d').date()
            except (ValueError, TypeError):
                lesson.lesson_date = None
        else:
            lesson.lesson_date = None
    
    # Traccia la modifica
    lesson.updated_at = datetime.utcnow()
    lesson.modification_count = (lesson.modification_count or 0) + 1
    
    db.session.commit()
    return jsonify({'message': 'Lezione aggiornata con successo'})

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>', methods=['DELETE'])
def delete_lesson(course_id, lesson_id):
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    db.session.delete(lesson)
    db.session.commit()
    return jsonify({'message': 'Lezione eliminata con successo'}), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/versions', methods=['GET'])
def get_lesson_versions(course_id, lesson_id):
    """Ottiene la cronologia delle versioni di una lezione"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    
    versions = LessonVersion.query.filter_by(lesson_id=lesson_id).order_by(LessonVersion.version_number.desc()).all()
    
    versions_data = [{
        'id': v.id,
        'version_number': v.version_number,
        'title': v.title,
        'created_at': v.created_at.isoformat() if v.created_at else None,
        'comment': v.comment,
        'created_by': v.created_by
    } for v in versions]
    
    return jsonify({
        'lesson_id': lesson_id,
        'current_version': {
            'title': lesson.title,
            'updated_at': lesson.updated_at.isoformat() if lesson.updated_at else None
        },
        'versions': versions_data
    }), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/versions/<int:version_id>', methods=['GET'])
def get_lesson_version(course_id, lesson_id, version_id):
    """Ottiene i dettagli di una versione specifica"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    version = LessonVersion.query.filter_by(id=version_id, lesson_id=lesson_id).first_or_404()
    
    return jsonify({
        'id': version.id,
        'version_number': version.version_number,
        'title': version.title,
        'description': version.description,
        'lesson_type': version.lesson_type,
        'duration_hours': version.duration_hours,
        'order': version.order,
        'content': version.content,
        'objectives': json.loads(version.objectives) if version.objectives else [],
        'materials': json.loads(version.materials) if version.materials else [],
        'exercises': json.loads(version.exercises) if version.exercises else [],
        'lesson_date': version.lesson_date.isoformat() if version.lesson_date else None,
        'comment': version.comment,
        'created_at': version.created_at.isoformat() if version.created_at else None,
        'created_by': version.created_by
    }), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/versions/compare', methods=['GET'])
def compare_lesson_versions(course_id, lesson_id):
    """Confronta due versioni di una lezione side-by-side"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    
    version1_id = request.args.get('version1')
    version2_id = request.args.get('version2', 'current')
    
    # Versione 1 (sempre una versione storica)
    if version1_id:
        v1 = LessonVersion.query.filter_by(id=version1_id, lesson_id=lesson_id).first_or_404()
        v1_data = {
            'id': v1.id,
            'version_number': v1.version_number,
            'title': v1.title,
            'description': v1.description,
            'content': v1.content,
            'objectives': json.loads(v1.objectives) if v1.objectives else [],
            'materials': json.loads(v1.materials) if v1.materials else [],
            'exercises': json.loads(v1.exercises) if v1.exercises else [],
            'created_at': v1.created_at.isoformat() if v1.created_at else None
        }
    else:
        return jsonify({'error': 'version1 è richiesto'}), 400
    
    # Versione 2 (può essere 'current' o un'altra versione)
    if version2_id == 'current':
        v2_data = {
            'id': 'current',
            'version_number': 'current',
            'title': lesson.title,
            'description': lesson.description,
            'content': lesson.content,
            'objectives': json.loads(lesson.objectives) if lesson.objectives else [],
            'materials': json.loads(lesson.materials) if lesson.materials else [],
            'exercises': json.loads(lesson.exercises) if lesson.exercises else [],
            'updated_at': lesson.updated_at.isoformat() if lesson.updated_at else None
        }
    else:
        v2 = LessonVersion.query.filter_by(id=version2_id, lesson_id=lesson_id).first_or_404()
        v2_data = {
            'id': v2.id,
            'version_number': v2.version_number,
            'title': v2.title,
            'description': v2.description,
            'content': v2.content,
            'objectives': json.loads(v2.objectives) if v2.objectives else [],
            'materials': json.loads(v2.materials) if v2.materials else [],
            'exercises': json.loads(v2.exercises) if v2.exercises else [],
            'created_at': v2.created_at.isoformat() if v2.created_at else None
        }
    
    return jsonify({
        'version1': v1_data,
        'version2': v2_data
    }), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/versions/<int:version_id>/restore', methods=['POST'])
def restore_lesson_version(course_id, lesson_id, version_id):
    """Ripristina una versione precedente della lezione"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    version = LessonVersion.query.filter_by(id=version_id, lesson_id=lesson_id).first_or_404()
    
    # Salva la versione corrente prima di ripristinare
    try:
        max_version = db.session.query(db.func.max(LessonVersion.version_number)).filter_by(lesson_id=lesson_id).scalar() or 0
        next_version = max_version + 1
        
        current_version = LessonVersion(
            lesson_id=lesson_id,
            version_number=next_version,
            title=lesson.title,
            description=lesson.description,
            lesson_type=lesson.lesson_type,
            duration_hours=lesson.duration_hours,
            order=lesson.order,
            content=lesson.content,
            objectives=lesson.objectives,
            materials=lesson.materials,
            exercises=lesson.exercises,
            lesson_date=lesson.lesson_date,
            comment='Backup prima del ripristino alla versione ' + str(version.version_number)
        )
        db.session.add(current_version)
    except Exception as e:
        logger.warning(f"Errore nel salvataggio versione corrente: {e}")
    
    # Ripristina i valori dalla versione
    lesson.title = version.title
    lesson.description = version.description
    lesson.lesson_type = version.lesson_type
    lesson.duration_hours = version.duration_hours
    lesson.order = version.order
    lesson.content = version.content
    lesson.objectives = version.objectives
    lesson.materials = version.materials
    lesson.exercises = version.exercises
    lesson.lesson_date = version.lesson_date
    lesson.updated_at = datetime.utcnow()
    lesson.modification_count = (lesson.modification_count or 0) + 1
    
    db.session.commit()
    return jsonify({'message': f'Lezione ripristinata alla versione {version.version_number}'}), 200

# ==================== API COLLABORAZIONE ====================

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/notes', methods=['GET'])
def get_lesson_notes(course_id, lesson_id):
    """Ottiene tutte le note di una lezione"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    notes = LessonNote.query.filter_by(lesson_id=lesson_id).order_by(LessonNote.created_at.desc()).all()
    
    return jsonify([{
        'id': n.id,
        'content': n.content,
        'is_private': n.is_private,
        'created_by': n.created_by,
        'created_at': n.created_at.isoformat() if n.created_at else None,
        'updated_at': n.updated_at.isoformat() if n.updated_at else None
    } for n in notes]), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/notes', methods=['POST'])
def create_lesson_note(course_id, lesson_id):
    """Crea una nuova nota per una lezione"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    data = request.json
    
    is_private = data.get('is_private', True)
    
    note = LessonNote(
        lesson_id=lesson_id,
        content=data.get('content', ''),
        is_private=is_private,
        created_by=data.get('created_by', 'Utente')
    )
    db.session.add(note)
    db.session.commit()
    
    # Crea notifica solo se la nota NON è privata
    if not is_private:
        notification = Notification(
            lesson_id=lesson_id,
            course_id=course_id,
            type='note',
            title='Nuova nota aggiunta',
            message=f'Una nuova nota pubblica è stata aggiunta alla lezione "{lesson.title}"'
        )
        db.session.add(notification)
        db.session.commit()
    
    return jsonify({
        'id': note.id,
        'message': 'Nota creata con successo'
    }), 201

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/notes/<int:note_id>', methods=['PUT'])
def update_lesson_note(course_id, lesson_id, note_id):
    """Aggiorna una nota"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    note = LessonNote.query.filter_by(id=note_id, lesson_id=lesson_id).first_or_404()
    data = request.json
    
    note.content = data.get('content', note.content)
    note.is_private = data.get('is_private', note.is_private)
    note.updated_at = datetime.utcnow()
    db.session.commit()
    
    return jsonify({'message': 'Nota aggiornata con successo'}), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/notes/<int:note_id>', methods=['DELETE'])
def delete_lesson_note(course_id, lesson_id, note_id):
    """Elimina una nota"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    note = LessonNote.query.filter_by(id=note_id, lesson_id=lesson_id).first_or_404()
    db.session.delete(note)
    db.session.commit()
    return jsonify({'message': 'Nota eliminata con successo'}), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/comments', methods=['GET'])
def get_lesson_comments(course_id, lesson_id):
    """Ottiene tutti i commenti di una lezione"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    comments = LessonComment.query.filter_by(lesson_id=lesson_id, parent_id=None).order_by(LessonComment.created_at.desc()).all()
    
    def serialize_comment(comment):
        return {
            'id': comment.id,
            'content': comment.content,
            'author': comment.author,
            'created_at': comment.created_at.isoformat() if comment.created_at else None,
            'updated_at': comment.updated_at.isoformat() if comment.updated_at else None,
            'replies': [serialize_comment(reply) for reply in comment.replies]
        }
    
    return jsonify([serialize_comment(c) for c in comments]), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/comments', methods=['POST'])
def create_lesson_comment(course_id, lesson_id):
    """Crea un nuovo commento per una lezione"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    data = request.json
    
    comment = LessonComment(
        lesson_id=lesson_id,
        content=data.get('content', ''),
        author=data.get('author', 'Utente'),
        parent_id=data.get('parent_id')
    )
    db.session.add(comment)
    db.session.commit()
    
    # Crea notifica
    notification = Notification(
        lesson_id=lesson_id,
        course_id=course_id,
        type='comment',
        title='Nuovo commento',
        message=f'{comment.author} ha commentato la lezione "{lesson.title}"'
    )
    db.session.add(notification)
    db.session.commit()
    
    return jsonify({
        'id': comment.id,
        'message': 'Commento creato con successo'
    }), 201

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/comments/<int:comment_id>', methods=['PUT'])
def update_lesson_comment(course_id, lesson_id, comment_id):
    """Aggiorna un commento"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    comment = LessonComment.query.filter_by(id=comment_id, lesson_id=lesson_id).first_or_404()
    data = request.json
    
    comment.content = data.get('content', comment.content)
    comment.updated_at = datetime.utcnow()
    db.session.commit()
    
    return jsonify({'message': 'Commento aggiornato con successo'}), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/comments/<int:comment_id>', methods=['DELETE'])
def delete_lesson_comment(course_id, lesson_id, comment_id):
    """Elimina un commento"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    comment = LessonComment.query.filter_by(id=comment_id, lesson_id=lesson_id).first_or_404()
    db.session.delete(comment)
    db.session.commit()
    return jsonify({'message': 'Commento eliminato con successo'}), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/reminders', methods=['GET'])
def get_lesson_reminders(course_id, lesson_id):
    """Ottiene tutti i promemoria di una lezione"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    reminders = LessonReminder.query.filter_by(lesson_id=lesson_id).order_by(LessonReminder.reminder_date.asc()).all()
    
    return jsonify([{
        'id': r.id,
        'title': r.title,
        'description': r.description,
        'reminder_date': r.reminder_date.isoformat() if r.reminder_date else None,
        'is_completed': r.is_completed,
        'created_by': r.created_by,
        'created_at': r.created_at.isoformat() if r.created_at else None
    } for r in reminders]), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/reminders', methods=['POST'])
def create_lesson_reminder(course_id, lesson_id):
    """Crea un nuovo promemoria per una lezione"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    data = request.json
    
    reminder_date = None
    if data.get('reminder_date'):
        try:
            reminder_date = datetime.fromisoformat(data['reminder_date'].replace('Z', '+00:00'))
        except (ValueError, TypeError):
            try:
                reminder_date = datetime.strptime(data['reminder_date'], '%Y-%m-%dT%H:%M:%S')
            except (ValueError, TypeError):
                pass
    
    reminder = LessonReminder(
        lesson_id=lesson_id,
        title=data.get('title', ''),
        description=data.get('description', ''),
        reminder_date=reminder_date or datetime.utcnow(),
        created_by=data.get('created_by', 'Utente')
    )
    db.session.add(reminder)
    db.session.commit()
    
    # Crea notifica
    notification = Notification(
        lesson_id=lesson_id,
        course_id=course_id,
        type='reminder',
        title='Nuovo promemoria',
        message=f'Promemoria creato per la lezione "{lesson.title}": {reminder.title}'
    )
    db.session.add(notification)
    db.session.commit()
    
    return jsonify({
        'id': reminder.id,
        'message': 'Promemoria creato con successo'
    }), 201

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/reminders/<int:reminder_id>', methods=['PUT'])
def update_lesson_reminder(course_id, lesson_id, reminder_id):
    """Aggiorna un promemoria"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    reminder = LessonReminder.query.filter_by(id=reminder_id, lesson_id=lesson_id).first_or_404()
    data = request.json
    
    reminder.title = data.get('title', reminder.title)
    reminder.description = data.get('description', reminder.description)
    reminder.is_completed = data.get('is_completed', reminder.is_completed)
    
    if data.get('reminder_date'):
        try:
            reminder.reminder_date = datetime.fromisoformat(data['reminder_date'].replace('Z', '+00:00'))
        except (ValueError, TypeError):
            try:
                reminder.reminder_date = datetime.strptime(data['reminder_date'], '%Y-%m-%dT%H:%M:%S')
            except (ValueError, TypeError):
                pass
    
    db.session.commit()
    return jsonify({'message': 'Promemoria aggiornato con successo'}), 200

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/reminders/<int:reminder_id>', methods=['DELETE'])
def delete_lesson_reminder(course_id, lesson_id, reminder_id):
    """Elimina un promemoria"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    reminder = LessonReminder.query.filter_by(id=reminder_id, lesson_id=lesson_id).first_or_404()
    db.session.delete(reminder)
    db.session.commit()
    return jsonify({'message': 'Promemoria eliminato con successo'}), 200

@app.route('/api/notifications', methods=['GET'])
def get_notifications():
    """Ottiene tutte le notifiche non lette"""
    unread_only = request.args.get('unread_only', 'false').lower() == 'true'
    
    query = Notification.query
    if unread_only:
        query = query.filter_by(is_read=False)
    
    notifications = query.order_by(Notification.created_at.desc()).limit(50).all()
    
    return jsonify([{
        'id': n.id,
        'lesson_id': n.lesson_id,
        'course_id': n.course_id,
        'type': n.type,
        'title': n.title,
        'message': n.message,
        'is_read': n.is_read,
        'created_at': n.created_at.isoformat() if n.created_at else None
    } for n in notifications]), 200

@app.route('/api/notifications/<int:notification_id>/read', methods=['PUT'])
def mark_notification_read(notification_id):
    """Segna una notifica come letta"""
    notification = Notification.query.get_or_404(notification_id)
    notification.is_read = True
    db.session.commit()
    return jsonify({'message': 'Notifica segnata come letta'}), 200

@app.route('/api/notifications/read-all', methods=['PUT'])
def mark_all_notifications_read():
    """Segna tutte le notifiche come lette"""
    Notification.query.update({'is_read': True})
    db.session.commit()
    return jsonify({'message': 'Tutte le notifiche segnate come lette'}), 200

@app.route('/api/notifications/<int:notification_id>', methods=['DELETE'])
def delete_notification(notification_id):
    """Elimina una notifica"""
    notification = Notification.query.get_or_404(notification_id)
    db.session.delete(notification)
    db.session.commit()
    return jsonify({'message': 'Notifica eliminata con successo'}), 200

@app.route('/api/courses/<int:course_id>/lessons/generate-objectives', methods=['POST'])
def generate_lesson_objectives(course_id):
    """Genera obiettivi formativi per una lezione basandosi sul titolo e descrizione"""
    course = Course.query.get_or_404(course_id)
    data = request.json
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata'}), 400
    
    try:
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        lesson_title = data.get('title', '')
        lesson_description = data.get('description', '')
        lesson_type = data.get('lesson_type', 'theory')
        lesson_duration = data.get('duration_hours', 2)
        
        lesson_type_text = "teorica" if lesson_type == 'theory' else "pratica"
        
        prompt = f"""Sei un docente esperto per un corso regionale di formazione professionale.

Il corso è: "{course.name}"
Durata totale: {course.total_hours} ore ({course.theory_hours} teoria, {course.practice_hours} pratica)

Genera 3-5 obiettivi formativi specifici e misurabili per questa lezione {lesson_type_text}:

**Titolo lezione:** {lesson_title}
**Descrizione:** {lesson_description or 'Nessuna descrizione'}
**Durata:** {lesson_duration} ore

**Istruzioni:**
1. Gli obiettivi devono essere specifici, misurabili e raggiungibili nella durata della lezione
2. Formulali in modo chiaro e professionale
3. Usa verbi all'infinito (es: "Comprendere...", "Applicare...", "Creare...")
4. Adattali al tipo di lezione ({lesson_type_text})
5. Restituisci SOLO gli obiettivi, uno per riga, senza numerazione o punti elenco

Restituisci SOLO gli obiettivi, uno per riga."""
        
        # Recupera preferenze AI
        ai_model = get_preference_value('aiModel', 'gpt-4o')
        ai_temperature = float(get_preference_value('aiTemperature', 0.7))
        ai_max_tokens = int(get_preference_value('aiMaxTokens', 500))
        
        response = client.chat.completions.create(
            model=ai_model,
            messages=[
                {"role": "system", "content": "Sei un assistente esperto nella creazione di obiettivi formativi per corsi di formazione professionale."},
                {"role": "user", "content": prompt}
            ],
            temperature=ai_temperature,
            max_tokens=ai_max_tokens
        )
        
        objectives_text = response.choices[0].message.content.strip()
        # Pulisci il testo rimuovendo numerazioni e punti elenco
        objectives_lines = [line.strip() for line in objectives_text.split('\n') if line.strip()]
        objectives_clean = [line.lstrip('0123456789.-) ').strip() for line in objectives_lines]
        objectives_clean = [obj for obj in objectives_clean if obj]
        
        return jsonify({
            'objectives': objectives_clean,
            'objectives_text': '\n'.join(objectives_clean)
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Errore durante la generazione obiettivi: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/lessons/generate-materials', methods=['POST'])
def generate_lesson_materials(course_id):
    """Genera suggerimenti materiali per una lezione"""
    course = Course.query.get_or_404(course_id)
    data = request.json
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata'}), 400
    
    try:
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        lesson_title = data.get('title', '')
        lesson_description = data.get('description', '')
        lesson_type = data.get('lesson_type', 'theory')
        lesson_content = data.get('content', '')
        
        lesson_type_text = "teorica" if lesson_type == 'theory' else "pratica"
        
        prompt = f"""Sei un docente esperto per un corso regionale di formazione professionale.

Il corso è: "{course.name}"

Genera 4-6 suggerimenti di materiali didattici (libri, articoli, video, strumenti online, software, ecc.) per questa lezione {lesson_type_text}:

**Titolo lezione:** {lesson_title}
**Descrizione:** {lesson_description or 'Nessuna descrizione'}
**Contenuto:** {lesson_content[:500] if lesson_content else 'Nessun contenuto ancora'}

**Istruzioni:**
1. I materiali devono essere pertinenti e utili per la lezione
2. Includi vari tipi: libri, articoli, video YouTube, strumenti online, software, risorse gratuite quando possibile
3. Sii specifico (titoli, autori, link se noti)
4. Adattali al tipo di lezione ({lesson_type_text})
5. Restituisci SOLO i materiali, uno per riga, senza numerazione

Restituisci SOLO i materiali, uno per riga."""
        
        # Recupera preferenze AI
        ai_model = get_preference_value('aiModel', 'gpt-4o')
        ai_temperature = float(get_preference_value('aiTemperature', 0.7))
        ai_max_tokens = int(get_preference_value('aiMaxTokens', 500))
        
        response = client.chat.completions.create(
            model=ai_model,
            messages=[
                {"role": "system", "content": "Sei un assistente esperto nella selezione di materiali didattici per corsi di formazione professionale."},
                {"role": "user", "content": prompt}
            ],
            temperature=ai_temperature,
            max_tokens=ai_max_tokens
        )
        
        materials_text = response.choices[0].message.content.strip()
        # Pulisci il testo
        materials_lines = [line.strip() for line in materials_text.split('\n') if line.strip()]
        materials_clean = [line.lstrip('0123456789.-) ').strip() for line in materials_lines]
        materials_clean = [mat for mat in materials_clean if mat]
        
        return jsonify({
            'materials': materials_clean,
            'materials_text': '\n'.join(materials_clean)
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Errore durante la generazione materiali: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/lessons/generate-exercises', methods=['POST'])
def generate_lesson_exercises(course_id):
    """Genera suggerimenti esercizi per una lezione"""
    course = Course.query.get_or_404(course_id)
    data = request.json
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata'}), 400
    
    try:
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        lesson_title = data.get('title', '')
        lesson_description = data.get('description', '')
        lesson_type = data.get('lesson_type', 'theory')
        lesson_content = data.get('content', '')
        lesson_duration = data.get('duration_hours', 2)
        
        lesson_type_text = "teorica" if lesson_type == 'theory' else "pratica"
        
        prompt = f"""Sei un docente esperto per un corso regionale di formazione professionale.

Il corso è: "{course.name}"

Genera 3-5 esercizi pratici o attività per questa lezione {lesson_type_text}:

**Titolo lezione:** {lesson_title}
**Descrizione:** {lesson_description or 'Nessuna descrizione'}
**Contenuto:** {lesson_content[:500] if lesson_content else 'Nessun contenuto ancora'}
**Durata:** {lesson_duration} ore

**Istruzioni:**
1. Gli esercizi devono essere pratici e applicabili
2. Adattali alla durata della lezione ({lesson_duration} ore)
3. Per lezioni teoriche: esercizi di comprensione, analisi, riflessione
4. Per lezioni pratiche: esercizi hands-on, progetti, simulazioni
5. Sii specifico e concreto
6. Restituisci SOLO gli esercizi, uno per riga, senza numerazione

Restituisci SOLO gli esercizi, uno per riga."""
        
        # Recupera preferenze AI
        ai_model = get_preference_value('aiModel', 'gpt-4o')
        ai_temperature = float(get_preference_value('aiTemperature', 0.7))
        ai_max_tokens = int(get_preference_value('aiMaxTokens', 500))
        
        response = client.chat.completions.create(
            model=ai_model,
            messages=[
                {"role": "system", "content": "Sei un assistente esperto nella creazione di esercizi pratici per corsi di formazione professionale."},
                {"role": "user", "content": prompt}
            ],
            temperature=ai_temperature,
            max_tokens=ai_max_tokens
        )
        
        exercises_text = response.choices[0].message.content.strip()
        # Pulisci il testo
        exercises_lines = [line.strip() for line in exercises_text.split('\n') if line.strip()]
        exercises_clean = [line.lstrip('0123456789.-) ').strip() for line in exercises_lines]
        exercises_clean = [ex for ex in exercises_clean if ex]
        
        return jsonify({
            'exercises': exercises_clean,
            'exercises_text': '\n'.join(exercises_clean)
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Errore durante la generazione esercizi: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/lessons/optimize-content', methods=['POST'])
def optimize_lesson_content(course_id):
    """Ottimizza i contenuti esistenti di una lezione"""
    course = Course.query.get_or_404(course_id)
    data = request.json
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata'}), 400
    
    try:
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        lesson_title = data.get('title', '')
        lesson_description = data.get('description', '')
        lesson_type = data.get('lesson_type', 'theory')
        lesson_content = data.get('content', '')
        
        if not lesson_content or len(lesson_content.strip()) < 50:
            return jsonify({'error': 'Il contenuto è troppo breve per essere ottimizzato. Aggiungi almeno 50 caratteri.'}), 400
        
        lesson_type_text = "teorica" if lesson_type == 'theory' else "pratica"
        
        prompt = f"""Sei un docente esperto per un corso regionale di formazione professionale.

Il corso è: "{course.name}"
Durata totale: {course.total_hours} ore ({course.theory_hours} teoria, {course.practice_hours} pratica)

Ottimizza e migliora il contenuto di questa lezione {lesson_type_text}:

**Titolo lezione:** {lesson_title}
**Descrizione:** {lesson_description or 'Nessuna descrizione'}

**Contenuto attuale da ottimizzare:**
{lesson_content}

**Istruzioni per l'ottimizzazione:**
1. Migliora la chiarezza e la struttura del contenuto
2. Assicurati che sia ben organizzato con sezioni e sottosezioni appropriate
3. Aggiungi esempi pratici e concreti dove utile
4. Migliora la formattazione Markdown
5. Verifica che il contenuto sia completo e coerente
6. Mantieni il tono professionale e adatto a un corso regionale
7. Assicurati che il contenuto sia adeguato per una lezione {lesson_type_text}
8. Non cambiare il significato o le informazioni principali, solo migliora presentazione e completezza

Restituisci SOLO il contenuto ottimizzato in formato Markdown, senza aggiungere commenti o spiegazioni."""
        
        # Recupera preferenze AI
        ai_model = get_preference_value('aiModel', 'gpt-4o')
        ai_temperature = float(get_preference_value('aiTemperature', 0.7))
        ai_max_tokens = int(get_preference_value('aiMaxTokens', 2000))
        
        response = client.chat.completions.create(
            model=ai_model,
            messages=[
                {"role": "system", "content": "Sei un assistente esperto nell'ottimizzazione di materiali didattici per corsi di formazione professionale."},
                {"role": "user", "content": prompt}
            ],
            temperature=ai_temperature,
            max_tokens=ai_max_tokens
        )
        
        optimized_content = response.choices[0].message.content.strip()
        
        return jsonify({
            'optimized_content': optimized_content
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Errore durante l\'ottimizzazione: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/expand', methods=['POST'])
def expand_lesson_with_chatgpt(course_id, lesson_id):
    """Espande e dettaglia una lezione usando ChatGPT"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    course = Course.query.get_or_404(course_id)
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata. Imposta la variabile d\'ambiente OPENAI_API_KEY'}), 400
    
    try:
        # Inizializza il client OpenAI
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        # Ottieni materiali di addestramento per il corso
        training_materials = TrainingMaterial.query.filter_by(course_id=course_id).all()
        knowledge_base = ""
        if training_materials:
            knowledge_base = "\n\n**CONOSCENZA AGGIUNTIVA DAL MATERIALE DI ADDESTRAMENTO:**\n\n"
            for tm in training_materials:
                if tm.summary:
                    source = tm.filename if tm.material_type == 'file' else tm.url
                    knowledge_base += f"---\n**Fonte: {source}**\n{tm.summary}\n\n"
        
        # Costruisci il prompt per ChatGPT
        lesson_type_text = "teorica" if lesson.lesson_type == 'theory' else "pratica"
        objectives_text = ""
        if lesson.objectives:
            objectives = json.loads(lesson.objectives) if isinstance(lesson.objectives, str) else lesson.objectives
            objectives_text = "\n".join([f"- {obj}" for obj in objectives])
        
        prompt = f"""Sei un docente esperto per un corso regionale di formazione professionale.

Il corso è: "{course.name}"
Durata totale: {course.total_hours} ore ({course.theory_hours} teoria, {course.practice_hours} pratica)

Devi dettagliare e approfondire questa lezione {lesson_type_text}:

**Titolo lezione:** {lesson.title}
**Descrizione:** {lesson.description or 'Nessuna descrizione'}
**Durata:** {lesson.duration_hours} ore
**Ordine:** {lesson.order}

**Obiettivi della lezione:**
{objectives_text if objectives_text else "Nessun obiettivo specificato"}

**Contenuto attuale (da espandere):**
{lesson.content or 'Contenuto minimo - da sviluppare completamente'}
{knowledge_base}
**Istruzioni:**
1. Espandi e dettaglia il contenuto della lezione in modo completo e professionale
2. Usa la conoscenza aggiuntiva fornita nei materiali di addestramento per arricchire il contenuto
3. Mantieni un formato Markdown ben strutturato
4. Aggiungi sezioni appropriate come:
   - Introduzione
   - Contenuti principali (con sottosezioni)
   - Esempi pratici (se lezione teorica) o esercizi guidati (se lezione pratica)
   - Riepilogo
5. Il contenuto deve essere adatto a un corso regionale di formazione professionale
6. Usa un linguaggio chiaro e professionale
7. Includi esempi concreti e pratici quando possibile

Restituisci SOLO il contenuto espanso in formato Markdown, senza aggiungere commenti o spiegazioni aggiuntive."""

        # Recupera preferenze AI
        ai_model = get_preference_value('aiModel', 'gpt-4o-mini')
        ai_temperature = float(get_preference_value('aiTemperature', 0.7))
        ai_max_tokens = int(get_preference_value('aiMaxTokens', 2000))
        
        # Chiama ChatGPT
        response = client.chat.completions.create(
            model=ai_model,
            messages=[
                {"role": "system", "content": "Sei un assistente esperto nella creazione di materiali didattici per corsi di formazione professionale."},
                {"role": "user", "content": prompt}
            ],
            temperature=ai_temperature,
            max_tokens=ai_max_tokens
        )
        
        expanded_content = response.choices[0].message.content.strip()
        
        # Aggiorna la lezione con il contenuto espanso
        # Mantieni il contenuto originale e aggiungi quello espanso
        if lesson.content:
            lesson.content = f"{lesson.content}\n\n---\n\n## Contenuto Dettagliato (Generato con ChatGPT)\n\n{expanded_content}"
        else:
            lesson.content = expanded_content
        
        db.session.commit()
        
        return jsonify({
            'message': 'Lezione espansa con successo',
            'expanded_content': expanded_content
        }), 200
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'Errore durante l\'espansione: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/training-materials', methods=['GET'])
def get_training_materials(course_id):
    """Ottieni tutti i materiali di addestramento per un corso"""
    materials = TrainingMaterial.query.filter_by(course_id=course_id).order_by(TrainingMaterial.created_at.desc()).all()
    return jsonify([{
        'id': m.id,
        'material_type': m.material_type,
        'filename': m.filename,
        'url': m.url,
        'summary': m.summary,
        'created_at': m.created_at.isoformat() if m.created_at else None
    } for m in materials])

@app.route('/api/courses/<int:course_id>/train-ai', methods=['POST'])
def train_ai_on_course(course_id):
    """Aggiungi materiali di addestramento (file o URL) per l'AI"""
    course = Course.query.get_or_404(course_id)
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata. Imposta la variabile d\'ambiente OPENAI_API_KEY'}), 400
    
    try:
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        # Gestisci upload file
        if 'file' in request.files:
            file = request.files['file']
            if file.filename:
                # Salva il file temporaneamente
                import tempfile
                import PyPDF2
                from docx import Document as DocxDocument
                
                filename = file.filename
                ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
                
                with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{ext}') as tmp_file:
                    file.save(tmp_file.name)
                    tmp_path = tmp_file.name
                
                try:
                    # Estrai contenuto in base al tipo
                    content = ""
                    if ext == 'pdf':
                        with open(tmp_path, 'rb') as f:
                            pdf_reader = PyPDF2.PdfReader(f)
                            for page in pdf_reader.pages:
                                content += page.extract_text() + "\n"
                    elif ext in ['docx', 'doc']:
                        doc = DocxDocument(tmp_path)
                        for para in doc.paragraphs:
                            content += para.text + "\n"
                    else:
                        # Prova a leggere come testo
                        with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                            content = f.read()
                    
                    # Genera riassunto con ChatGPT
                    content_limited = content[:8000]  # Limita a 8000 caratteri per evitare token eccessivi
                    prompt = f"""Analizza questo documento e crea un riassunto strutturato dei concetti chiave che possono essere utili per un corso di formazione professionale.

Documento: {filename}

Contenuto:
{content_limited}

Crea un riassunto che includa:
1. Argomenti principali trattati
2. Concetti chiave
3. Informazioni utili per la didattica
4. Esempi o casi studio rilevanti

Formato il riassunto in Markdown."""

                    # Recupera preferenze AI
                    ai_model = get_preference_value('aiModel', 'gpt-4o-mini')
                    ai_temperature = float(get_preference_value('aiTemperature', 0.7))
                    ai_max_tokens = int(get_preference_value('aiMaxTokens', 1500))
                    
                    response = client.chat.completions.create(
                        model=ai_model,
                        messages=[
                            {"role": "system", "content": "Sei un assistente esperto nell'analisi di documenti didattici."},
                            {"role": "user", "content": prompt}
                        ],
                        temperature=ai_temperature,
                        max_tokens=ai_max_tokens
                    )
                    
                    summary = response.choices[0].message.content.strip()
                    
                    # Salva il materiale
                    material = TrainingMaterial(
                        course_id=course_id,
                        material_type='file',
                        filename=filename,
                        content_extracted=content[:50000],  # Limita a 50k caratteri
                        summary=summary
                    )
                    db.session.add(material)
                    db.session.commit()
                    
                    os.unlink(tmp_path)  # Elimina file temporaneo
                    
                    return jsonify({
                        'message': f'Documento "{filename}" caricato e processato con successo',
                        'material_id': material.id,
                        'summary': summary
                    }), 200
                    
                except Exception as e:
                    if os.path.exists(tmp_path):
                        os.unlink(tmp_path)
                    raise e
        
        # Gestisci URL
        elif request.json and 'url' in request.json:
            url = request.json['url']
            
            try:
                import requests
                from bs4 import BeautifulSoup
                
                # Scarica il contenuto della pagina
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                }
                response_http = requests.get(url, headers=headers, timeout=30)
                response_http.raise_for_status()
                
                # Estrai testo dalla pagina
                soup = BeautifulSoup(response_http.text, 'html.parser')
                # Rimuovi script e style
                for script in soup(["script", "style"]):
                    script.decompose()
                content = soup.get_text()
                # Pulisci il testo
                lines = (line.strip() for line in content.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                content = '\n'.join(chunk for chunk in chunks if chunk)
                
                # Genera riassunto con ChatGPT
                prompt = f"""Analizza questo contenuto web e crea un riassunto strutturato dei concetti chiave che possono essere utili per un corso di formazione professionale.

URL: {url}

Contenuto:
{content[:8000]}

Crea un riassunto che includa:
1. Argomenti principali trattati
2. Concetti chiave
3. Informazioni utili per la didattica
4. Esempi o casi studio rilevanti

Formato il riassunto in Markdown."""

                # Recupera preferenze AI
                ai_model = get_preference_value('aiModel', 'gpt-4o-mini')
                ai_temperature = float(get_preference_value('aiTemperature', 0.7))
                ai_max_tokens = int(get_preference_value('aiMaxTokens', 1500))
                
                response = client.chat.completions.create(
                    model=ai_model,
                    messages=[
                        {"role": "system", "content": "Sei un assistente esperto nell'analisi di contenuti web didattici."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=ai_temperature,
                    max_tokens=ai_max_tokens
                )
                
                summary = response.choices[0].message.content.strip()
                
                # Salva il materiale
                material = TrainingMaterial(
                    course_id=course_id,
                    material_type='url',
                    url=url,
                    content_extracted=content[:50000],
                    summary=summary
                )
                db.session.add(material)
                db.session.commit()
                
                return jsonify({
                    'message': f'URL processato con successo',
                    'material_id': material.id,
                    'summary': summary
                }), 200
                
            except Exception as e:
                return jsonify({'error': f'Errore nel processare l\'URL: {str(e)}'}), 500
        
        else:
            return jsonify({'error': 'Fornisci un file o un URL'}), 400
        
    except Exception as e:
        db.session.rollback()
        import traceback
        return jsonify({'error': f'Errore durante l\'addestramento: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses/<int:course_id>/training-materials/<int:material_id>', methods=['DELETE'])
def delete_training_material(course_id, material_id):
    """Elimina un materiale di addestramento"""
    material = TrainingMaterial.query.filter_by(id=material_id, course_id=course_id).first_or_404()
    db.session.delete(material)
    db.session.commit()
    return jsonify({'message': 'Materiale eliminato con successo'}), 200

@app.route('/api/courses/<int:course_id>/questions', methods=['GET'])
def get_questions(course_id):
    """Ottieni tutte le domande per un corso"""
    questions = Question.query.filter_by(course_id=course_id).order_by(Question.question_number).all()
    return jsonify([{
        'id': q.id,
        'question_number': q.question_number,
        'question_text': q.question_text,
        'option_a': q.option_a,
        'option_b': q.option_b,
        'option_c': q.option_c,
        'option_d': q.option_d,
        'correct_answer': q.correct_answer
    } for q in questions])

@app.route('/api/courses/<int:course_id>/generate-questions', methods=['POST'])
def generate_questions(course_id):
    """Genera 30 domande a risposta multipla per un corso usando ChatGPT"""
    course = Course.query.get_or_404(course_id)
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata. Imposta la variabile d\'ambiente OPENAI_API_KEY'}), 400
    
    try:
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        # Ottieni tutte le lezioni del corso
        lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
        
        if not lessons:
            return jsonify({'error': 'Nessuna lezione trovata per questo corso'}), 400
        
        # Prepara il contenuto delle lezioni
        lessons_content = ""
        for lesson in lessons:
            lessons_content += f"\n\n### Lezione {lesson.order}: {lesson.title}\n"
            lessons_content += f"Tipo: {'Teorica' if lesson.lesson_type == 'theory' else 'Pratica'}\n"
            if lesson.description:
                lessons_content += f"Descrizione: {lesson.description}\n"
            if lesson.content:
                # Prendi solo i primi 2000 caratteri per lezione per evitare token eccessivi
                lessons_content += f"Contenuto: {lesson.content[:2000]}\n"
        
        # Ottieni materiali di addestramento
        training_materials = TrainingMaterial.query.filter_by(course_id=course_id).all()
        knowledge_base = ""
        if training_materials:
            knowledge_base = "\n\n**Materiali di addestramento aggiuntivi:**\n"
            for tm in training_materials:
                if tm.summary:
                    source = tm.filename if tm.material_type == 'file' else tm.url
                    knowledge_base += f"\n- {source}: {tm.summary[:500]}\n"
        
        prompt = f"""Sei un docente esperto per un corso regionale di formazione professionale.

Corso: "{course.name}"
Durata: {course.total_hours} ore ({course.theory_hours} teoria, {course.practice_hours} pratica)

Contenuto delle lezioni:
{lessons_content[:10000]}
{knowledge_base[:2000]}

**ISTRUZIONI CRITICHE - LEGGI CON ATTENZIONE:**
Genera esattamente {get_preference_value('questionsCount', 30)} domande a risposta multipla con {get_preference_value('questionsOptions', 4)} opzioni ciascuna (A, B, C, D).

REQUISITI OBBLIGATORI:
1. Le domande devono coprire tutti gli argomenti trattati nelle lezioni
2. Ogni domanda deve avere esattamente 4 opzioni (A, B, C, D)
3. **DISTRIBUZIONE RISPOSTE CORRETTE (OBBLIGATORIA):**
   - Deve essere esattamente: 7-8 risposte corrette per A, 7-8 per B, 7-8 per C, 7-8 per D
   - Totale: 30 domande = 7+8+7+8 oppure 8+7+8+7 (variazioni accettabili)
   - NON puoi avere più di 8 risposte corrette per la stessa lettera
   - NON puoi avere meno di 7 risposte corrette per la stessa lettera
   - DISTRIBUISCI LE RISPOSTE CORRETTE IN MODO CASUALE MA EQUILIBRATO
4. Le domande devono essere chiare, precise e pertinenti al contenuto del corso
5. Le opzioni sbagliate devono essere plausibili ma chiaramente errate
6. Le domande devono testare la comprensione, non solo la memorizzazione

**ESEMPIO DI DISTRIBUZIONE CORRETTA:**
- Domande 1-8: risposte corrette A, B, C, D, A, B, C, D (2 per lettera)
- Domande 9-16: risposte corrette A, B, C, D, A, B, C, D (altre 2 per lettera)
- Domande 17-24: risposte corrette A, B, C, D, A, B, C, D (altre 2 per lettera)
- Domande 25-30: risposte corrette A, B, C, D, A, B (ultime 2 per A e B, 1 per C e D)
Totale: 8 A, 8 B, 7 C, 7 D = 30 domande

**IMPORTANTE:** Prima di generare, pianifica mentalmente la distribuzione delle risposte corrette per assicurarti che sia equilibrata!

FORMATO DI RISPOSTA (JSON):
{{
  "questions": [
    {{
      "number": 1,
      "question": "Testo della domanda?",
      "options": {{
        "A": "Opzione A",
        "B": "Opzione B",
        "C": "Opzione C",
        "D": "Opzione D"
      }},
      "correct": "A"
    }},
    ...
  ]
}}

Restituisci SOLO il JSON, senza commenti o spiegazioni aggiuntive."""

        # Recupera preferenze AI
        ai_model = get_preference_value('aiModel', 'gpt-4o-mini')
        ai_temperature = float(get_preference_value('aiTemperature', 0.7))
        ai_max_tokens = int(get_preference_value('aiMaxTokens', 4000))
        
        response = client.chat.completions.create(
            model=ai_model,
            messages=[
                {"role": "system", "content": "Sei un esperto nella creazione di test e quiz per corsi di formazione professionale. Restituisci sempre JSON valido."},
                {"role": "user", "content": prompt}
            ],
            temperature=ai_temperature,
            max_tokens=ai_max_tokens,
            response_format={"type": "json_object"}
        )
        
        result = json.loads(response.choices[0].message.content)
        
        # Elimina le domande esistenti per questo corso
        Question.query.filter_by(course_id=course_id).delete()
        
        # Salva le nuove domande
        questions_data = result.get('questions', [])
        if len(questions_data) != 30:
            return jsonify({'error': f'Errore: sono state generate {len(questions_data)} domande invece di 30'}), 500
        
        # Verifica distribuzione prima di salvare
        temp_distribution = {'A': 0, 'B': 0, 'C': 0, 'D': 0}
        for q_data in questions_data:
            correct = q_data.get('correct', '').upper()
            if correct in temp_distribution:
                temp_distribution[correct] += 1
        
        # Controlla se la distribuzione è accettabile (min 6, max 9 per lettera)
        min_count = min(temp_distribution.values())
        max_count = max(temp_distribution.values())
        
        if max_count > 9 or min_count < 6:
            # Ribilancia le risposte corrette
            target_distribution = [7, 8, 7, 8]  # Totale: 30
            letters = ['A', 'B', 'C', 'D']
            balanced_answers = []
            for i, letter in enumerate(letters):
                balanced_answers.extend([letter] * target_distribution[i])
            
            # Mescola per randomizzare
            import random
            random.shuffle(balanced_answers)
            
            # Assegna le risposte corrette bilanciate
            for i, q_data in enumerate(questions_data):
                q_data['correct'] = balanced_answers[i]
        
        # Salva le domande
        for q_data in questions_data:
            question = Question(
                course_id=course_id,
                question_number=q_data['number'],
                question_text=q_data['question'],
                option_a=q_data['options']['A'],
                option_b=q_data['options']['B'],
                option_c=q_data['options']['C'],
                option_d=q_data['options']['D'],
                correct_answer=q_data['correct'].upper()
            )
            db.session.add(question)
        
        db.session.commit()
        
        # Verifica distribuzione risposte corrette finale
        questions = Question.query.filter_by(course_id=course_id).all()
        answer_distribution = {'A': 0, 'B': 0, 'C': 0, 'D': 0}
        for q in questions:
            answer_distribution[q.correct_answer] += 1
        
        return jsonify({
            'message': f'30 domande generate con successo',
            'total_questions': len(questions),
            'answer_distribution': answer_distribution
        }), 200
        
    except json.JSONDecodeError as e:
        db.session.rollback()
        return jsonify({'error': f'Errore nel parsing della risposta JSON: {str(e)}'}), 500
    except Exception as e:
        db.session.rollback()
        import traceback
        return jsonify({'error': f'Errore durante la generazione: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses/<int:course_id>/questions/export/<format_type>', methods=['GET'])
def export_questions(course_id, format_type):
    """Esporta le domande in PDF o DOCX"""
    course = Course.query.get_or_404(course_id)
    questions = Question.query.filter_by(course_id=course_id).order_by(Question.question_number).all()
    
    if not questions:
        return jsonify({'error': 'Nessuna domanda trovata per questo corso'}), 404
    
    if format_type not in ['pdf', 'docx']:
        return jsonify({'error': 'Formato non supportato. Usa "pdf" o "docx"'}), 400
    
    # Recupera preferenze esportazione
    export_author = get_preference_value('exportAuthor', '')
    export_company = get_preference_value('exportCompany', '')
    
    try:
        if format_type == 'pdf':
            from weasyprint import HTML
            
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                {f'<meta name="author" content="{export_author}">' if export_author else ''}
                {f'<meta name="company" content="{export_company}">' if export_company else ''}
                <style>
                    @page {{
                        size: A4;
                        margin: 2cm;
                    }}
                    body {{
                        font-family: 'DejaVu Sans', Arial, sans-serif;
                        font-size: 11pt;
                        line-height: 1.6;
                        color: #333;
                    }}
                    h1 {{
                        font-size: 20pt;
                        color: #2c3e50;
                        border-bottom: 3px solid #3498db;
                        padding-bottom: 10px;
                        margin-bottom: 20px;
                    }}
                    .question {{
                        margin-bottom: 20px;
                        padding: 15px;
                        border-left: 4px solid #3498db;
                        background-color: #f8f9fa;
                    }}
                    .question-number {{
                        font-weight: bold;
                        font-size: 12pt;
                        margin-bottom: 8px;
                    }}
                    .question-text {{
                        margin-bottom: 12px;
                        font-weight: 500;
                    }}
                    .option {{
                        margin: 5px 0;
                        padding-left: 20px;
                    }}
                </style>
            </head>
            <body>
                <h1>Domande - {course.name}</h1>
                <p><strong>Corso:</strong> {course.code} | <strong>Totale domande:</strong> {len(questions)}</p>
                <hr>
            """
            
            for q in questions:
                html_content += f"""
                <div class="question">
                    <div class="question-number">Domanda {q.question_number}</div>
                    <div class="question-text">{q.question_text}</div>
                    <div class="option"><strong>A)</strong> {q.option_a}</div>
                    <div class="option"><strong>B)</strong> {q.option_b}</div>
                    <div class="option"><strong>C)</strong> {q.option_c}</div>
                    <div class="option"><strong>D)</strong> {q.option_d}</div>
                </div>
                """
            
            html_content += """
                <div style="margin-top: 40px; padding-top: 20px; border-top: 1px solid #ddd; font-size: 9pt; color: #7f8c8d; text-align: center;">
                    <p>Generato il """ + datetime.now().strftime('%d/%m/%Y alle %H:%M') + """</p>
                </div>
            </body>
            </html>
            """
            
            pdf_bytes = HTML(string=html_content).write_pdf()
            
            # Traccia l'esportazione
            course.export_count = (course.export_count or 0) + 1
            db.session.commit()
            
            filename = f"{course.code}_Domande.pdf"
            return send_file(
                BytesIO(pdf_bytes),
                mimetype='application/pdf',
                as_attachment=True,
                download_name=filename
            )
        
        elif format_type == 'docx':
            doc = Document()
            
            # Imposta metadati
            if export_author:
                doc.core_properties.author = export_author
            if export_company:
                doc.core_properties.company = export_company
            
            title = doc.add_heading(f'Domande - {course.name}', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph(f'Corso: {course.code}')
            doc.add_paragraph(f'Totale domande: {len(questions)}')
            doc.add_paragraph('')
            
            for q in questions:
                doc.add_heading(f'Domanda {q.question_number}', level=2)
                doc.add_paragraph(q.question_text)
                doc.add_paragraph(f'A) {q.option_a}', style='List Bullet')
                doc.add_paragraph(f'B) {q.option_b}', style='List Bullet')
                doc.add_paragraph(f'C) {q.option_c}', style='List Bullet')
                doc.add_paragraph(f'D) {q.option_d}', style='List Bullet')
                doc.add_paragraph('')
            
            doc_bytes = BytesIO()
            doc.save(doc_bytes)
            doc_bytes.seek(0)
            
            # Traccia l'esportazione
            course.export_count = (course.export_count or 0) + 1
            db.session.commit()
            
            filename = f"{course.code}_Domande.docx"
            return send_file(
                doc_bytes,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                as_attachment=True,
                download_name=filename
            )
    
    except Exception as e:
        import traceback
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses/<int:course_id>/questions/answer-key', methods=['GET'])
def export_answer_key(course_id):
    """Esporta il correttore con le risposte esatte"""
    course = Course.query.get_or_404(course_id)
    questions = Question.query.filter_by(course_id=course_id).order_by(Question.question_number).all()
    
    if not questions:
        return jsonify({'error': 'Nessuna domanda trovata per questo corso'}), 404
    
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from io import BytesIO
        
        # Recupera preferenze esportazione
        export_author = get_preference_value('exportAuthor', '')
        export_company = get_preference_value('exportCompany', '')
        
        doc = Document()
        
        # Imposta metadati
        if export_author:
            doc.core_properties.author = export_author
        if export_company:
            doc.core_properties.company = export_company
        
        title = doc.add_heading(f'Correttore - {course.name}', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph(f'Corso: {course.code}')
        doc.add_paragraph(f'Totale domande: {len(questions)}')
        doc.add_paragraph('')
        
        # Tabella con risposte
        table = doc.add_table(rows=1, cols=2)
        table.style = 'Light Grid Accent 1'
        
        # Intestazione
        header_cells = table.rows[0].cells
        header_cells[0].text = 'Domanda'
        header_cells[1].text = 'Risposta Corretta'
        header_cells[0].paragraphs[0].runs[0].bold = True
        header_cells[1].paragraphs[0].runs[0].bold = True
        
        # Aggiungi righe
        for q in questions:
            row_cells = table.add_row().cells
            row_cells[0].text = str(q.question_number)
            row_cells[1].text = q.correct_answer
        
        doc_bytes = BytesIO()
        doc.save(doc_bytes)
        doc_bytes.seek(0)
        
        # Traccia l'esportazione
        course.export_count = (course.export_count or 0) + 1
        db.session.commit()
        
        filename = f"{course.code}_Correttore.docx"
        return send_file(
            doc_bytes,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name=filename
        )
    
    except Exception as e:
        import traceback
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses/<int:course_id>/final-report', methods=['POST'])
def generate_final_report(course_id):
    """Genera una relazione finale del corso usando ChatGPT"""
    course = Course.query.get_or_404(course_id)
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata. Imposta la variabile d\'ambiente OPENAI_API_KEY'}), 400
    
    data = request.get_json()
    user_prompt = data.get('prompt', '')
    
    if not user_prompt:
        return jsonify({'error': 'Prompt descrittivo richiesto'}), 400
    
    try:
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        # Ottieni tutte le lezioni del corso
        lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
        
        if not lessons:
            return jsonify({'error': 'Nessuna lezione trovata per questo corso'}), 400
        
        # Prepara il contenuto delle lezioni
        lessons_content = ""
        for lesson in lessons:
            lessons_content += f"\n\n### Lezione {lesson.order}: {lesson.title}\n"
            lessons_content += f"Tipo: {'Teorica' if lesson.lesson_type == 'theory' else 'Pratica'}\n"
            lessons_content += f"Durata: {lesson.duration_hours} ore\n"
            if lesson.description:
                lessons_content += f"Descrizione: {lesson.description}\n"
            if lesson.content:
                # Prendi tutto il contenuto della lezione
                lessons_content += f"Contenuto:\n{lesson.content}\n"
        
        # Ottieni materiali di addestramento
        training_materials = TrainingMaterial.query.filter_by(course_id=course_id).all()
        knowledge_base = ""
        if training_materials:
            knowledge_base = "\n\n**Materiali di addestramento aggiuntivi:**\n"
            for tm in training_materials:
                if tm.summary:
                    source = tm.filename if tm.material_type == 'file' else tm.url
                    knowledge_base += f"\n- {source}: {tm.summary[:500]}\n"
        
        prompt = f"""Sei un docente esperto e devi scrivere una relazione finale professionale per un corso regionale di formazione professionale.

**INFORMAZIONI SUL CORSO:**
- Nome: {course.name}
- Codice: {course.code}
- Durata totale: {course.total_hours} ore ({course.theory_hours} teoria, {course.practice_hours} pratica)
- Numero lezioni: {len(lessons)}

**CONTENUTO DELLE LEZIONI:**
{lessons_content[:15000]}
{knowledge_base[:2000]}

**PROMPT PERSONALIZZATO DELL'UTENTE:**
{user_prompt}

**ISTRUZIONI:**
Scrivi una relazione finale professionale e dettagliata che:
1. Segua le indicazioni del prompt personalizzato fornito dall'utente
2. Sia strutturata in modo chiaro e professionale
3. Faccia riferimento specifico alle lezioni svolte e agli argomenti trattati
4. Sia appropriata per un corso regionale di formazione professionale
5. Utilizzi un linguaggio formale ma accessibile
6. Includa sezioni logiche (es: Introduzione, Obiettivi, Contenuti, Metodologia, Risultati, Conclusioni)

La relazione deve essere completa, ben strutturata e professionale. Restituisci SOLO il testo della relazione, senza commenti aggiuntivi."""

        # Recupera preferenze AI
        ai_model = get_preference_value('aiModel', 'gpt-4o-mini')
        ai_temperature = float(get_preference_value('aiTemperature', 0.7))
        ai_max_tokens = int(get_preference_value('aiMaxTokens', 4000))
        
        response = client.chat.completions.create(
            model=ai_model,
            messages=[
                {"role": "system", "content": "Sei un esperto nella scrittura di relazioni finali per corsi di formazione professionale. Scrivi sempre relazioni professionali, strutturate e dettagliate."},
                {"role": "user", "content": prompt}
            ],
            temperature=ai_temperature,
            max_tokens=ai_max_tokens,
            timeout=120
        )
        
        report_content = response.choices[0].message.content
        
        # Salva la relazione nel database (opzionale, per ora la restituiamo solo)
        return jsonify({
            'message': 'Relazione finale generata con successo',
            'report': report_content
        }), 200
        
    except Exception as e:
        import traceback
        return jsonify({'error': f'Errore durante la generazione: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses/<int:course_id>/final-report/export/<format_type>', methods=['POST'])
def export_final_report(course_id, format_type):
    """Esporta la relazione finale in PDF o DOCX"""
    course = Course.query.get_or_404(course_id)
    
    data = request.get_json()
    report_content = data.get('report', '')
    
    if not report_content:
        return jsonify({'error': 'Contenuto della relazione mancante'}), 400
    
    if format_type not in ['pdf', 'docx']:
        return jsonify({'error': 'Formato non supportato. Usa "pdf" o "docx"'}), 400
    
    try:
        if format_type == 'pdf':
            from weasyprint import HTML
            
            # Converti Markdown a HTML
            html_content_md = markdown.markdown(report_content, extensions=['extra', 'nl2br'])
            
            html_doc = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                {f'<meta name="author" content="{export_author}">' if export_author else ''}
                {f'<meta name="company" content="{export_company}">' if export_company else ''}
                <style>
                    @page {{
                        size: A4;
                        margin: 2.5cm;
                    }}
                    body {{
                        font-family: 'DejaVu Sans', Arial, sans-serif;
                        font-size: 11pt;
                        line-height: 1.8;
                        color: #333;
                    }}
                    h1 {{
                        font-size: 22pt;
                        color: #2c3e50;
                        border-bottom: 3px solid #3498db;
                        padding-bottom: 10px;
                        margin-bottom: 25px;
                        text-align: center;
                    }}
                    h2 {{
                        font-size: 16pt;
                        color: #34495e;
                        margin-top: 25px;
                        margin-bottom: 15px;
                        border-left: 4px solid #3498db;
                        padding-left: 10px;
                    }}
                    h3 {{
                        font-size: 13pt;
                        color: #555;
                        margin-top: 20px;
                        margin-bottom: 10px;
                    }}
                    p {{
                        margin-bottom: 12px;
                        text-align: justify;
                    }}
                    ul, ol {{
                        margin-bottom: 15px;
                        padding-left: 25px;
                    }}
                    li {{
                        margin-bottom: 8px;
                    }}
                    strong {{
                        color: #2c3e50;
                    }}
                </style>
            </head>
            <body>
                <h1>Relazione Finale - {course.name}</h1>
                <p style="text-align: center; color: #7f8c8d; margin-bottom: 30px;">
                    <strong>Corso:</strong> {course.code} | 
                    <strong>Durata:</strong> {course.total_hours} ore | 
                    <strong>Data generazione:</strong> {datetime.now().strftime('%d/%m/%Y')}
                </p>
                <hr style="margin-bottom: 30px;">
                {html_content_md}
                <div style="margin-top: 50px; padding-top: 20px; border-top: 1px solid #ddd; font-size: 9pt; color: #7f8c8d; text-align: center;">
                    <p>Relazione generata il {datetime.now().strftime('%d/%m/%Y alle %H:%M')}</p>
                </div>
            </body>
            </html>
            """
            
            pdf_bytes = HTML(string=html_doc).write_pdf()
            
            # Traccia l'esportazione
            course.export_count = (course.export_count or 0) + 1
            db.session.commit()
            
            filename = f"{course.code}_RelazioneFinale.pdf"
            return send_file(
                BytesIO(pdf_bytes),
                mimetype='application/pdf',
                as_attachment=True,
                download_name=filename
            )
        
        elif format_type == 'docx':
            doc = Document()
            
            # Imposta metadati
            if export_author:
                doc.core_properties.author = export_author
            if export_company:
                doc.core_properties.company = export_company
            
            title = doc.add_heading(f'Relazione Finale - {course.name}', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph(f'Corso: {course.code}')
            doc.add_paragraph(f'Durata: {course.total_hours} ore')
            doc.add_paragraph(f'Data generazione: {datetime.now().strftime("%d/%m/%Y")}')
            doc.add_paragraph('')
            
            # Converti Markdown a paragrafi Word
            lines = report_content.split('\n')
            current_paragraph = None
            
            for line in lines:
                line = line.strip()
                if not line:
                    if current_paragraph:
                        doc.add_paragraph('')
                    current_paragraph = None
                    continue
                
                if line.startswith('# '):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith('### '):
                    doc.add_heading(line[4:], level=3)
                elif line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line[2:], style='List Bullet')
                elif line.startswith(tuple(str(i) + '. ' for i in range(1, 100))):
                    doc.add_paragraph(line, style='List Number')
                else:
                    # Rimuovi markdown bold/italic
                    line = line.replace('**', '').replace('*', '').replace('__', '').replace('_', '')
                    doc.add_paragraph(line)
            
            doc_bytes = BytesIO()
            doc.save(doc_bytes)
            doc_bytes.seek(0)
            
            # Traccia l'esportazione
            course.export_count = (course.export_count or 0) + 1
            db.session.commit()
            
            filename = f"{course.code}_RelazioneFinale.docx"
            return send_file(
                doc_bytes,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                as_attachment=True,
                download_name=filename
            )
    
    except Exception as e:
        import traceback
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses/<int:course_id>/calendar/ics', methods=['GET'])
def export_course_calendar(course_id):
    """Esporta tutte le date delle lezioni in formato ICS (iCalendar)"""
    course = Course.query.get_or_404(course_id)
    lessons = Lesson.query.filter_by(course_id=course_id).filter(Lesson.lesson_date.isnot(None)).order_by(Lesson.order).all()
    
    if not lessons:
        return jsonify({'error': 'Nessuna lezione con data impostata'}), 404
    
    try:
        from io import StringIO
        
        # Recupera preferenza orario predefinito
        default_time_str = get_preference_value('defaultLessonTime', '09:00')
        try:
            time_parts = str(default_time_str).split(':')
            default_hour = int(time_parts[0])
            default_minute = int(time_parts[1]) if len(time_parts) > 1 else 0
        except (ValueError, AttributeError):
            default_hour, default_minute = 9, 0
        
        # Crea il contenuto ICS
        ics_content = StringIO()
        ics_content.write('BEGIN:VCALENDAR\r\n')
        ics_content.write('VERSION:2.0\r\n')
        ics_content.write('PRODID:-//Learning Management//EN\r\n')
        ics_content.write('CALSCALE:GREGORIAN\r\n')
        ics_content.write('METHOD:PUBLISH\r\n')
        
        for lesson in lessons:
            if not lesson.lesson_date:
                continue
                
            # Calcola data e ora di inizio (usa preferenza o default: 09:00)
            start_datetime = datetime.combine(lesson.lesson_date, time(default_hour, default_minute))
            
            # Calcola data e ora di fine (inizio + durata in ore)
            duration_hours = int(lesson.duration_hours)
            duration_minutes = int((lesson.duration_hours - duration_hours) * 60)
            end_datetime = start_datetime + timedelta(hours=duration_hours, minutes=duration_minutes)
            
            # Formatta le date in formato UTC (YYYYMMDDTHHMMSSZ)
            dtstart = start_datetime.strftime('%Y%m%dT%H%M%SZ')
            dtend = end_datetime.strftime('%Y%m%dT%H%M%SZ')
            
            # Crea l'evento
            ics_content.write('BEGIN:VEVENT\r\n')
            ics_content.write(f'UID:lesson-{lesson.id}@learning-management\r\n')
            ics_content.write(f'DTSTART:{dtstart}\r\n')
            ics_content.write(f'DTEND:{dtend}\r\n')
            ics_content.write(f'SUMMARY:{lesson.title}\r\n')
            
            # Descrizione con dettagli della lezione
            description = f'Lezione {lesson.order} - {lesson.lesson_type}\n'
            if lesson.description:
                description += f'\n{lesson.description}\n'
            description += f'\nDurata: {lesson.duration_hours}h'
            
            # Escape caratteri speciali per ICS
            description = description.replace('\\', '\\\\').replace(',', '\\,').replace(';', '\\;').replace('\n', '\\n')
            ics_content.write(f'DESCRIPTION:{description}\r\n')
            
            # Location (opzionale, potrebbe essere il nome del corso)
            ics_content.write(f'LOCATION:{course.name}\r\n')
            
            # Timestamp di creazione
            created = lesson.created_at or datetime.utcnow()
            ics_content.write(f'DTSTAMP:{created.strftime("%Y%m%dT%H%M%SZ")}\r\n')
            
            ics_content.write('END:VEVENT\r\n')
        
        ics_content.write('END:VCALENDAR\r\n')
        
        # Converti in bytes
        ics_bytes = BytesIO(ics_content.getvalue().encode('utf-8'))
        
        filename = f"{course.code or 'Corso'}_{course.id}_calendario.ics"
        return send_file(
            ics_bytes,
            mimetype='text/calendar',
            as_attachment=True,
            download_name=filename
        )
    
    except Exception as e:
        import traceback
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses/<int:course_id>/slides', methods=['GET'])
def create_slides(course_id):
    """Genera una presentazione PowerPoint basata sulle lezioni del corso"""
    course = Course.query.get_or_404(course_id)
    lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
    
    if not lessons:
        return jsonify({'error': 'Nessuna lezione trovata per questo corso'}), 404
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from io import BytesIO
        
        # Crea una nuova presentazione
        prs = Presentation()
        
        # Imposta dimensioni slide (16:9)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Slide 1: Titolo
        slide_layout = prs.slide_layouts[0]  # Layout titolo
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = course.name
        subtitle.text = f"{course.code} - {course.total_hours} ore ({course.theory_hours} teoria, {course.practice_hours} pratica)"
        
        # Slide 2: Panoramica del corso
        slide_layout = prs.slide_layouts[1]  # Layout titolo e contenuto
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Panoramica del Corso"
        tf = content.text_frame
        tf.text = f"Codice: {course.code}"
        p = tf.add_paragraph()
        p.text = f"Durata totale: {course.total_hours} ore"
        p = tf.add_paragraph()
        p.text = f"Ore teoriche: {course.theory_hours}"
        p = tf.add_paragraph()
        p.text = f"Ore pratiche: {course.practice_hours}"
        p = tf.add_paragraph()
        p.text = f"Numero lezioni: {len(lessons)}"
        
        if course.description:
            p = tf.add_paragraph()
            p.text = f"\n{course.description}"
        
        # Slide per ogni lezione
        for lesson in lessons:
            # Slide titolo lezione
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"Lezione {lesson.order}: {lesson.title}"
            tf = content.text_frame
            tf.word_wrap = True
            
            # Informazioni lezione
            info_text = f"Tipo: {'Teorica' if lesson.lesson_type == 'theory' else 'Pratica'}\n"
            info_text += f"Durata: {lesson.duration_hours} ore\n"
            
            if lesson.description:
                info_text += f"\n{lesson.description}\n"
            
            tf.text = info_text
            
            # Se c'è contenuto, crea slide aggiuntive
            if lesson.content:
                # Dividi il contenuto in paragrafi/slide
                content_lines = lesson.content.split('\n')
                current_slide_content = []
                char_count = 0
                max_chars_per_slide = 800  # Limite caratteri per slide
                
                for line in content_lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    # Se la riga è un titolo (## o ###)
                    if line.startswith('##'):
                        # Salva slide precedente se c'è contenuto
                        if current_slide_content:
                            slide_layout = prs.slide_layouts[1]
                            slide = prs.slides.add_slide(slide_layout)
                            title = slide.shapes.title
                            content = slide.placeholders[1]
                            
                            title.text = f"Lezione {lesson.order}: {lesson.title}"
                            tf = content.text_frame
                            tf.word_wrap = True
                            tf.text = '\n'.join(current_slide_content)
                            
                            current_slide_content = []
                            char_count = 0
                        
                        # Nuova slide con titolo sezione
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        content = slide.placeholders[1]
                        
                        section_title = line.lstrip('#').strip()
                        title.text = section_title
                        tf = content.text_frame
                        tf.word_wrap = True
                        tf.text = ""  # Contenuto verrà aggiunto dopo
                        current_slide_content = []
                        char_count = 0
                    else:
                        # Aggiungi alla slide corrente
                        if char_count + len(line) > max_chars_per_slide and current_slide_content:
                            # Crea nuova slide
                            slide_layout = prs.slide_layouts[1]
                            slide = prs.slides.add_slide(slide_layout)
                            title = slide.shapes.title
                            content = slide.placeholders[1]
                            
                            title.text = f"Lezione {lesson.order}: {lesson.title} (continua)"
                            tf = content.text_frame
                            tf.word_wrap = True
                            tf.text = '\n'.join(current_slide_content)
                            
                            current_slide_content = [line]
                            char_count = len(line)
                        else:
                            current_slide_content.append(line)
                            char_count += len(line) + 1
                
                # Aggiungi ultima slide se c'è contenuto rimasto
                if current_slide_content:
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = f"Lezione {lesson.order}: {lesson.title}"
                    tf = content.text_frame
                    tf.word_wrap = True
                    tf.text = '\n'.join(current_slide_content)
        
        # Slide finale: Riepilogo
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Riepilogo"
        tf = content.text_frame
        tf.text = f"Corso completato: {course.name}\n"
        p = tf.add_paragraph()
        p.text = f"Totale lezioni: {len(lessons)}\n"
        p = tf.add_paragraph()
        p.text = f"Durata totale: {course.total_hours} ore"
        
        # Salva in memoria
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        filename = f"{course.code}_Presentazione.pptx"
        return send_file(
            pptx_bytes,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )
    
    except Exception as e:
        import traceback
        return jsonify({'error': f'Errore durante la creazione delle slide: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/canva/auth', methods=['GET', 'POST'])
def canva_auth():
    """Endpoint di autorizzazione OAuth - Canva chiama questo endpoint"""
    # Funzione in standby
    return jsonify({'error': 'Funzione Canva attualmente in standby'}), 503
    
    course_id = request.args.get('course_id') or request.form.get('course_id')
    
    if not app.config['CANVA_CLIENT_ID']:
        return jsonify({'error': 'Canva Client ID non configurato. Imposta la variabile d\'ambiente CANVA_CLIENT_ID'}), 400
    
    if course_id:
        # Salva il course_id nella sessione
        session['canva_course_id'] = int(course_id)
    
    # Canva si aspetta una risposta HTML con un form di autorizzazione
    # o un redirect diretto a Canva per l'autorizzazione
    client_id = app.config['CANVA_CLIENT_ID']
    redirect_uri = 'https://www.canva.com/apps/oauth/authorized'  # URL fisso di Canva
    
    # URL di autorizzazione Canva - reindirizza l'utente a Canva
    # Formato corretto secondo documentazione Canva
    from urllib.parse import urlencode
    
    params = {
        'client_id': client_id,
        'redirect_uri': redirect_uri,
        'response_type': 'code',
        'scope': 'design:content:read design:content:write',
        'state': f'course_{course_id if course_id else "default"}'  # State per sicurezza
    }
    
    auth_url = f'https://www.canva.com/api/oauth/authorize?{urlencode(params)}'
    
    # Se è una richiesta GET da Canva, mostra la pagina di autorizzazione
    if request.method == 'GET' and 'course_id' not in request.args:
        return f'''
        <html>
        <head><title>Autorizzazione Canva</title></head>
        <body>
            <h1>Autorizzazione Richiesta</h1>
            <p>Stai per essere reindirizzato a Canva per autorizzare l'applicazione.</p>
            <script>
                window.location.href = '{auth_url}';
            </script>
            <p><a href="{auth_url}">Clicca qui se non vieni reindirizzato automaticamente</a></p>
        </body>
        </html>
        '''
    
    # Se è una richiesta da JavaScript, restituisci l'URL
    return jsonify({'auth_url': auth_url}), 200

@app.route('/api/canva/token', methods=['POST'])
def canva_token_exchange():
    """Endpoint per lo scambio del codice di autorizzazione con il token - Canva chiama questo"""
    # Funzione in standby
    return jsonify({'error': 'Funzione Canva attualmente in standby'}), 503
    
    try:
        # Canva invia il codice di autorizzazione
        code = request.form.get('code') or request.json.get('code') if request.is_json else None
        redirect_uri = request.form.get('redirect_uri') or request.json.get('redirect_uri') if request.is_json else 'https://www.canva.com/apps/oauth/authorized'
        
        if not code:
            return jsonify({'error': 'Codice di autorizzazione mancante'}), 400
        
        # Scambia il codice con un token di accesso
        token_url = 'https://api.canva.com/rest/v1/oauth/token'
        
        client_id = app.config['CANVA_CLIENT_ID']
        client_secret = app.config['CANVA_CLIENT_SECRET']
        
        # Codifica client_id:client_secret in base64
        credentials = base64.b64encode(f'{client_id}:{client_secret}'.encode()).decode()
        
        token_response = requests.post(
            token_url,
            headers={
                'Authorization': f'Basic {credentials}',
                'Content-Type': 'application/x-www-form-urlencoded'
            },
            data={
                'grant_type': 'authorization_code',
                'code': code,
                'redirect_uri': redirect_uri
            }
        )
        
        if token_response.status_code != 200:
            return jsonify({
                'error': 'Errore nell\'ottenere il token',
                'details': token_response.text
            }), 400
        
        token_data = token_response.json()
        access_token = token_data.get('access_token')
        refresh_token = token_data.get('refresh_token')
        
        if not access_token:
            return jsonify({'error': 'Token di accesso non ricevuto'}), 400
        
        # Salva il token (potresti volerlo salvare nel database associato al corso)
        course_id = session.get('canva_course_id')
        
        # Restituisci i token a Canva
        return jsonify({
            'access_token': access_token,
            'refresh_token': refresh_token,
            'token_type': 'Bearer',
            'expires_in': token_data.get('expires_in', 3600),
            'course_id': course_id
        }), 200
    
    except Exception as e:
        import traceback
        return jsonify({
            'error': f'Errore durante lo scambio del token: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500

@app.route('/api/courses/<int:course_id>/canva/create', methods=['POST'])
def create_canva_presentation(course_id):
    """Crea una presentazione su Canva"""
    # Funzione in standby
    return jsonify({'error': 'Funzione Canva attualmente in standby'}), 503
    
    course = Course.query.get_or_404(course_id)
    lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
    
    if not lessons:
        return jsonify({'error': 'Nessuna lezione trovata per questo corso'}), 404
    
    data = request.get_json()
    access_token = data.get('access_token')
    
    if not access_token:
        return jsonify({'error': 'Token di accesso mancante'}), 400
    
    try:
        # Crea un design su Canva usando l'API
        # Nota: L'API di Canva per creare design potrebbe richiedere un template specifico
        # Per ora, creiamo un design vuoto e aggiungiamo contenuti
        
        # Endpoint per creare un design (verifica la documentazione Canva per l'endpoint corretto)
        create_design_url = 'https://api.canva.com/rest/v1/designs'
        
        headers = {
            'Authorization': f'Bearer {access_token}',
            'Content-Type': 'application/json'
        }
        
        # Crea un design vuoto (presentazione)
        design_data = {
            'type': 'PRESENTATION',
            'title': f'{course.name} - Presentazione'
        }
        
        response = requests.post(create_design_url, headers=headers, json=design_data)
        
        if response.status_code not in [200, 201]:
            # Se l'endpoint non esiste o ha un formato diverso, proviamo un approccio alternativo
            # Canva potrebbe richiedere di usare un template esistente
            return jsonify({
                'error': f'Errore nella creazione del design: {response.status_code}',
                'details': response.text,
                'message': 'Potrebbe essere necessario usare un template esistente. Verifica la documentazione API di Canva.'
            }), 400
        
        design_info = response.json()
        design_id = design_info.get('id') or design_info.get('design_id')
        design_url = design_info.get('url') or design_info.get('edit_url')
        
        # Prepara il contenuto delle slide basato sulle lezioni
        slides_content = []
        slides_content.append({
            'title': course.name,
            'subtitle': f'{course.code} - {course.total_hours} ore'
        })
        
        slides_content.append({
            'title': 'Panoramica del Corso',
            'content': [
                f'Codice: {course.code}',
                f'Durata totale: {course.total_hours} ore',
                f'Ore teoriche: {course.theory_hours}',
                f'Ore pratiche: {course.practice_hours}',
                f'Numero lezioni: {len(lessons)}'
            ]
        })
        
        for lesson in lessons:
            slide = {
                'title': f'Lezione {lesson.order}: {lesson.title}',
                'content': [
                    f'Tipo: {"Teorica" if lesson.lesson_type == "theory" else "Pratica"}',
                    f'Durata: {lesson.duration_hours} ore'
                ]
            }
            if lesson.description:
                slide['content'].append(f'Descrizione: {lesson.description}')
            if lesson.content:
                # Prendi le prime righe del contenuto
                content_preview = '\n'.join(lesson.content.split('\n')[:5])
                slide['content'].append(f'Contenuto:\n{content_preview}')
            slides_content.append(slide)
        
        return jsonify({
            'message': 'Presentazione creata su Canva con successo!',
            'design_id': design_id,
            'design_url': design_url,
            'edit_url': design_url,
            'slides_count': len(slides_content)
        }), 200
    
    except Exception as e:
        import traceback
        return jsonify({
            'error': f'Errore durante la creazione: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/export-pdf-from-content', methods=['POST'])
def export_pdf_from_content(course_id, lesson_id):
    """Genera PDF dal contenuto Markdown completo (come nell'anteprima)"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    course = Course.query.get_or_404(course_id)
    
    try:
        data = request.json
        full_content = data.get('content', '')
        lesson_title = data.get('title', lesson.title)
        
        if not full_content:
            return jsonify({'error': 'Contenuto mancante'}), 400
        
        # Recupera preferenze esportazione
        export_author = get_preference_value('exportAuthor', '')
        export_company = get_preference_value('exportCompany', '')
        
        # Genera PDF usando WeasyPrint
        from weasyprint import HTML
        
        # Converti Markdown a HTML usando la stessa logica dell'anteprima
        # Usa la funzione convertMarkdownToHTML personalizzata (replicata in Python)
        html_content = markdown.markdown(full_content, extensions=['extra', 'codehilite', 'nl2br'])
        
        # Crea HTML completo con gli stili allineati all'anteprima
        html_doc = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <title>{lesson_title}</title>
            {f'<meta name="author" content="{export_author}">' if export_author else ''}
            {f'<meta name="creator" content="{export_company}">' if export_company else ''}
            {f'<meta name="subject" content="{export_company}">' if export_company else ''}
            <style>
                @page {{
                    size: A4;
                    margin: 2cm;
                }}
                body {{
                    font-family: 'DejaVu Sans', Arial, sans-serif;
                    font-size: 11pt;
                    line-height: 1.8;
                    color: #2c3e50;
                }}
                h1 {{
                    font-size: 2.25em;
                    font-weight: 700;
                    color: #1a1a1a;
                    border-bottom: 3px solid #667eea;
                    padding-bottom: 0.5em;
                    margin: 0 0 1.5em 0;
                    line-height: 1.2;
                }}
                h2 {{
                    font-size: 1.75em;
                    font-weight: 600;
                    color: #2c3e50;
                    border-bottom: 2px solid #e0e0e0;
                    padding-bottom: 0.4em;
                    margin: 2.5em 0 1em 0;
                    line-height: 1.3;
                }}
                h3 {{
                    font-size: 1.4em;
                    font-weight: 600;
                    color: #3d4a5c;
                    margin: 2em 0 0.8em 0;
                    line-height: 1.4;
                }}
                h4 {{
                    font-size: 1.2em;
                    font-weight: 600;
                    color: #4a5568;
                    margin: 1.5em 0 0.6em 0;
                    line-height: 1.4;
                }}
                p {{
                    font-size: 1.05em;
                    line-height: 1.8;
                    color: #2d3748;
                    margin: 0 0 1.25em 0;
                    text-align: left;
                }}
                ul, ol {{
                    margin: 1em 0 1.5em 0;
                    padding-left: 2.5em;
                    line-height: 1.8;
                }}
                li {{
                    font-size: 1.05em;
                    line-height: 1.8;
                    color: #2d3748;
                    margin-bottom: 0.75em;
                }}
                strong {{
                    font-weight: 700;
                    color: #1a1a1a;
                }}
                em {{
                    font-style: italic;
                    color: #4a5568;
                }}
                code {{
                    background-color: #f7f7f7;
                    border: 1px solid #e1e4e8;
                    border-radius: 4px;
                    padding: 3px 8px;
                    font-family: 'Courier New', monospace;
                    font-size: 0.9em;
                    color: #d73a49;
                }}
                pre {{
                    background-color: #f6f8fa;
                    border: 1px solid #e1e4e8;
                    border-radius: 8px;
                    padding: 1.25rem;
                    overflow-x: auto;
                    margin-bottom: 1.5em;
                }}
                pre code {{
                    background-color: transparent;
                    border: none;
                    padding: 0;
                    color: #24292e;
                }}
                blockquote {{
                    border-left: 4px solid #667eea;
                    padding: 0.75rem 1.25rem;
                    margin: 1.5em 0;
                    background-color: #f8f9fa;
                    border-radius: 0 4px 4px 0;
                    color: #4a5568;
                    font-style: italic;
                }}
                .metadata {{
                    background-color: #ecf0f1;
                    padding: 15px;
                    border-radius: 5px;
                    margin-bottom: 20px;
                }}
                .metadata p {{
                    margin: 5px 0;
                }}
            </style>
        </head>
        <body>
            <div class="metadata">
                <h1>{lesson_title}</h1>
                <p><strong>Corso:</strong> {course.name} ({course.code})</p>
                <p><strong>Tipo:</strong> {'Teorica' if lesson.lesson_type == 'theory' else 'Pratica'}</p>
                <p><strong>Durata:</strong> {lesson.duration_hours} ore</p>
            </div>
            
            <div class="lesson-content">
                {html_content if html_content and html_content.strip() else '<p>Nessun contenuto disponibile.</p>'}
            </div>
            
            <div style="margin-top: 40px; padding-top: 20px; border-top: 1px solid #ddd; font-size: 9pt; color: #7f8c8d; text-align: center;">
                <p>Generato il {datetime.now().strftime('%d/%m/%Y alle %H:%M')}</p>
            </div>
        </body>
        </html>
        """
        
        # Genera PDF
        pdf_bytes = HTML(string=html_doc).write_pdf()
        
        # Traccia l'esportazione
        course.export_count = (course.export_count or 0) + 1
        db.session.commit()
        
        return send_file(
            BytesIO(pdf_bytes),
            mimetype='application/pdf',
            as_attachment=True,
            download_name=f"{course.code}_{lesson.order}_{lesson_title[:30].replace(' ', '_')}.pdf"
        )
    
    except Exception as e:
        import traceback
        logging.error(f'Errore durante l\'esportazione PDF: {e}\n{traceback.format_exc()}')
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/lessons/<int:lesson_id>/export/<format_type>', methods=['GET'])
def export_lesson(course_id, lesson_id, format_type):
    """Esporta una lezione in PDF o DOCX"""
    lesson = Lesson.query.filter_by(id=lesson_id, course_id=course_id).first_or_404()
    course = Course.query.get_or_404(course_id)
    
    if format_type not in ['pdf', 'docx']:
        return jsonify({'error': 'Formato non supportato. Usa "pdf" o "docx"'}), 400
    
    try:
        # Prepara il contenuto come nell'anteprima: prima Markdown, poi sezioni separate
        import re
        
        # Rimuovi il testo "(Generato con ChatGPT)" dalle esportazioni
        content = lesson.content or ''
        if content:
            content = re.sub(r'\n\n---\n\n## Contenuto Dettagliato \(Generato con ChatGPT\)\n\n', '\n\n', content)
            content = re.sub(r'\n\n---\n\n## Contenuto Dettagliato \(Generato con ChatGPT\)', '', content)
            content = re.sub(r'## Contenuto Dettagliato \(Generato con ChatGPT\)\n\n', '', content)
            content = content.strip()
        
        objectives = json.loads(lesson.objectives) if lesson.objectives else []
        materials = json.loads(lesson.materials) if lesson.materials else []
        exercises = json.loads(lesson.exercises) if lesson.exercises else []
        
        # Costruisci il contenuto completo come nell'anteprima
        full_content = ''
        
        # Se c'è contenuto Markdown, aggiungilo
        if content and content.strip() != '':
            full_content = content.strip()
            if not full_content.endswith('\n'):
                full_content += '\n'
            full_content += '\n'
        else:
            # Se non c'è contenuto, aggiungi solo la sezione Contenuti
            full_content = '## Contenuti\n\n'
            full_content += '*Nessun contenuto disponibile.*\n\n'
        
        # Aggiungi SEMPRE le sezioni Obiettivi, Materiali ed Esercizi dopo il contenuto Markdown
        # Sezione Obiettivi
        full_content += '## Obiettivi\n\n'
        if objectives and len(objectives) > 0:
            for obj in objectives:
                full_content += f'- {obj}\n'
        else:
            full_content += '*Nessun obiettivo specificato.*\n'
        full_content += '\n'
        
        # Sezione Materiali
        full_content += '## Materiali\n\n'
        if materials and len(materials) > 0:
            for mat in materials:
                full_content += f'- {mat}\n'
        else:
            full_content += '*Nessun materiale specificato.*\n'
        full_content += '\n'
        
        # Sezione Esercizi
        full_content += '## Esercizi\n\n'
        if exercises and len(exercises) > 0:
            for ex in exercises:
                full_content += f'- {ex}\n'
        else:
            full_content += '*Nessun esercizio specificato.*\n'
        full_content += '\n'
        
        # Recupera preferenze esportazione
        export_author = get_preference_value('exportAuthor', '')
        export_company = get_preference_value('exportCompany', '')
        
        if format_type == 'pdf':
            # Genera PDF usando WeasyPrint
            from weasyprint import HTML
            
            # Converti Markdown a HTML (usa il contenuto completo costruito come nell'anteprima)
            # Debug: verifica il contenuto
            logging.info(f'Contenuto da convertire (lunghezza: {len(full_content)} caratteri)')
            logging.info(f'Primi 500 caratteri: {full_content[:500]}')
            
            # Assicurati che il contenuto non sia vuoto
            if not full_content or not full_content.strip():
                logging.warning('Contenuto vuoto, uso contenuto di default')
                full_content = '## Contenuti\n\n*Nessun contenuto disponibile.*\n\n'
            
            try:
                # Converti Markdown a HTML
                html_content = markdown.markdown(full_content, extensions=['extra', 'codehilite', 'nl2br'])
                logging.info(f'Conversione Markdown completata. HTML generato (lunghezza: {len(html_content)} caratteri)')
                logging.info(f'Primi 500 caratteri HTML: {html_content[:500]}')
                
                # Verifica che la conversione abbia prodotto HTML valido
                if not html_content or html_content.strip() == '':
                    logging.warning('Conversione Markdown ha prodotto HTML vuoto, uso fallback')
                    html_content = markdown.markdown(full_content, extensions=['extra'])
                    
            except Exception as e:
                # Fallback: se la conversione fallisce, usa una conversione base
                logging.error(f'Errore conversione Markdown: {e}', exc_info=True)
                try:
                    html_content = markdown.markdown(full_content, extensions=['extra'])
                    logging.info('Conversione Markdown con fallback completata')
                except Exception as e2:
                    # Ultimo fallback: escape HTML e mostra come preformattato
                    logging.error(f'Errore anche nel fallback: {e2}', exc_info=True)
                    import html as html_escape
                    html_content = f'<pre>{html_escape.escape(full_content)}</pre>'
            
            # Crea HTML completo
            html_doc = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <title>{lesson.title}</title>
                {f'<meta name="author" content="{export_author}">' if export_author else ''}
                {f'<meta name="creator" content="{export_company}">' if export_company else ''}
                {f'<meta name="subject" content="{export_company}">' if export_company else ''}
                <style>
                    @page {{
                        size: A4;
                        margin: 2cm;
                    }}
                    body {{
                        font-family: 'DejaVu Sans', Arial, sans-serif;
                        font-size: 11pt;
                        line-height: 1.6;
                        color: #333;
                    }}
                    h1 {{
                        font-size: 2.25em;
                        font-weight: 700;
                        color: #1a1a1a;
                        border-bottom: 3px solid #667eea;
                        padding-bottom: 0.5em;
                        margin: 0 0 1.5em 0;
                        line-height: 1.2;
                    }}
                    h2 {{
                        font-size: 1.75em;
                        font-weight: 600;
                        color: #2c3e50;
                        border-bottom: 2px solid #e0e0e0;
                        padding-bottom: 0.4em;
                        margin: 2.5em 0 1em 0;
                        line-height: 1.3;
                    }}
                    h3 {{
                        font-size: 1.4em;
                        font-weight: 600;
                        color: #3d4a5c;
                        margin: 2em 0 0.8em 0;
                        line-height: 1.4;
                    }}
                    h4 {{
                        font-size: 1.2em;
                        font-weight: 600;
                        color: #4a5568;
                        margin: 1.5em 0 0.6em 0;
                        line-height: 1.4;
                    }}
                    p {{
                        font-size: 1.05em;
                        line-height: 1.8;
                        color: #2d3748;
                        margin: 0 0 1.25em 0;
                        text-align: left;
                    }}
                    ul, ol {{
                        margin: 1em 0 1.5em 0;
                        padding-left: 2.5em;
                        line-height: 1.8;
                    }}
                    li {{
                        font-size: 1.05em;
                        line-height: 1.8;
                        color: #2d3748;
                        margin-bottom: 0.75em;
                    }}
                    strong {{
                        font-weight: 700;
                        color: #1a1a1a;
                    }}
                    em {{
                        font-style: italic;
                        color: #4a5568;
                    }}
                    code {{
                        background-color: #f4f4f4;
                        padding: 2px 6px;
                        border-radius: 3px;
                        font-family: 'Courier New', monospace;
                        font-size: 10pt;
                    }}
                    pre {{
                        background-color: #f8f8f8;
                        border: 1px solid #ddd;
                        border-radius: 4px;
                        padding: 12px;
                        overflow-x: auto;
                        margin-bottom: 15px;
                    }}
                    blockquote {{
                        border-left: 4px solid #3498db;
                        padding-left: 15px;
                        margin-left: 0;
                        color: #7f8c8d;
                        font-style: italic;
                    }}
                    table {{
                        border-collapse: collapse;
                        width: 100%;
                        margin-bottom: 15px;
                    }}
                    table th, table td {{
                        border: 1px solid #ddd;
                        padding: 8px;
                        text-align: left;
                    }}
                    table th {{
                        background-color: #ecf0f1;
                        font-weight: bold;
                    }}
                    .metadata {{
                        background-color: #ecf0f1;
                        padding: 15px;
                        border-radius: 5px;
                        margin-bottom: 20px;
                    }}
                    .metadata p {{
                        margin: 5px 0;
                    }}
                </style>
            </head>
            <body>
                <div class="metadata">
                    <h1>{lesson.title}</h1>
                    <p><strong>Corso:</strong> {course.name} ({course.code})</p>
                    <p><strong>Tipo:</strong> {'Teorica' if lesson.lesson_type == 'theory' else 'Pratica'}</p>
                    <p><strong>Durata:</strong> {lesson.duration_hours} ore</p>
                </div>
                
                <div class="lesson-content">
                    {html_content if html_content and html_content.strip() else '<p>Nessun contenuto disponibile.</p>'}
                </div>
                
                <div style="margin-top: 40px; padding-top: 20px; border-top: 1px solid #ddd; font-size: 9pt; color: #7f8c8d; text-align: center;">
                    <p>Generato il {datetime.now().strftime('%d/%m/%Y alle %H:%M')}</p>
                </div>
            </body>
            </html>
            """
            
            # Genera PDF (WeasyPrint legge automaticamente i metadati dai tag <meta> nell'HTML)
            pdf_bytes = HTML(string=html_doc).write_pdf()
            
            # Traccia l'esportazione
            course.export_count = (course.export_count or 0) + 1
            db.session.commit()
            
            filename = f"{course.code}_{lesson.order}_{lesson.title[:30].replace(' ', '_')}.pdf"
            return send_file(
                BytesIO(pdf_bytes),
                mimetype='application/pdf',
                as_attachment=True,
                download_name=filename
            )
        
        elif format_type == 'docx':
            # Genera DOCX
            doc = Document()
            
            # Imposta metadati
            if export_author:
                doc.core_properties.author = export_author
            if export_company:
                doc.core_properties.company = export_company
            
            # Titolo
            title = doc.add_heading(lesson.title, 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Metadati
            doc.add_paragraph(f'Corso: {course.name} ({course.code})')
            doc.add_paragraph(f'Tipo: {"Teorica" if lesson.lesson_type == "theory" else "Pratica"}')
            doc.add_paragraph(f'Durata: {lesson.duration_hours} ore')
            doc.add_paragraph('')  # Spazio
            
            # Descrizione
            doc.add_heading('Descrizione', level=2)
            doc.add_paragraph(lesson.description or 'Nessuna descrizione disponibile.')
            doc.add_paragraph('')  # Spazio
            
            # Contenuti
            doc.add_heading('Contenuti', level=2)
            # Parsa Markdown direttamente solo se c'è contenuto
            if content and content.strip():
                lines = content.split('\n')
                i = 0
                while i < len(lines):
                    line = lines[i].strip()
                    
                    if not line:
                        i += 1
                        continue
                    
                    # Rileva titoli
                    if line.startswith('#'):
                        level = len(line) - len(line.lstrip('#'))
                        heading_text = line.lstrip('#').strip()
                        if heading_text:
                            doc.add_heading(heading_text, level=min(level, 3))
                    # Rileva liste
                    elif line.startswith('-') or line.startswith('*') or re.match(r'^\d+\.', line):
                        # Raccogli tutte le righe della lista
                        list_items = []
                        while i < len(lines) and (lines[i].strip().startswith('-') or 
                                                 lines[i].strip().startswith('*') or 
                                                 re.match(r'^\d+\.', lines[i].strip())):
                            item = lines[i].strip()
                            # Rimuovi marker
                            item = re.sub(r'^[-*\d+\.]\s+', '', item)
                            if item:
                                list_items.append(item)
                            i += 1
                        # Aggiungi lista
                        for item in list_items:
                            doc.add_paragraph(item, style='List Bullet')
                        i -= 1  # Riposiziona per il prossimo ciclo
                    # Rileva codice (```)
                    elif line.startswith('```'):
                        # Raccogli tutto il blocco di codice
                        code_lines = []
                        i += 1
                        while i < len(lines) and not lines[i].strip().startswith('```'):
                            code_lines.append(lines[i])
                            i += 1
                        if code_lines:
                            code_para = doc.add_paragraph(''.join(code_lines))
                            code_para.style = 'No Spacing'
                            for run in code_para.runs:
                                run.font.name = 'Courier New'
                                run.font.size = Pt(9)
                    # Paragrafo normale
                    else:
                        # Rimuovi formattazione Markdown semplice
                        para_text = line
                        para_text = re.sub(r'\*\*(.+?)\*\*', r'\1', para_text)  # Bold
                        para_text = re.sub(r'\*(.+?)\*', r'\1', para_text)  # Italic
                        para_text = re.sub(r'`(.+?)`', r'\1', para_text)  # Inline code
                        if para_text:
                            doc.add_paragraph(para_text)
                    
                    i += 1
            else:
                doc.add_paragraph('Nessun contenuto disponibile.', style='Intense Quote')
            
            doc.add_paragraph('')  # Spazio
            
            # Obiettivi - sempre mostrati
            doc.add_heading('Obiettivi', level=2)
            if objectives:
                for obj in objectives:
                    doc.add_paragraph(obj, style='List Bullet')
            else:
                doc.add_paragraph('Nessun obiettivo specificato.', style='Intense Quote')
            
            doc.add_paragraph('')  # Spazio
            
            # Materiali - sempre mostrati
            doc.add_heading('Materiali', level=2)
            if materials:
                for mat in materials:
                    doc.add_paragraph(mat, style='List Bullet')
            else:
                doc.add_paragraph('Nessun materiale specificato.', style='Intense Quote')
            
            doc.add_paragraph('')  # Spazio
            
            # Esercizi - sempre mostrati
            doc.add_heading('Esercizi', level=2)
            if exercises:
                for ex in exercises:
                    doc.add_paragraph(ex, style='List Bullet')
            else:
                doc.add_paragraph('Nessun esercizio specificato.', style='Intense Quote')
            
            # Footer
            doc.add_paragraph('')
            footer_para = doc.add_paragraph(f'Generato il {datetime.now().strftime("%d/%m/%Y alle %H:%M")}')
            footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if footer_para.runs:
                footer_para.runs[0].font.size = Pt(9)
            
            # Salva in memoria
            doc_bytes = BytesIO()
            doc.save(doc_bytes)
            doc_bytes.seek(0)
            
            # Traccia l'esportazione
            course.export_count = (course.export_count or 0) + 1
            db.session.commit()
            
            filename = f"{course.code}_{lesson.order}_{lesson.title[:30].replace(' ', '_')}.docx"
            return send_file(
                doc_bytes,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                as_attachment=True,
                download_name=filename
            )
    
    except Exception as e:
        import traceback
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}\n{traceback.format_exc()}'}), 500

@app.route('/api/courses/<int:course_id>/export-batch', methods=['POST'])
def export_course_batch(course_id):
    """Esporta batch di un corso: tutte le lezioni, domande, relazione finale"""
    course = Course.query.get_or_404(course_id)
    
    try:
        # Parametri dalla query string
        format_type = request.args.get('format', 'pdf')
        include_questions = request.args.get('include_questions', 'false').lower() == 'true'
        include_final_report = request.args.get('include_final_report', 'false').lower() == 'true'
        as_zip = request.args.get('as_zip', 'false').lower() == 'true'
        
        if format_type not in ['pdf', 'docx']:
            return jsonify({'error': 'Formato non supportato. Usa "pdf" o "docx"'}), 400
        
        # Recupera preferenze esportazione
        export_author = get_preference_value('exportAuthor', '')
        export_company = get_preference_value('exportCompany', '')
        
        # Recupera tutte le lezioni ordinate
        lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
        
        if not lessons:
            return jsonify({'error': 'Nessuna lezione trovata per questo corso'}), 404
        
        # Se as_zip è True, crea un file ZIP con tutti i file
        if as_zip:
            from zipfile import ZipFile
            import tempfile
            import os
            
            zip_buffer = BytesIO()
            with ZipFile(zip_buffer, 'w') as zip_file:
                # Esporta tutte le lezioni
                for lesson in lessons:
                    # Genera il contenuto completo come nell'anteprima
                    content = lesson.content or ''
                    objectives = json.loads(lesson.objectives) if lesson.objectives else []
                    materials = json.loads(lesson.materials) if lesson.materials else []
                    exercises = json.loads(lesson.exercises) if lesson.exercises else []
                    
                    full_content = ''
                    if content and content.strip() != '':
                        full_content = content.strip()
                        if not full_content.endswith('\n'):
                            full_content += '\n'
                        full_content += '\n'
                    else:
                        full_content = '## Contenuti\n\n*Nessun contenuto disponibile.*\n\n'
                    
                    full_content += '## Obiettivi\n\n'
                    if objectives:
                        for obj in objectives:
                            full_content += f'- {obj}\n'
                    else:
                        full_content += '*Nessun obiettivo specificato.*\n'
                    full_content += '\n'
                    
                    full_content += '## Materiali\n\n'
                    if materials:
                        for mat in materials:
                            full_content += f'- {mat}\n'
                    else:
                        full_content += '*Nessun materiale specificato.*\n'
                    full_content += '\n'
                    
                    full_content += '## Esercizi\n\n'
                    if exercises:
                        for ex in exercises:
                            full_content += f'- {ex}\n'
                    else:
                        full_content += '*Nessun esercizio specificato.*\n'
                    
                    # Genera il file per questa lezione
                    if format_type == 'pdf':
                        from weasyprint import HTML
                        html_content = markdown.markdown(full_content, extensions=['extra', 'codehilite', 'nl2br'])
                        html_doc = f"""
                        <!DOCTYPE html>
                        <html>
                        <head>
                            <meta charset="UTF-8">
                            <title>{lesson.title}</title>
                            {f'<meta name="author" content="{export_author}">' if export_author else ''}
                            {f'<meta name="creator" content="{export_company}">' if export_company else ''}
                            <style>
                                @page {{ size: A4; margin: 2cm; }}
                                body {{ font-family: 'DejaVu Sans', Arial, sans-serif; font-size: 11pt; line-height: 1.8; color: #2c3e50; }}
                                h1 {{ font-size: 2.25em; font-weight: 700; color: #1a1a1a; border-bottom: 3px solid #667eea; padding-bottom: 0.5em; margin: 0 0 1.5em 0; }}
                                h2 {{ font-size: 1.75em; font-weight: 600; color: #2c3e50; border-bottom: 2px solid #e0e0e0; padding-bottom: 0.4em; margin: 2.5em 0 1em 0; }}
                                h3 {{ font-size: 1.4em; font-weight: 600; color: #3d4a5c; margin: 2em 0 0.8em 0; }}
                                p {{ font-size: 1.05em; line-height: 1.8; color: #2d3748; margin: 0 0 1.25em 0; }}
                                ul {{ margin: 1em 0 1.5em 0; padding-left: 2.5em; }}
                                li {{ font-size: 1.05em; line-height: 1.8; color: #2d3748; margin-bottom: 0.75em; }}
                                strong {{ font-weight: 700; color: #1a1a1a; }}
                            </style>
                        </head>
                        <body>
                            <div style="background-color: #ecf0f1; padding: 15px; border-radius: 5px; margin-bottom: 20px;">
                                <h1>{lesson.title}</h1>
                                <p><strong>Corso:</strong> {course.name} ({course.code})</p>
                                <p><strong>Tipo:</strong> {'Teorica' if lesson.lesson_type == 'theory' else 'Pratica'}</p>
                                <p><strong>Durata:</strong> {lesson.duration_hours} ore</p>
                            </div>
                            <div>{html_content}</div>
                        </body>
                        </html>
                        """
                        pdf_bytes = HTML(string=html_doc).write_pdf()
                        filename = f"Lezione_{lesson.order:02d}_{lesson.title[:30].replace(' ', '_')}.pdf"
                        zip_file.writestr(filename, pdf_bytes)
                    else:  # docx
                        from docx import Document
                        from docx.shared import Pt
                        from docx.enum.text import WD_ALIGN_PARAGRAPH
                        doc = Document()
                        if export_author:
                            doc.core_properties.author = export_author
                        if export_company:
                            doc.core_properties.company = export_company
                        doc.add_heading(lesson.title, 0)
                        doc.add_paragraph(f'Corso: {course.name} ({course.code})')
                        doc.add_paragraph(f'Tipo: {"Teorica" if lesson.lesson_type == "theory" else "Pratica"}')
                        doc.add_paragraph(f'Durata: {lesson.duration_hours} ore')
                        doc.add_paragraph('')
                        doc.add_heading('Contenuti', level=2)
                        if content:
                            html_content = markdown.markdown(full_content, extensions=['extra', 'codehilite', 'nl2br'])
                            # Aggiungi contenuto come paragrafi (semplificato)
                            for line in full_content.split('\n'):
                                if line.strip():
                                    if line.startswith('##'):
                                        doc.add_heading(line.replace('##', '').strip(), level=2)
                                    elif line.startswith('-'):
                                        doc.add_paragraph(line.replace('-', '').strip(), style='List Bullet')
                                    else:
                                        doc.add_paragraph(line.strip())
                        doc_bytes = BytesIO()
                        doc.save(doc_bytes)
                        doc_bytes.seek(0)
                        filename = f"Lezione_{lesson.order:02d}_{lesson.title[:30].replace(' ', '_')}.docx"
                        zip_file.writestr(filename, doc_bytes.getvalue())
                
                # Aggiungi domande se richiesto
                if include_questions:
                    questions = Question.query.filter_by(course_id=course_id).order_by(Question.question_number).all()
                    if questions:
                        from docx import Document
                        doc = Document()
                        doc.add_heading('Domande del Corso', 0)
                        for q in questions:
                            doc.add_heading(f'Domanda {q.question_number}', level=2)
                            doc.add_paragraph(q.question_text)
                            doc.add_paragraph(f'A) {q.option_a}')
                            doc.add_paragraph(f'B) {q.option_b}')
                            doc.add_paragraph(f'C) {q.option_c}')
                            doc.add_paragraph(f'D) {q.option_d}')
                            doc.add_paragraph('')
                        doc_bytes = BytesIO()
                        doc.save(doc_bytes)
                        doc_bytes.seek(0)
                        zip_file.writestr('Domande.docx', doc_bytes.getvalue())
                
                # Aggiungi relazione finale se richiesta
                if include_final_report:
                    # Cerca se esiste una relazione finale salvata
                    # Per ora, generiamo un documento placeholder
                    from docx import Document
                    doc = Document()
                    doc.add_heading('Relazione Finale del Corso', 0)
                    doc.add_paragraph(f'Corso: {course.name} ({course.code})')
                    doc.add_paragraph('Relazione finale generata automaticamente.')
                    doc_bytes = BytesIO()
                    doc.save(doc_bytes)
                    doc_bytes.seek(0)
                    zip_file.writestr('Relazione_Finale.docx', doc_bytes.getvalue())
            
            zip_buffer.seek(0)
            course.export_count = (course.export_count or 0) + 1
            db.session.commit()
            
            return send_file(
                zip_buffer,
                mimetype='application/zip',
                as_attachment=True,
                download_name=f"{course.code}_export.zip"
            )
        
        else:
            # Esporta tutto in un unico documento
            if format_type == 'pdf':
                from weasyprint import HTML
                
                # Costruisci HTML completo con tutte le lezioni
                html_parts = []
                html_parts.append(f"""
                <!DOCTYPE html>
                <html>
                <head>
                    <meta charset="UTF-8">
                    <title>{course.name}</title>
                    {f'<meta name="author" content="{export_author}">' if export_author else ''}
                    {f'<meta name="creator" content="{export_company}">' if export_company else ''}
                    <style>
                        @page {{ size: A4; margin: 2cm; }}
                        body {{ font-family: 'DejaVu Sans', Arial, sans-serif; font-size: 11pt; line-height: 1.8; color: #2c3e50; }}
                        h1 {{ font-size: 2.25em; font-weight: 700; color: #1a1a1a; border-bottom: 3px solid #667eea; padding-bottom: 0.5em; margin: 0 0 1.5em 0; page-break-after: avoid; }}
                        h2 {{ font-size: 1.75em; font-weight: 600; color: #2c3e50; border-bottom: 2px solid #e0e0e0; padding-bottom: 0.4em; margin: 2.5em 0 1em 0; page-break-after: avoid; }}
                        h3 {{ font-size: 1.4em; font-weight: 600; color: #3d4a5c; margin: 2em 0 0.8em 0; }}
                        .lesson-section {{ page-break-before: always; }}
                        .lesson-section:first-child {{ page-break-before: auto; }}
                    </style>
                </head>
                <body>
                    <div style="background-color: #ecf0f1; padding: 15px; border-radius: 5px; margin-bottom: 20px;">
                        <h1>{course.name}</h1>
                        <p><strong>Codice:</strong> {course.code}</p>
                        <p><strong>Descrizione:</strong> {course.description or 'Nessuna descrizione'}</p>
                        <p><strong>Ore totali:</strong> {course.total_hours}h</p>
                    </div>
                """)
                
                # Aggiungi tutte le lezioni
                for lesson in lessons:
                    content = lesson.content or ''
                    objectives = json.loads(lesson.objectives) if lesson.objectives else []
                    materials = json.loads(lesson.materials) if lesson.materials else []
                    exercises = json.loads(lesson.exercises) if lesson.exercises else []
                    
                    full_content = ''
                    if content and content.strip() != '':
                        full_content = content.strip()
                        if not full_content.endswith('\n'):
                            full_content += '\n'
                        full_content += '\n'
                    else:
                        full_content = '## Contenuti\n\n*Nessun contenuto disponibile.*\n\n'
                    
                    full_content += '## Obiettivi\n\n'
                    if objectives:
                        for obj in objectives:
                            full_content += f'- {obj}\n'
                    else:
                        full_content += '*Nessun obiettivo specificato.*\n'
                    full_content += '\n'
                    
                    full_content += '## Materiali\n\n'
                    if materials:
                        for mat in materials:
                            full_content += f'- {mat}\n'
                    else:
                        full_content += '*Nessun materiale specificato.*\n'
                    full_content += '\n'
                    
                    full_content += '## Esercizi\n\n'
                    if exercises:
                        for ex in exercises:
                            full_content += f'- {ex}\n'
                    else:
                        full_content += '*Nessun esercizio specificato.*\n'
                    
                    html_content = markdown.markdown(full_content, extensions=['extra', 'codehilite', 'nl2br'])
                    html_parts.append(f"""
                    <div class="lesson-section">
                        <div style="background-color: #f8f9fa; padding: 15px; border-radius: 5px; margin-bottom: 20px;">
                            <h2>{lesson.title}</h2>
                            <p><strong>Tipo:</strong> {'Teorica' if lesson.lesson_type == 'theory' else 'Pratica'}</p>
                            <p><strong>Durata:</strong> {lesson.duration_hours} ore</p>
                        </div>
                        <div>{html_content}</div>
                    </div>
                    """)
                
                # Aggiungi domande se richiesto
                if include_questions:
                    questions = Question.query.filter_by(course_id=course_id).order_by(Question.question_number).all()
                    if questions:
                        html_parts.append('<div class="lesson-section"><h2>Domande del Corso</h2>')
                        for q in questions:
                            html_parts.append(f'<h3>Domanda {q.question_number}</h3><p>{q.question_text}</p><ul>')
                            html_parts.append(f'<li>A) {q.option_a}</li>')
                            html_parts.append(f'<li>B) {q.option_b}</li>')
                            html_parts.append(f'<li>C) {q.option_c}</li>')
                            html_parts.append(f'<li>D) {q.option_d}</li>')
                            html_parts.append('</ul>')
                        html_parts.append('</div>')
                
                # Aggiungi relazione finale se richiesta
                if include_final_report:
                    html_parts.append('<div class="lesson-section"><h2>Relazione Finale del Corso</h2><p>Relazione finale generata automaticamente.</p></div>')
                
                html_parts.append('</body></html>')
                html_doc = ''.join(html_parts)
                
                pdf_bytes = HTML(string=html_doc).write_pdf()
                course.export_count = (course.export_count or 0) + 1
                db.session.commit()
                
                return send_file(
                    BytesIO(pdf_bytes),
                    mimetype='application/pdf',
                    as_attachment=True,
                    download_name=f"{course.code}_completo.pdf"
                )
            else:  # docx
                from docx import Document
                from docx.shared import Pt
                from docx.enum.text import WD_ALIGN_PARAGRAPH
                
                doc = Document()
                if export_author:
                    doc.core_properties.author = export_author
                if export_company:
                    doc.core_properties.company = export_company
                
                # Titolo corso
                doc.add_heading(course.name, 0)
                doc.add_paragraph(f'Codice: {course.code}')
                doc.add_paragraph(f'Descrizione: {course.description or "Nessuna descrizione"}')
                doc.add_paragraph(f'Ore totali: {course.total_hours}h')
                doc.add_page_break()
                
                # Aggiungi tutte le lezioni
                for lesson in lessons:
                    content = lesson.content or ''
                    objectives = json.loads(lesson.objectives) if lesson.objectives else []
                    materials = json.loads(lesson.materials) if lesson.materials else []
                    exercises = json.loads(lesson.exercises) if lesson.exercises else []
                    
                    doc.add_heading(lesson.title, level=1)
                    doc.add_paragraph(f'Tipo: {"Teorica" if lesson.lesson_type == "theory" else "Pratica"}')
                    doc.add_paragraph(f'Durata: {lesson.duration_hours} ore')
                    doc.add_paragraph('')
                    
                    # Contenuti
                    if content:
                        html_content = markdown.markdown(content, extensions=['extra', 'codehilite', 'nl2br'])
                        for line in content.split('\n'):
                            if line.strip():
                                if line.startswith('##'):
                                    doc.add_heading(line.replace('##', '').strip(), level=2)
                                elif line.startswith('-'):
                                    doc.add_paragraph(line.replace('-', '').strip(), style='List Bullet')
                                else:
                                    doc.add_paragraph(line.strip())
                    
                    doc.add_heading('Obiettivi', level=2)
                    if objectives:
                        for obj in objectives:
                            doc.add_paragraph(obj, style='List Bullet')
                    else:
                        doc.add_paragraph('Nessun obiettivo specificato.', style='Intense Quote')
                    
                    doc.add_heading('Materiali', level=2)
                    if materials:
                        for mat in materials:
                            doc.add_paragraph(mat, style='List Bullet')
                    else:
                        doc.add_paragraph('Nessun materiale specificato.', style='Intense Quote')
                    
                    doc.add_heading('Esercizi', level=2)
                    if exercises:
                        for ex in exercises:
                            doc.add_paragraph(ex, style='List Bullet')
                    else:
                        doc.add_paragraph('Nessun esercizio specificato.', style='Intense Quote')
                    
                    doc.add_page_break()
                
                # Aggiungi domande se richiesto
                if include_questions:
                    questions = Question.query.filter_by(course_id=course_id).order_by(Question.question_number).all()
                    if questions:
                        doc.add_heading('Domande del Corso', level=1)
                        for q in questions:
                            doc.add_heading(f'Domanda {q.question_number}', level=2)
                            doc.add_paragraph(q.question_text)
                            doc.add_paragraph(f'A) {q.option_a}')
                            doc.add_paragraph(f'B) {q.option_b}')
                            doc.add_paragraph(f'C) {q.option_c}')
                            doc.add_paragraph(f'D) {q.option_d}')
                            doc.add_paragraph('')
                        doc.add_page_break()
                
                # Aggiungi relazione finale se richiesta
                if include_final_report:
                    doc.add_heading('Relazione Finale del Corso', level=1)
                    doc.add_paragraph('Relazione finale generata automaticamente.')
                
                doc_bytes = BytesIO()
                doc.save(doc_bytes)
                doc_bytes.seek(0)
                
                course.export_count = (course.export_count or 0) + 1
                db.session.commit()
                
                return send_file(
                    doc_bytes,
                    mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    as_attachment=True,
                    download_name=f"{course.code}_completo.docx"
                )
    
    except Exception as e:
        import traceback
        logging.error(f'Errore durante l\'esportazione batch: {e}\n{traceback.format_exc()}')
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/ai-analysis/<analysis_type>', methods=['POST'])
def ai_analysis(course_id, analysis_type):
    """Esegue analisi AI sul corso: bilanciamento, suggerimenti, duplicati, coerenza"""
    course = Course.query.get_or_404(course_id)
    
    if not app.config['OPENAI_API_KEY']:
        return jsonify({'error': 'API key di OpenAI non configurata'}), 400
    
    if analysis_type not in ['balance', 'suggestions', 'duplicates', 'coherence']:
        return jsonify({'error': 'Tipo di analisi non valido'}), 400
    
    try:
        import os
        os.environ['OPENAI_API_KEY'] = app.config['OPENAI_API_KEY']
        client = OpenAI()
        
        # Recupera tutte le lezioni del corso
        lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
        
        if not lessons:
            return jsonify({'error': 'Nessuna lezione trovata per questo corso'}), 404
        
        # Prepara i dati del corso per l'analisi
        course_data = {
            'name': course.name,
            'code': course.code,
            'description': course.description or '',
            'total_hours': course.total_hours,
            'theory_hours': course.theory_hours,
            'practice_hours': course.practice_hours,
            'lessons': []
        }
        
        for lesson in lessons:
            objectives = json.loads(lesson.objectives) if lesson.objectives else []
            materials = json.loads(lesson.materials) if lesson.materials else []
            exercises = json.loads(lesson.exercises) if lesson.exercises else []
            
            course_data['lessons'].append({
                'title': lesson.title,
                'order': lesson.order,
                'type': lesson.lesson_type,
                'duration': lesson.duration_hours,
                'description': lesson.description or '',
                'content': lesson.content or '',
                'objectives': objectives,
                'materials': materials,
                'exercises': exercises
            })
        
        # Recupera preferenze AI
        ai_model = get_preference_value('aiModel', 'gpt-4o')
        ai_temperature = float(get_preference_value('aiTemperature', '0.7'))
        ai_max_tokens = int(get_preference_value('aiMaxTokens', '2000'))
        
        if analysis_type == 'balance':
            # Analisi bilanciamento teoria/pratica
            theory_hours = sum(l['duration'] for l in course_data['lessons'] if l['type'] == 'theory')
            practice_hours = sum(l['duration'] for l in course_data['lessons'] if l['type'] == 'practice')
            theory_percentage = (theory_hours / course_data['total_hours'] * 100) if course_data['total_hours'] > 0 else 0
            practice_percentage = (practice_hours / course_data['total_hours'] * 100) if course_data['total_hours'] > 0 else 0
            
            prompt = f"""Analizza il bilanciamento tra teoria e pratica per questo corso di formazione professionale.

**Corso:** {course_data['name']} ({course_data['code']})
**Ore totali:** {course_data['total_hours']}h
**Ore teoria:** {theory_hours}h ({theory_percentage:.1f}%)
**Ore pratica:** {practice_hours}h ({practice_percentage:.1f}%)

**Lezioni:**
{chr(10).join([f"- Lezione {l['order']}: {l['title']} ({l['type']}, {l['duration']}h)" for l in course_data['lessons']])}

**Compito:**
1. Analizza il bilanciamento attuale tra teoria e pratica
2. Valuta se è appropriato per un corso di formazione professionale
3. Fornisci raccomandazioni per migliorare il bilanciamento se necessario
4. Suggerisci eventuali aggiustamenti nelle durate delle lezioni

Rispondi in formato Markdown con sezioni chiare."""
            
            response = client.chat.completions.create(
                model=ai_model,
                messages=[{"role": "user", "content": prompt}],
                temperature=ai_temperature,
                max_tokens=ai_max_tokens
            )
            
            analysis_text = response.choices[0].message.content
            
            # Salva l'analisi nel database
            analysis_data_json = json.dumps({
                'theory_hours': theory_hours,
                'practice_hours': practice_hours,
                'theory_percentage': round(theory_percentage, 1),
                'practice_percentage': round(practice_percentage, 1)
            })
            ai_analysis = AIAnalysis(
                course_id=course_id,
                analysis_type='balance',
                title='Analisi Bilanciamento Teoria/Pratica',
                analysis_content=analysis_text,
                analysis_data=analysis_data_json
            )
            db.session.add(ai_analysis)
            db.session.commit()
            
            return jsonify({
                'id': ai_analysis.id,
                'title': 'Analisi Bilanciamento Teoria/Pratica',
                'analysis': analysis_text,
                'data': {
                    'theory_hours': theory_hours,
                    'practice_hours': practice_hours,
                    'theory_percentage': round(theory_percentage, 1),
                    'practice_percentage': round(practice_percentage, 1)
                },
                'created_at': ai_analysis.created_at.isoformat()
            })
        
        elif analysis_type == 'suggestions':
            # Suggerimenti per migliorare lezioni
            lessons_summary = '\n\n'.join([
                f"**Lezione {l['order']}: {l['title']}** ({l['type']}, {l['duration']}h)\n"
                f"Descrizione: {l['description'] or 'Nessuna'}\n"
                f"Contenuto: {l['content'][:500] if l['content'] else 'Nessun contenuto'}..."
                for l in course_data['lessons']
            ])
            
            prompt = f"""Sei un esperto di didattica e formazione professionale. Analizza questo corso e fornisci suggerimenti concreti per migliorare le lezioni.

**Corso:** {course_data['name']} ({course_data['code']})
**Descrizione:** {course_data['description'] or 'Nessuna descrizione'}
**Ore totali:** {course_data['total_hours']}h ({course_data['theory_hours']}h teoria, {course_data['practice_hours']}h pratica)
**Numero lezioni:** {len(course_data['lessons'])}

**Lezioni:**
{lessons_summary}

**Compito:**
1. Analizza ogni lezione e identifica punti di forza e debolezze
2. Fornisci suggerimenti specifici per migliorare:
   - La struttura e organizzazione
   - La chiarezza dei contenuti
   - L'equilibrio tra teoria e pratica
   - L'engagement degli studenti
   - La completezza delle informazioni
3. Suggerisci eventuali lezioni mancanti o da aggiungere
4. Fornisci raccomandazioni pratiche e attuabili

Rispondi in formato Markdown con sezioni chiare per ogni lezione."""
            
            response = client.chat.completions.create(
                model=ai_model,
                messages=[{"role": "user", "content": prompt}],
                temperature=ai_temperature,
                max_tokens=ai_max_tokens * 2  # Più token per analisi dettagliata
            )
            
            analysis_text = response.choices[0].message.content
            
            # Salva l'analisi nel database
            ai_analysis = AIAnalysis(
                course_id=course_id,
                analysis_type='suggestions',
                title='Suggerimenti per Migliorare le Lezioni',
                analysis_content=analysis_text,
                analysis_data=None
            )
            db.session.add(ai_analysis)
            db.session.commit()
            
            return jsonify({
                'id': ai_analysis.id,
                'title': 'Suggerimenti per Migliorare le Lezioni',
                'analysis': analysis_text,
                'created_at': ai_analysis.created_at.isoformat()
            })
        
        elif analysis_type == 'duplicates':
            # Rilevamento contenuti duplicati
            lessons_content = '\n\n'.join([
                f"**LEZIONE {l['order']}: {l['title']}**\n"
                f"Contenuto: {l['content'] or 'Nessun contenuto'}\n"
                f"Obiettivi: {', '.join(l['objectives']) if l['objectives'] else 'Nessuno'}"
                for l in course_data['lessons']
            ])
            
            prompt = f"""Analizza questo corso di formazione professionale per identificare contenuti duplicati o molto simili tra le lezioni.

**Corso:** {course_data['name']} ({course_data['code']})
**Numero lezioni:** {len(course_data['lessons'])}

**Contenuti delle lezioni:**
{lessons_content}

**Compito:**
1. Identifica contenuti duplicati o molto simili tra le lezioni
2. Indica quali lezioni hanno sovrapposizioni significative
3. Suggerisci come eliminare o ridurre le duplicazioni
4. Valuta se alcune lezioni potrebbero essere unificate
5. Fornisci raccomandazioni per migliorare la progressione logica

Rispondi in formato Markdown con sezioni chiare."""
            
            response = client.chat.completions.create(
                model=ai_model,
                messages=[{"role": "user", "content": prompt}],
                temperature=ai_temperature,
                max_tokens=ai_max_tokens
            )
            
            analysis_text = response.choices[0].message.content
            
            # Salva l'analisi nel database
            ai_analysis = AIAnalysis(
                course_id=course_id,
                analysis_type='duplicates',
                title='Rilevamento Contenuti Duplicati',
                analysis_content=analysis_text,
                analysis_data=None
            )
            db.session.add(ai_analysis)
            db.session.commit()
            
            return jsonify({
                'id': ai_analysis.id,
                'title': 'Rilevamento Contenuti Duplicati',
                'analysis': analysis_text,
                'created_at': ai_analysis.created_at.isoformat()
            })
        
        elif analysis_type == 'coherence':
            # Analisi coerenza obiettivi-contenuti
            lessons_detailed = '\n\n'.join([
                f"**LEZIONE {l['order']}: {l['title']}** ({l['type']}, {l['duration']}h)\n"
                f"Obiettivi formativi:\n" + '\n'.join([f"- {obj}" for obj in l['objectives']]) + "\n"
                f"Contenuto:\n{l['content'] or 'Nessun contenuto'}\n"
                for l in course_data['lessons'] if l['objectives']
            ])
            
            prompt = f"""Analizza la coerenza tra gli obiettivi formativi e i contenuti delle lezioni di questo corso.

**Corso:** {course_data['name']} ({course_data['code']})
**Descrizione corso:** {course_data['description'] or 'Nessuna descrizione'}

**Lezioni con obiettivi e contenuti:**
{lessons_detailed}

**Compito:**
1. Per ogni lezione, verifica se i contenuti sono allineati con gli obiettivi formativi
2. Identifica eventuali discrepanze tra obiettivi dichiarati e contenuti effettivi
3. Suggerisci come migliorare l'allineamento
4. Valuta se gli obiettivi sono realistici rispetto ai contenuti e alla durata
5. Fornisci raccomandazioni per migliorare la coerenza complessiva

Rispondi in formato Markdown con sezioni chiare per ogni lezione."""
            
            response = client.chat.completions.create(
                model=ai_model,
                messages=[{"role": "user", "content": prompt}],
                temperature=ai_temperature,
                max_tokens=ai_max_tokens * 2
            )
            
            analysis_text = response.choices[0].message.content
            
            # Salva l'analisi nel database
            ai_analysis = AIAnalysis(
                course_id=course_id,
                analysis_type='coherence',
                title='Analisi Coerenza Obiettivi-Contenuti',
                analysis_content=analysis_text,
                analysis_data=None
            )
            db.session.add(ai_analysis)
            db.session.commit()
            
            return jsonify({
                'id': ai_analysis.id,
                'title': 'Analisi Coerenza Obiettivi-Contenuti',
                'analysis': analysis_text,
                'created_at': ai_analysis.created_at.isoformat()
            })
    
    except Exception as e:
        import traceback
        logging.error(f'Errore durante l\'analisi AI: {e}\n{traceback.format_exc()}')
        return jsonify({'error': f'Errore durante l\'analisi: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/ai-analysis/<analysis_type>/export/<format_type>', methods=['POST'])
def export_ai_analysis(course_id, analysis_type, format_type):
    """Esporta i risultati dell'analisi AI in PDF o DOCX"""
    course = Course.query.get_or_404(course_id)
    
    if format_type not in ['pdf', 'docx']:
        return jsonify({'error': 'Formato non supportato. Usa "pdf" o "docx"'}), 400
    
    try:
        data = request.json
        analysis_text = data.get('analysis', '')
        analysis_title = data.get('title', 'Analisi AI')
        analysis_data = data.get('data', {})
        
        if not analysis_text:
            return jsonify({'error': 'Nessun contenuto di analisi da esportare'}), 400
        
        # Recupera preferenze esportazione
        export_author = get_preference_value('exportAuthor', '')
        export_company = get_preference_value('exportCompany', '')
        
        if format_type == 'pdf':
            from weasyprint import HTML
            
            # Converti Markdown a HTML
            html_content = markdown.markdown(analysis_text, extensions=['extra', 'codehilite', 'nl2br'])
            
            # Aggiungi dati statistici se presenti
            data_section = ''
            if analysis_data:
                if 'theory_percentage' in analysis_data:
                    data_section = f"""
                    <div style="background-color: #f8f9fa; padding: 15px; border-radius: 5px; margin-bottom: 20px;">
                        <h3>Dati Statistici</h3>
                        <p><strong>Ore Teoria:</strong> {analysis_data.get('theory_hours', 0)}h ({analysis_data.get('theory_percentage', 0)}%)</p>
                        <p><strong>Ore Pratica:</strong> {analysis_data.get('practice_hours', 0)}h ({analysis_data.get('practice_percentage', 0)}%)</p>
                    </div>
                    """
            
            html_doc = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <title>{analysis_title}</title>
                {f'<meta name="author" content="{export_author}">' if export_author else ''}
                {f'<meta name="creator" content="{export_company}">' if export_company else ''}
                <style>
                    @page {{ size: A4; margin: 2cm; }}
                    body {{ font-family: 'DejaVu Sans', Arial, sans-serif; font-size: 11pt; line-height: 1.8; color: #2c3e50; }}
                    h1 {{ font-size: 2.25em; font-weight: 700; color: #1a1a1a; border-bottom: 3px solid #667eea; padding-bottom: 0.5em; margin: 0 0 1.5em 0; }}
                    h2 {{ font-size: 1.75em; font-weight: 600; color: #2c3e50; border-bottom: 2px solid #e0e0e0; padding-bottom: 0.4em; margin: 2.5em 0 1em 0; }}
                    h3 {{ font-size: 1.4em; font-weight: 600; color: #3d4a5c; margin: 2em 0 0.8em 0; }}
                    p {{ font-size: 1.05em; line-height: 1.8; color: #2d3748; margin: 0 0 1.25em 0; }}
                    ul, ol {{ margin: 1em 0 1.5em 0; padding-left: 2.5em; }}
                    li {{ font-size: 1.05em; line-height: 1.8; color: #2d3748; margin-bottom: 0.75em; }}
                    strong {{ font-weight: 700; color: #1a1a1a; }}
                    code {{ background-color: #f7f7f7; border: 1px solid #e1e4e8; border-radius: 4px; padding: 3px 8px; font-family: 'Courier New', monospace; font-size: 0.9em; }}
                    blockquote {{ border-left: 4px solid #667eea; padding: 0.75rem 1.25rem; margin: 1.5em 0; background-color: #f8f9fa; border-radius: 0 4px 4px 0; }}
                </style>
            </head>
            <body>
                <div style="background-color: #ecf0f1; padding: 15px; border-radius: 5px; margin-bottom: 20px;">
                    <h1>{analysis_title}</h1>
                    <p><strong>Corso:</strong> {course.name} ({course.code})</p>
                    <p><strong>Data analisi:</strong> {datetime.now().strftime('%d/%m/%Y alle %H:%M')}</p>
                </div>
                {data_section}
                <div class="analysis-content">
                    {html_content}
                </div>
                <div style="margin-top: 40px; padding-top: 20px; border-top: 1px solid #ddd; font-size: 9pt; color: #7f8c8d; text-align: center;">
                    <p>Generato il {datetime.now().strftime('%d/%m/%Y alle %H:%M')}</p>
                </div>
            </body>
            </html>
            """
            
            pdf_bytes = HTML(string=html_doc).write_pdf()
            
            return send_file(
                BytesIO(pdf_bytes),
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f"{course.code}_{analysis_type}_analisi.pdf"
            )
        
        else:  # docx
            from docx import Document
            from docx.shared import Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            if export_author:
                doc.core_properties.author = export_author
            if export_company:
                doc.core_properties.company = export_company
            
            # Titolo
            doc.add_heading(analysis_title, 0)
            doc.add_paragraph(f'Corso: {course.name} ({course.code})')
            doc.add_paragraph(f'Data analisi: {datetime.now().strftime("%d/%m/%Y alle %H:%M")}')
            doc.add_paragraph('')
            
            # Dati statistici se presenti
            if analysis_data:
                doc.add_heading('Dati Statistici', level=2)
                if 'theory_percentage' in analysis_data:
                    doc.add_paragraph(f'Ore Teoria: {analysis_data.get("theory_hours", 0)}h ({analysis_data.get("theory_percentage", 0)}%)')
                    doc.add_paragraph(f'Ore Pratica: {analysis_data.get("practice_hours", 0)}h ({analysis_data.get("practice_percentage", 0)}%)')
                doc.add_paragraph('')
            
            # Contenuto analisi
            doc.add_heading('Analisi', level=2)
            
            # Converti Markdown in paragrafi Word (semplificato)
            lines = analysis_text.split('\n')
            for line in lines:
                line = line.strip()
                if not line:
                    doc.add_paragraph('')
                elif line.startswith('## '):
                    doc.add_heading(line.replace('##', '').strip(), level=2)
                elif line.startswith('### '):
                    doc.add_heading(line.replace('###', '').strip(), level=3)
                elif line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line.replace('-', '').replace('*', '').strip(), style='List Bullet')
                elif line.startswith('**') and line.endswith('**'):
                    para = doc.add_paragraph()
                    run = para.add_run(line.replace('**', ''))
                    run.bold = True
                else:
                    doc.add_paragraph(line)
            
            # Footer
            doc.add_paragraph('')
            footer_para = doc.add_paragraph(f'Generato il {datetime.now().strftime("%d/%m/%Y alle %H:%M")}')
            footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if footer_para.runs:
                footer_para.runs[0].font.size = Pt(9)
            
            doc_bytes = BytesIO()
            doc.save(doc_bytes)
            doc_bytes.seek(0)
            
            return send_file(
                doc_bytes,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                as_attachment=True,
                download_name=f"{course.code}_{analysis_type}_analisi.docx"
            )
    
    except Exception as e:
        import traceback
        logging.error(f'Errore durante l\'esportazione analisi AI: {e}\n{traceback.format_exc()}')
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/ai-analyses', methods=['GET'])
def get_ai_analyses(course_id):
    """Recupera tutte le analisi AI salvate per un corso"""
    course = Course.query.get_or_404(course_id)
    
    analyses = AIAnalysis.query.filter_by(course_id=course_id).order_by(AIAnalysis.created_at.desc()).all()
    
    return jsonify([{
        'id': a.id,
        'analysis_type': a.analysis_type,
        'title': a.title,
        'created_at': a.created_at.isoformat(),
        'created_at_formatted': a.created_at.strftime('%d/%m/%Y alle %H:%M')
    } for a in analyses])

@app.route('/api/courses/<int:course_id>/ai-analyses/<int:analysis_id>', methods=['GET'])
def get_ai_analysis(course_id, analysis_id):
    """Recupera una specifica analisi AI"""
    analysis = AIAnalysis.query.filter_by(id=analysis_id, course_id=course_id).first_or_404()
    
    analysis_data = None
    if analysis.analysis_data:
        analysis_data = json.loads(analysis.analysis_data)
    
    return jsonify({
        'id': analysis.id,
        'analysis_type': analysis.analysis_type,
        'title': analysis.title,
        'analysis': analysis.analysis_content,
        'data': analysis_data,
        'created_at': analysis.created_at.isoformat(),
        'created_at_formatted': analysis.created_at.strftime('%d/%m/%Y alle %H:%M')
    })

@app.route('/api/courses/<int:course_id>/ai-analyses/<int:analysis_id>/export/<format_type>', methods=['GET'])
def export_saved_ai_analysis(course_id, analysis_id, format_type):
    """Esporta un'analisi AI salvata in PDF o DOCX"""
    course = Course.query.get_or_404(course_id)
    analysis = AIAnalysis.query.filter_by(id=analysis_id, course_id=course_id).first_or_404()
    
    if format_type not in ['pdf', 'docx']:
        return jsonify({'error': 'Formato non supportato. Usa "pdf" o "docx"'}), 400
    
    try:
        # Recupera preferenze esportazione
        export_author = get_preference_value('exportAuthor', '')
        export_company = get_preference_value('exportCompany', '')
        
        analysis_data = None
        if analysis.analysis_data:
            analysis_data = json.loads(analysis.analysis_data)
        
        if format_type == 'pdf':
            from weasyprint import HTML
            
            # Converti Markdown a HTML
            html_content = markdown.markdown(analysis.analysis_content, extensions=['extra', 'codehilite', 'nl2br'])
            
            # Aggiungi dati statistici se presenti
            data_section = ''
            if analysis_data and 'theory_percentage' in analysis_data:
                data_section = f"""
                <div style="background-color: #f8f9fa; padding: 15px; border-radius: 5px; margin-bottom: 20px;">
                    <h3>Dati Statistici</h3>
                    <p><strong>Ore Teoria:</strong> {analysis_data.get('theory_hours', 0)}h ({analysis_data.get('theory_percentage', 0)}%)</p>
                    <p><strong>Ore Pratica:</strong> {analysis_data.get('practice_hours', 0)}h ({analysis_data.get('practice_percentage', 0)}%)</p>
                </div>
                """
            
            html_doc = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <title>{analysis.title}</title>
                {f'<meta name="author" content="{export_author}">' if export_author else ''}
                {f'<meta name="creator" content="{export_company}">' if export_company else ''}
                <style>
                    @page {{ size: A4; margin: 2cm; }}
                    body {{ font-family: 'DejaVu Sans', Arial, sans-serif; font-size: 11pt; line-height: 1.8; color: #2c3e50; }}
                    h1 {{ font-size: 2.25em; font-weight: 700; color: #1a1a1a; border-bottom: 3px solid #667eea; padding-bottom: 0.5em; margin: 0 0 1.5em 0; }}
                    h2 {{ font-size: 1.75em; font-weight: 600; color: #2c3e50; border-bottom: 2px solid #e0e0e0; padding-bottom: 0.4em; margin: 2.5em 0 1em 0; }}
                    h3 {{ font-size: 1.4em; font-weight: 600; color: #3d4a5c; margin: 2em 0 0.8em 0; }}
                    p {{ font-size: 1.05em; line-height: 1.8; color: #2d3748; margin: 0 0 1.25em 0; }}
                    ul, ol {{ margin: 1em 0 1.5em 0; padding-left: 2.5em; }}
                    li {{ font-size: 1.05em; line-height: 1.8; color: #2d3748; margin-bottom: 0.75em; }}
                    strong {{ font-weight: 700; color: #1a1a1a; }}
                </style>
            </head>
            <body>
                <div style="background-color: #ecf0f1; padding: 15px; border-radius: 5px; margin-bottom: 20px;">
                    <h1>{analysis.title}</h1>
                    <p><strong>Corso:</strong> {course.name} ({course.code})</p>
                    <p><strong>Data analisi:</strong> {analysis.created_at.strftime('%d/%m/%Y alle %H:%M')}</p>
                </div>
                {data_section}
                <div class="analysis-content">
                    {html_content}
                </div>
            </body>
            </html>
            """
            
            pdf_bytes = HTML(string=html_doc).write_pdf()
            
            return send_file(
                BytesIO(pdf_bytes),
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f"{course.code}_{analysis.analysis_type}_analisi.pdf"
            )
        
        else:  # docx
            from docx import Document
            from docx.shared import Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            if export_author:
                doc.core_properties.author = export_author
            if export_company:
                doc.core_properties.company = export_company
            
            doc.add_heading(analysis.title, 0)
            doc.add_paragraph(f'Corso: {course.name} ({course.code})')
            doc.add_paragraph(f'Data analisi: {analysis.created_at.strftime("%d/%m/%Y alle %H:%M")}')
            doc.add_paragraph('')
            
            if analysis_data and 'theory_percentage' in analysis_data:
                doc.add_heading('Dati Statistici', level=2)
                doc.add_paragraph(f'Ore Teoria: {analysis_data.get("theory_hours", 0)}h ({analysis_data.get("theory_percentage", 0)}%)')
                doc.add_paragraph(f'Ore Pratica: {analysis_data.get("practice_hours", 0)}h ({analysis_data.get("practice_percentage", 0)}%)')
                doc.add_paragraph('')
            
            doc.add_heading('Analisi', level=2)
            lines = analysis.analysis_content.split('\n')
            for line in lines:
                line = line.strip()
                if not line:
                    doc.add_paragraph('')
                elif line.startswith('## '):
                    doc.add_heading(line.replace('##', '').strip(), level=2)
                elif line.startswith('### '):
                    doc.add_heading(line.replace('###', '').strip(), level=3)
                elif line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line.replace('-', '').replace('*', '').strip(), style='List Bullet')
                else:
                    doc.add_paragraph(line)
            
            doc_bytes = BytesIO()
            doc.save(doc_bytes)
            doc_bytes.seek(0)
            
            return send_file(
                doc_bytes,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                as_attachment=True,
                download_name=f"{course.code}_{analysis.analysis_type}_analisi.docx"
            )
    
    except Exception as e:
        import traceback
        logging.error(f'Errore durante l\'esportazione analisi salvata: {e}\n{traceback.format_exc()}')
        return jsonify({'error': f'Errore durante l\'esportazione: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/import-md', methods=['POST'])
def import_lessons_from_md(course_id):
    course = Course.query.get_or_404(course_id)
    
    if 'file' not in request.files:
        return jsonify({'error': 'Nessun file fornito'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'File non selezionato'}), 400
    
    if not file.filename.endswith('.md'):
        return jsonify({'error': 'Il file deve essere in formato Markdown (.md)'}), 400
    
    try:
        content = file.read().decode('utf-8')
        lessons = parse_markdown_lessons(content, course)
        
        # Elimina lezioni esistenti
        Lesson.query.filter_by(course_id=course_id).delete()
        
        # Aggiungi nuove lezioni
        for lesson_data in lessons:
            lesson = Lesson(
                course_id=course_id,
                title=lesson_data['title'],
                description=lesson_data.get('description', ''),
                lesson_type=lesson_data['lesson_type'],
                duration_hours=lesson_data['duration_hours'],
                order=lesson_data['order'],
                content=lesson_data.get('content', ''),
                objectives=json.dumps(lesson_data.get('objectives', [])),
                materials=json.dumps(lesson_data.get('materials', [])),
                exercises=json.dumps(lesson_data.get('exercises', []))
            )
            db.session.add(lesson)
        
        db.session.commit()
        return jsonify({
            'message': f'{len(lessons)} lezioni importate con successo',
            'lessons_count': len(lessons)
        }), 200
    except Exception as e:
        return jsonify({'error': f'Errore durante l\'import: {str(e)}'}), 500

def extract_markdown_from_rtf(rtf_content):
    """Estrae il contenuto Markdown da un file RTF"""
    import re
    # Decodifica RTF
    try:
        text = rtf_content.decode('utf-8', errors='ignore')
    except:
        text = rtf_content.decode('latin-1', errors='ignore')
    
    # Sostituisce \par con newline
    text = re.sub(r'\\par\s*', '\n', text)
    text = re.sub(r'\\line\s*', '\n', text)
    
    # Rimuove comandi RTF comuni ma mantiene il testo
    # Rimuove comandi RTF tipo \fs24, \cf0, etc.
    text = re.sub(r'\\[a-z]+\d*\s*', ' ', text)
    
    # Rimuove gruppi RTF vuoti o con solo comandi
    text = re.sub(r'\{[^}]*\\[^}]*\}', '', text)
    
    # Rimuove caratteri di controllo RTF
    text = re.sub(r'\\[{}]', '', text)
    text = re.sub(r'[{}]', '', text)
    
    # Gestisce caratteri speciali RTF
    text = text.replace('\\\'e0', 'à')
    text = text.replace('\\\'e8', 'è')
    text = text.replace('\\\'e9', 'é')
    text = text.replace('\\\'ec', 'ì')
    text = text.replace('\\\'f2', 'ò')
    text = text.replace('\\\'f9', 'ù')
    text = text.replace('\\\'92', "'")
    text = text.replace('\\\'94', '"')
    text = text.replace('\\\'96', '–')
    
    # Pulisce spazi multipli
    text = re.sub(r'[ \t]+', ' ', text)
    # Pulisce newline multiple ma mantiene struttura
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    
    # Rimuove caratteri di controllo residui
    text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', '', text)
    
    return text.strip()

def parse_course_from_markdown(content, filename):
    """Parsa un file Markdown completo e estrae metadati corso + lezioni"""
    import re
    
    lines = content.split('\n')
    course_data = {
        'code': filename.replace('.md', '').upper(),
        'name': '',
        'description': '',
        'total_hours': 80,
        'theory_hours': 40,
        'practice_hours': 40
    }
    lessons_data = []
    
    i = 0
    current_module = None
    current_section = None
    module_order = 0
    
    while i < len(lines):
        line = lines[i].strip()
        
        # Estrai metadati corso
        # Pulisci la riga dai backslash residui RTF
        clean_line = line.replace('\\', '').strip()
        
        # Titolo corso (può essere nella prima riga o in una sezione)
        if i < 5 and ('Corso di' in clean_line or 'Social Media Marketing' in clean_line):
            # Estrai titolo dalla prima riga con #
            if clean_line.startswith('#'):
                title = clean_line.lstrip('#').strip()
                if title and not course_data['name']:
                    course_data['name'] = title
        
        if 'Titolo corso:' in clean_line or '**Titolo corso:**' in clean_line:
            match = re.search(r'[Tt]itolo corso[:\s]+(.+)', clean_line)
            if match:
                course_data['name'] = match.group(1).strip()
        
        # Durata (pattern: "80 ore (40 teoria, 40 laboratorio)")
        if 'Durata:' in clean_line and 'ore' in clean_line:
            # Pattern completo: "80 ore (40 teoria, 40 laboratorio)"
            match = re.search(r'(\d+)\s*ore.*?\((\d+).*?teoria.*?(\d+).*?laboratorio', clean_line, re.IGNORECASE)
            if match:
                course_data['total_hours'] = int(match.group(1))
                course_data['theory_hours'] = int(match.group(2))
                course_data['practice_hours'] = int(match.group(3))
            else:
                # Pattern alternativo: "80 ore (40 teoria, 40 pratica)"
                match = re.search(r'(\d+)\s*ore.*?\((\d+).*?teoria.*?(\d+).*?pratica', clean_line, re.IGNORECASE)
                if match:
                    course_data['total_hours'] = int(match.group(1))
                    course_data['theory_hours'] = int(match.group(2))
                    course_data['practice_hours'] = int(match.group(3))
                else:
                    # Solo totale
                    match = re.search(r'(\d+)\s*ore', clean_line)
                    if match:
                        course_data['total_hours'] = int(match.group(1))
        
        if 'Durata totale:' in clean_line or '**Durata totale:**' in clean_line:
            match = re.search(r'(\d+)\s*ore', clean_line)
            if match:
                course_data['total_hours'] = int(match.group(1))
        
        if 'Teoria:' in clean_line and 'ore' in clean_line:
            match = re.search(r'(\d+)\s*ore', clean_line)
            if match:
                course_data['theory_hours'] = int(match.group(1))
        
        if ('Laboratorio' in clean_line or 'pratica' in clean_line.lower()) and 'ore' in clean_line:
            match = re.search(r'(\d+)\s*ore', clean_line)
            if match:
                course_data['practice_hours'] = int(match.group(1))
        
        # Cerca tabella moduli (sezione 3)
        if 'Struttura del corso' in line or 'moduli e ore' in line.lower():
            # Cerca la tabella nei prossimi 20 righe
            for j in range(i, min(i+20, len(lines))):
                table_line = lines[j].strip()
                # Pattern tabella: | 1 | Titolo | 12 | 8 | 4 |
                match = re.match(r'\|\s*(\d+)\s*\|\s*(.+?)\s*\|\s*(\d+)\s*\|\s*(\d+)\s*\|\s*(\d+)\s*\|', table_line)
                if match:
                    mod_num = int(match.group(1))
                    mod_title = match.group(2).strip().rstrip('\\').strip()
                    mod_total = int(match.group(3))
                    mod_theory = int(match.group(4))
                    mod_practice = int(match.group(5))
                    
                    # Cerca modulo corrispondente nei dettagli
                    module_found = False
                    for existing_mod in lessons_data:
                        if existing_mod.get('order') == mod_num:
                            existing_mod['duration_total'] = mod_total
                            existing_mod['duration_theory'] = mod_theory
                            existing_mod['duration_practice'] = mod_practice
                            module_found = True
                            break
                    
                    if not module_found:
                        lessons_data.append({
                            'title': mod_title,
                            'order': mod_num,
                            'duration_total': mod_total,
                            'duration_theory': mod_theory,
                            'duration_practice': mod_practice,
                            'objectives': [],
                            'content_theory': '',
                            'content_practice': '',
                            'activities': []
                        })
        
        # Rileva moduli (### Modulo X)
        if re.match(r'^###\s*Modulo\s+\d+', line, re.IGNORECASE):
            # Estrai numero modulo
            mod_match = re.search(r'Modulo\s+(\d+)', line, re.IGNORECASE)
            mod_num = int(mod_match.group(1)) if mod_match else module_order + 1
            
            # Cerca modulo esistente o creane uno nuovo
            current_module = None
            for mod in lessons_data:
                if mod.get('order') == mod_num:
                    current_module = mod
                    break
            
            if not current_module:
                module_title = re.sub(r'^###\s*Modulo\s+\d+\s*[–-]\s*', '', line, flags=re.IGNORECASE).strip().rstrip('\\').strip()
                current_module = {
                    'title': module_title,
                    'order': mod_num,
                    'duration_total': 0,
                    'duration_theory': 0,
                    'duration_practice': 0,
                    'objectives': [],
                    'content_theory': '',
                    'content_practice': '',
                    'activities': []
                }
                lessons_data.append(current_module)
            
            current_section = None
            i += 1
            continue
        
        # Durata modulo
        if current_module and ('Durata:' in line or '**Durata:**' in line):
            # Cerca pattern tipo "12 ore (8 teoria, 4 laboratorio)"
            match = re.search(r'(\d+)\s*ore.*?(\d+).*?teoria.*?(\d+).*?laboratorio', line, re.IGNORECASE)
            if match:
                current_module['duration_total'] = int(match.group(1))
                current_module['duration_theory'] = int(match.group(2))
                current_module['duration_practice'] = int(match.group(3))
            else:
                match = re.search(r'(\d+)\s*ore', line)
                if match:
                    current_module['duration_total'] = int(match.group(1))
            i += 1
            continue
        
        # Sezioni del modulo
        if current_module and line.startswith('####'):
            section_name = line.lstrip('#').strip().lower()
            if 'obiettivi' in section_name:
                current_section = 'objectives'
            elif 'contenuti teorici' in section_name or 'contenuto teorico' in section_name:
                current_section = 'content_theory'
            elif 'attività' in section_name or 'laboratorio' in section_name:
                current_section = 'content_practice'
            elif 'risultati' in section_name:
                current_section = 'results'
            i += 1
            continue
        
        # Processa contenuto
        if current_module and line:
            if current_section == 'objectives':
                if line.startswith('-') or line.startswith('*') or re.match(r'^\d+\.', line):
                    obj = re.sub(r'^[-*\d+\.]\s*', '', line).strip().rstrip('\\').strip()
                    if obj:
                        current_module['objectives'].append(obj)
            elif current_section == 'content_theory':
                current_module['content_theory'] += line + '\n'
            elif current_section == 'content_practice':
                current_module['content_practice'] += line + '\n'
        
        i += 1
    
    # Ordina moduli per ordine
    lessons_data.sort(key=lambda x: x.get('order', 999))
    
    # Crea lezioni dai moduli
    final_lessons = []
    lesson_order = 0
    
    for module in lessons_data:
        # Lezione teorica
        if module.get('duration_theory', 0) > 0:
            lesson_order += 1
            final_lessons.append({
                'title': f"{module['title']} - Parte Teorica",
                'description': f"Parte teorica del modulo {module['order']}",
                'lesson_type': 'theory',
                'duration_hours': float(module.get('duration_theory', 0)),
                'order': lesson_order,
                'content': module.get('content_theory', '').strip(),
                'objectives': module.get('objectives', []),
                'materials': [],
                'exercises': []
            })
        
        # Lezione pratica
        if module.get('duration_practice', 0) > 0:
            lesson_order += 1
            final_lessons.append({
                'title': f"{module['title']} - Parte Pratica",
                'description': f"Parte pratica del modulo {module['order']}",
                'lesson_type': 'practice',
                'duration_hours': float(module.get('duration_practice', 0)),
                'order': lesson_order,
                'content': module.get('content_practice', '').strip(),
                'objectives': [],
                'materials': [],
                'exercises': module.get('activities', [])
            })
    
    # Se non ci sono moduli strutturati, usa il parser semplice
    if not final_lessons:
        # Prova a parsare come lezioni semplici
        temp_course = type('Course', (), {'id': 1})()
        final_lessons = parse_markdown_lessons(content, temp_course)
    
    return course_data, final_lessons

def parse_markdown_lessons(content, course):
    """Parsa un file Markdown e estrae le lezioni"""
    import re
    lessons = []
    lines = content.split('\n')
    
    current_lesson = None
    current_section = None
    i = 0
    
    while i < len(lines):
        line = lines[i]
        line_stripped = line.strip()
        
        # Rileva intestazione lezione con pattern "# Lezione N – Titolo" o "# Lezione N: Titolo"
        lesson_match = re.match(r'^#\s*Lezione\s+(\d+)[\s–:\-]+\s*(.+)$', line_stripped, re.IGNORECASE)
        if lesson_match:
            # Salva lezione precedente
            if current_lesson:
                lessons.append(current_lesson)
            
            # Nuova lezione
            lesson_num = int(lesson_match.group(1))
            title = lesson_match.group(2).strip()
            current_lesson = {
                'title': title,
                'description': '',
                'lesson_type': 'theory',  # default
                'duration_hours': 2.0,  # default
                'order': lesson_num,
                'content': '',
                'objectives': [],
                'materials': [],
                'exercises': []
            }
            current_section = 'content'  # Di default, tutto il contenuto va nel campo content
            i += 1
            continue
        
        # Rileva anche intestazioni ## come titoli di lezione (fallback)
        if line_stripped.startswith('##') and not line_stripped.startswith('###'):
            # Salva lezione precedente solo se non è già stata salvata
            if current_lesson and current_lesson['order'] == len(lessons):
                # Non salvare ancora, potrebbe essere una sezione della lezione
                pass
            else:
                if current_lesson:
                    lessons.append(current_lesson)
            
            # Controlla se è una nuova lezione o una sezione
            title = line_stripped.lstrip('#').strip()
            # Se non c'è una lezione corrente o se il titolo sembra essere una nuova lezione
            if not current_lesson or (current_lesson and 'Lezione' in title):
                current_lesson = {
                    'title': title,
                    'description': '',
                    'lesson_type': 'theory',
                    'duration_hours': 2.0,
                    'order': len(lessons) + 1,
                    'content': '',
                    'objectives': [],
                    'materials': [],
                    'exercises': []
                }
                current_section = 'content'
            else:
                # È una sezione della lezione corrente, aggiungila al contenuto
                if current_lesson and current_section == 'content':
                    current_lesson['content'] += line + '\n'
            i += 1
            continue
        
        # Rileva sezioni (###)
        if line_stripped.startswith('###'):
            section_name = line_stripped.lstrip('#').strip().lower()
            if 'obiettivi' in section_name or 'obiettivo' in section_name:
                current_section = 'objectives'
            elif 'materiali' in section_name or 'risorse' in section_name:
                current_section = 'materials'
            elif 'esercizi' in section_name or 'esercizio' in section_name:
                current_section = 'exercises'
            elif 'contenuti' in section_name or 'contenuto' in section_name:
                current_section = 'content'
            elif 'descrizione' in section_name:
                current_section = 'description'
            elif 'tipo' in section_name:
                current_section = 'type'
            elif 'durata' in section_name:
                current_section = 'duration'
            else:
                # Sezione sconosciuta, aggiungila al contenuto
                current_section = 'content'
                if current_lesson:
                    current_lesson['content'] += line + '\n'
            i += 1
            continue
        
        # Processa contenuto in base alla sezione
        if current_lesson:
            if current_section == 'objectives':
                if line_stripped.startswith('-') or line_stripped.startswith('*'):
                    current_lesson['objectives'].append(line_stripped.lstrip('-* ').strip())
            elif current_section == 'materials':
                if line_stripped.startswith('-') or line_stripped.startswith('*'):
                    current_lesson['materials'].append(line_stripped.lstrip('-* ').strip())
            elif current_section == 'exercises':
                if line_stripped.startswith('-') or line_stripped.startswith('*'):
                    current_lesson['exercises'].append(line_stripped.lstrip('-* ').strip())
            elif current_section == 'content':
                # Aggiungi tutto al contenuto, inclusi i titoli delle sezioni ##
                current_lesson['content'] += line + '\n'
            elif current_section == 'description':
                current_lesson['description'] += line_stripped + ' '
            elif current_section == 'type':
                if 'pratic' in line_stripped.lower():
                    current_lesson['lesson_type'] = 'practice'
                elif 'teoric' in line_stripped.lower():
                    current_lesson['lesson_type'] = 'theory'
            elif current_section == 'duration':
                # Estrai numero da line (es: "2 ore" -> 2.0)
                match = re.search(r'(\d+\.?\d*)', line_stripped)
                if match:
                    current_lesson['duration_hours'] = float(match.group(1))
            elif not current_section:
                # Se non c'è una sezione specifica, aggiungi al contenuto
                current_lesson['content'] += line + '\n'
        
        i += 1
    
    # Aggiungi ultima lezione
    if current_lesson:
        lessons.append(current_lesson)
    
    # Pulisci descrizioni e contenuti
    for lesson in lessons:
        lesson['description'] = lesson['description'].strip()
        lesson['content'] = lesson['content'].strip()
        # Se non c'è contenuto ma c'è una descrizione, usa la descrizione come contenuto
        if not lesson['content'] and lesson['description']:
            lesson['content'] = lesson['description']
    
    return lessons

@app.route('/api/courses/<int:course_id>/generate', methods=['POST'])
def generate_course_files(course_id):
    course = Course.query.get_or_404(course_id)
    lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
    
    course_dir = os.path.join(app.config['COURSES_DIR'], course.code)
    os.makedirs(course_dir, exist_ok=True)
    
    # Genera README principale del corso
    readme_content = generate_course_readme(course, lessons)
    readme_path = os.path.join(course_dir, 'README.md')
    with open(readme_path, 'w', encoding='utf-8') as f:
        f.write(readme_content)
    
    # Genera file per ogni lezione
    for lesson in lessons:
        lesson_content = generate_lesson_markdown(course, lesson)
        lesson_filename = f"{lesson.order:02d}_{lesson.title.lower().replace(' ', '_')}.md"
        lesson_path = os.path.join(course_dir, lesson_filename)
        with open(lesson_path, 'w', encoding='utf-8') as f:
            f.write(lesson_content)
    
    return jsonify({
        'message': f'File generati con successo in {course_dir}',
        'path': course_dir
    })

def generate_course_readme(course, lessons):
    theory_lessons = [l for l in lessons if l.lesson_type == 'theory']
    practice_lessons = [l for l in lessons if l.lesson_type == 'practice']
    
    content = f"""# {course.name}

**Codice Corso:** {course.code}  
**Durata Totale:** {course.total_hours} ore  
**Ore Teoriche:** {course.theory_hours} ore  
**Ore Pratiche:** {course.practice_hours} ore

## Descrizione

{course.description or 'Descrizione del corso da completare.'}

## Programma del Corso

### Parte Teorica ({course.theory_hours} ore)

"""
    for i, lesson in enumerate(theory_lessons, 1):
        content += f"{i}. **{lesson.title}** ({lesson.duration_hours}h) - {lesson.description or ''}\n"
    
    content += f"\n### Parte Pratica ({course.practice_hours} ore)\n\n"
    for i, lesson in enumerate(practice_lessons, 1):
        content += f"{i}. **{lesson.title}** ({lesson.duration_hours}h) - {lesson.description or ''}\n"
    
    content += "\n## Struttura Lezioni\n\n"
    for lesson in lessons:
        content += f"- [{lesson.order:02d}. {lesson.title}](./{lesson.order:02d}_{lesson.title.lower().replace(' ', '_')}.md)\n"
    
    return content

def generate_lesson_markdown(course, lesson):
    objectives = json.loads(lesson.objectives) if lesson.objectives else []
    materials = json.loads(lesson.materials) if lesson.materials else []
    exercises = json.loads(lesson.exercises) if lesson.exercises else []
    
    content = f"""# {lesson.title}

**Corso:** {course.name} ({course.code})  
**Tipo:** {'Teorica' if lesson.lesson_type == 'theory' else 'Pratica'}  
**Durata:** {lesson.duration_hours} ore  
**Ordine:** {lesson.order}

## Descrizione

{lesson.description or 'Descrizione della lezione da completare.'}

## Obiettivi di Apprendimento

"""
    for obj in objectives:
        content += f"- {obj}\n"
    
    if not objectives:
        content += "- Obiettivi da definire\n"
    
    content += "\n## Contenuti\n\n"
    content += lesson.content or "Contenuti della lezione da completare.\n"
    
    if materials:
        content += "\n## Materiali e Risorse\n\n"
        for material in materials:
            content += f"- {material}\n"
    
    if exercises:
        content += "\n## Esercizi Pratici\n\n"
        for i, exercise in enumerate(exercises, 1):
            content += f"### Esercizio {i}\n\n{exercise}\n\n"
    
    content += f"\n---\n\n*Generato il {datetime.now().strftime('%d/%m/%Y %H:%M')}*\n"
    
    return content

def get_preference_value(key, default=None):
    """Helper per recuperare una preferenza con default"""
    pref = Preference.query.filter_by(key=key).first()
    if pref:
        try:
            return json.loads(pref.value)
        except (json.JSONDecodeError, TypeError):
            return pref.value
    return default

@app.route('/api/preferences', methods=['GET'])
def get_preferences():
    """Ottiene tutte le preferenze"""
    try:
        # Prova a caricare le preferenze - se fallisce, inizializza il database
        try:
            preferences = Preference.query.all()
        except Exception as db_error:
            # Se il database non esiste, inizializzalo
            error_str = str(db_error)
            if 'no such table' in error_str.lower() or 'operationalerror' in error_str.lower():
                logger.warning(f"Database non inizializzato in get_preferences, inizializzazione in corso... Errore: {db_error}")
                try:
                    init_database()
                    logger.info("Database inizializzato con successo in get_preferences")
                    # Riprova la query dopo l'inizializzazione
                    preferences = Preference.query.all()
                except Exception as init_error:
                    logger.error(f"Errore durante inizializzazione database in get_preferences: {init_error}")
                    # Se l'inizializzazione fallisce, ritorna comunque i default
                    # così l'app può funzionare anche senza database
                    defaults = {
                        'defaultLessonDuration': 2.0,
                        'defaultLessonType': 'theory',
                        'defaultLessonTime': '09:00',
                        'exportAuthor': '',
                        'exportCompany': '',
                        'hourlyRate': 25.0,
                        'includeFooter': True,
                        'aiModel': 'gpt-4o',
                        'aiTemperature': 0.7,
                        'aiMaxTokens': 2000,
                        'autoExpandLessons': True,
                        'questionsCount': 30,
                        'questionsOptions': 4,
                        'slidesPerLesson': 10,
                        'includeImages': True,
                        'coursesDir': app.config['COURSES_DIR'],
                        'mdSourceDir': app.config['MD_SOURCE_DIR'],
                        'enableDebugLogs': False
                    }
                    return jsonify(defaults), 200
            else:
                # Altro tipo di errore, rilancia
                raise
        prefs_dict = {}
        for pref in preferences:
            # Prova a parsare come JSON, altrimenti usa il valore diretto
            try:
                prefs_dict[pref.key] = json.loads(pref.value)
            except (json.JSONDecodeError, TypeError):
                prefs_dict[pref.key] = pref.value
        
        # Aggiungi valori di default se non esistono
        defaults = {
            'defaultLessonDuration': 2.0,
            'defaultLessonType': 'theory',
            'defaultLessonTime': '09:00',
            'exportAuthor': '',
            'exportCompany': '',
            'hourlyRate': 25.0,
            'includeFooter': True,
            'aiModel': 'gpt-4o',
            'aiTemperature': 0.7,
            'aiMaxTokens': 2000,
            'autoExpandLessons': True,
            'questionsCount': 30,
            'questionsOptions': 4,
            'slidesPerLesson': 10,
            'includeImages': True,
            'coursesDir': app.config['COURSES_DIR'],
            'mdSourceDir': app.config['MD_SOURCE_DIR'],
            'enableDebugLogs': False
        }
        
        # Unisci con i default (le preferenze salvate hanno priorità)
        result = {**defaults, **prefs_dict}
        
        return jsonify(result), 200
    except Exception as e:
        return jsonify({'error': f'Errore nel recupero delle preferenze: {str(e)}'}), 500

@app.route('/api/preferences', methods=['POST'])
def save_preferences():
    """Salva le preferenze"""
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'Nessun dato fornito'}), 400
        
        # Prova a salvare - se fallisce per database non esistente, inizializzalo
        try:
            # Prova una query per vedere se il database funziona
            Preference.query.limit(1).all()
        except Exception as db_error:
            error_str = str(db_error)
            if 'no such table' in error_str.lower() or 'operationalerror' in error_str.lower():
                logger.warning(f"Database non inizializzato in save_preferences, inizializzazione in corso... Errore: {db_error}")
                try:
                    init_database()
                    logger.info("Database inizializzato con successo in save_preferences")
                except Exception as init_error:
                    logger.error(f"Errore durante inizializzazione database in save_preferences: {init_error}")
                    return jsonify({'error': f'Database non inizializzato: {str(init_error)}'}), 500
            else:
                raise
        
        # Salva ogni preferenza
        for key, value in data.items():
            # Converti valori complessi in JSON
            if isinstance(value, (dict, list, bool)):
                value_str = json.dumps(value)
            else:
                value_str = str(value)
            
            # Cerca preferenza esistente o crea nuova
            pref = Preference.query.filter_by(key=key).first()
            if pref:
                pref.value = value_str
                pref.updated_at = datetime.utcnow()
            else:
                pref = Preference(key=key, value=value_str)
                db.session.add(pref)
        
        db.session.commit()
        return jsonify({'message': 'Preferenze salvate con successo'}), 200
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'Errore nel salvataggio delle preferenze: {str(e)}'}), 500

@app.route('/api/dashboard/stats', methods=['GET'])
def get_dashboard_stats():
    """Ottiene tutte le statistiche per la dashboard"""
    try:
        from sqlalchemy import func
        from sqlalchemy.orm import joinedload
        
        # Parametri paginazione
        recent_page = int(request.args.get('recent_page', 1))
        recent_per_page = int(request.args.get('recent_per_page', 5))
        upcoming_page = int(request.args.get('upcoming_page', 1))
        upcoming_per_page = int(request.args.get('upcoming_per_page', 5))
        
        # Statistiche generali
        total_courses = Course.query.count()
        total_lessons = Lesson.query.count()
        
        # Calcola ore totali
        total_hours_result = db.session.query(func.sum(Course.total_hours)).scalar() or 0
        total_hours = float(total_hours_result)
        
        # Calcola importo totale
        hourly_rate = float(get_preference_value('hourlyRate', 25.0))
        total_revenue = total_hours * hourly_rate
        
        # Ore per corso (per grafico)
        courses_hours = db.session.query(
            Course.name,
            Course.total_hours
        ).all()
        hours_per_course = [{'name': c[0], 'hours': float(c[1] or 0)} for c in courses_hours]
        
        # Distribuzione teoria/pratica
        total_theory = db.session.query(func.sum(Course.theory_hours)).scalar() or 0
        total_practice = db.session.query(func.sum(Course.practice_hours)).scalar() or 0
        theory_practice_dist = {
            'theory': float(total_theory),
            'practice': float(total_practice)
        }
        
        # Corsi più recenti con paginazione
        total_recent = Course.query.count()
        recent_offset = (recent_page - 1) * recent_per_page
        recent_courses = Course.query.order_by(Course.id.desc()).offset(recent_offset).limit(recent_per_page).all()
        recent_courses_data = [{
            'id': c.id,
            'name': c.name,
            'code': c.code,
            'total_hours': c.total_hours,
            'created_at': c.created_at.isoformat() if c.created_at else None
        } for c in recent_courses]
        recent_total_pages = (total_recent + recent_per_page - 1) // recent_per_page if total_recent > 0 else 1
        
        # Prossime lezioni programmate con paginazione
        today = date.today()
        total_upcoming = db.session.query(Lesson).join(Course).filter(
            Lesson.lesson_date >= today
        ).count()
        upcoming_offset = (upcoming_page - 1) * upcoming_per_page
        upcoming_lessons = db.session.query(Lesson, Course).join(Course).filter(
            Lesson.lesson_date >= today
        ).order_by(Lesson.lesson_date.asc()).offset(upcoming_offset).limit(upcoming_per_page).all()
        
        upcoming_lessons_data = []
        for lesson, course in upcoming_lessons:
            upcoming_lessons_data.append({
                'id': lesson.id,
                'title': lesson.title,
                'course_name': course.name,
                'course_code': course.code,
                'lesson_date': lesson.lesson_date.isoformat() if lesson.lesson_date else None,
                'duration_hours': lesson.duration_hours,
                'lesson_type': lesson.lesson_type
            })
        upcoming_total_pages = (total_upcoming + upcoming_per_page - 1) // upcoming_per_page if total_upcoming > 0 else 1
        
        return jsonify({
            'total_courses': total_courses,
            'total_lessons': total_lessons,
            'total_hours': total_hours,
            'total_revenue': total_revenue,
            'hours_per_course': hours_per_course,
            'theory_practice_dist': theory_practice_dist,
            'recent_courses': recent_courses_data,
            'recent_pagination': {
                'page': recent_page,
                'per_page': recent_per_page,
                'total': total_recent,
                'total_pages': recent_total_pages,
                'has_prev': recent_page > 1,
                'has_next': recent_page < recent_total_pages
            },
            'upcoming_lessons': upcoming_lessons_data,
            'upcoming_pagination': {
                'page': upcoming_page,
                'per_page': upcoming_per_page,
                'total': total_upcoming,
                'total_pages': upcoming_total_pages,
                'has_prev': upcoming_page > 1,
                'has_next': upcoming_page < upcoming_total_pages
            }
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Errore nel recupero delle statistiche: {str(e)}'}), 500

@app.route('/api/dashboard/revenue', methods=['GET'])
def get_revenue_by_period():
    """Ottiene gli importi per periodo"""
    try:
        from sqlalchemy import func
        
        period = request.args.get('period', 'year')
        hourly_rate = float(get_preference_value('hourlyRate', 25.0))
        
        # Calcola la data di inizio in base al periodo
        today = date.today()
        if period == 'month':
            start_date = today - timedelta(days=30)
        elif period == 'quarter':
            start_date = today - timedelta(days=90)
        elif period == 'year':
            start_date = today - timedelta(days=365)
        else:  # all
            start_date = None
        
        # Raggruppa per mese usando i dati dei corsi
        # Per semplicità, usiamo il mese di creazione del corso
        # In una versione più avanzata, potresti raggruppare per mese delle lezioni
        revenue_by_month = {}
        courses = Course.query.all()
        
        for course in courses:
            # Se c'è un filtro per data, verifica che il corso sia nel periodo
            if start_date and course.created_at and course.created_at.date() < start_date:
                continue
            
            if course.created_at:
                month_key = course.created_at.strftime('%Y-%m')
                if month_key not in revenue_by_month:
                    revenue_by_month[month_key] = 0
                revenue_by_month[month_key] += (course.total_hours or 0) * hourly_rate
        
        # Ordina per mese
        sorted_months = sorted(revenue_by_month.items())
        
        return jsonify({
            'period': period,
            'revenue_by_month': [{'month': m[0], 'revenue': m[1]} for m in sorted_months]
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Errore nel recupero degli importi: {str(e)}'}), 500

@app.route('/api/templates', methods=['GET'])
def get_templates():
    """Ottiene tutti i template disponibili"""
    try:
        templates = CourseTemplate.query.order_by(CourseTemplate.created_at.desc()).all()
        result = []
        for t in templates:
            result.append({
                'id': t.id,
                'name': t.name,
                'description': t.description,
                'code_prefix': t.code_prefix,
                'is_predefined': t.is_predefined,
                'created_at': t.created_at.isoformat() if t.created_at else None,
                'updated_at': t.updated_at.isoformat() if t.updated_at else None
            })
        return jsonify(result), 200
    except Exception as e:
        logger.error(f"Errore nel recupero template: {e}")
        return jsonify({'error': f'Errore nel recupero template: {str(e)}'}), 500

@app.route('/api/templates', methods=['POST'])
def create_template():
    """Salva un corso come template"""
    try:
        data = request.json
        course_id = data.get('course_id')
        
        if not course_id:
            return jsonify({'error': 'course_id richiesto'}), 400
        
        # Carica il corso con tutte le lezioni
        course = Course.query.get_or_404(course_id)
        lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
        
        # Prepara i dati del corso
        course_data = {
            'name': course.name,
            'description': course.description,
            'total_hours': course.total_hours,
            'theory_hours': course.theory_hours,
            'practice_hours': course.practice_hours,
            'num_lessons': course.num_lessons or 0
        }
        
        # Prepara i dati delle lezioni
        lessons_data = []
        for lesson in lessons:
            lessons_data.append({
                'title': lesson.title,
                'description': lesson.description,
                'lesson_type': lesson.lesson_type,
                'duration_hours': lesson.duration_hours,
                'order': lesson.order,
                'content': lesson.content,
                'objectives': json.loads(lesson.objectives) if lesson.objectives else [],
                'materials': json.loads(lesson.materials) if lesson.materials else [],
                'exercises': json.loads(lesson.exercises) if lesson.exercises else []
            })
        
        # Crea il template
        template = CourseTemplate(
            name=data.get('template_name', f"Template: {course.name}"),
            description=data.get('template_description', course.description),
            code_prefix=data.get('code_prefix', 'TEMPLATE_'),
            course_data=json.dumps(course_data),
            lessons_data=json.dumps(lessons_data),
            is_predefined=False
        )
        db.session.add(template)
        db.session.commit()
        
        return jsonify({
            'id': template.id,
            'message': 'Template creato con successo'
        }), 201
    except Exception as e:
        db.session.rollback()
        logger.error(f"Errore nella creazione template: {e}")
        return jsonify({'error': f'Errore nella creazione template: {str(e)}'}), 500

@app.route('/api/templates/<int:template_id>', methods=['GET'])
def get_template(template_id):
    """Ottiene i dettagli di un template"""
    try:
        template = CourseTemplate.query.get_or_404(template_id)
        return jsonify({
            'id': template.id,
            'name': template.name,
            'description': template.description,
            'code_prefix': template.code_prefix,
            'is_predefined': template.is_predefined,
            'course_data': json.loads(template.course_data),
            'lessons_data': json.loads(template.lessons_data) if template.lessons_data else [],
            'created_at': template.created_at.isoformat() if template.created_at else None
        }), 200
    except Exception as e:
        logger.error(f"Errore nel recupero template: {e}")
        return jsonify({'error': f'Errore nel recupero template: {str(e)}'}), 500

@app.route('/api/templates/<int:template_id>', methods=['DELETE'])
def delete_template(template_id):
    """Elimina un template"""
    try:
        template = CourseTemplate.query.get_or_404(template_id)
        if template.is_predefined:
            return jsonify({'error': 'Non è possibile eliminare template predefiniti'}), 403
        db.session.delete(template)
        db.session.commit()
        return jsonify({'message': 'Template eliminato con successo'}), 200
    except Exception as e:
        db.session.rollback()
        logger.error(f"Errore nell'eliminazione template: {e}")
        return jsonify({'error': f'Errore nell\'eliminazione template: {str(e)}'}), 500

@app.route('/api/courses/<int:course_id>/duplicate', methods=['POST'])
def duplicate_course(course_id):
    """Duplica un corso con tutte le lezioni"""
    try:
        data = request.json
        original_course = Course.query.get_or_404(course_id)
        original_lessons = Lesson.query.filter_by(course_id=course_id).order_by(Lesson.order).all()
        
        # Genera un nuovo codice univoco
        base_code = data.get('new_code', f"{original_course.code}_COPY")
        new_code = base_code
        counter = 1
        while Course.query.filter_by(code=new_code).first():
            new_code = f"{base_code}_{counter}"
            counter += 1
        
        # Crea il nuovo corso
        new_course = Course(
            code=new_code,
            name=data.get('new_name', f"{original_course.name} (Copia)"),
            description=original_course.description,
            total_hours=original_course.total_hours,
            theory_hours=original_course.theory_hours,
            practice_hours=original_course.practice_hours,
            num_lessons=original_course.num_lessons
        )
        db.session.add(new_course)
        db.session.flush()
        
        # Duplica le lezioni
        for lesson in original_lessons:
            new_lesson = Lesson(
                course_id=new_course.id,
                title=lesson.title,
                description=lesson.description,
                lesson_type=lesson.lesson_type,
                duration_hours=lesson.duration_hours,
                order=lesson.order,
                content=lesson.content,
                objectives=lesson.objectives,
                materials=lesson.materials,
                exercises=lesson.exercises
            )
            db.session.add(new_lesson)
        
        db.session.commit()
        return jsonify({
            'id': new_course.id,
            'code': new_course.code,
            'message': 'Corso duplicato con successo'
        }), 201
    except Exception as e:
        db.session.rollback()
        logger.error(f"Errore nella duplicazione corso: {e}")
        return jsonify({'error': f'Errore nella duplicazione corso: {str(e)}'}), 500

@app.route('/api/courses/from-template', methods=['POST'])
def create_course_from_template():
    """Crea un nuovo corso da un template"""
    try:
        data = request.json
        template_id = data.get('template_id')
        
        if not template_id:
            return jsonify({'error': 'template_id richiesto'}), 400
        
        template = CourseTemplate.query.get_or_404(template_id)
        course_data = json.loads(template.course_data)
        lessons_data = json.loads(template.lessons_data) if template.lessons_data else []
        
        # Genera un nuovo codice
        code_prefix = template.code_prefix or 'TEMPLATE_'
        base_code = data.get('code', f"{code_prefix}{course_data['name'][:20].upper().replace(' ', '_')}")
        new_code = base_code
        counter = 1
        while Course.query.filter_by(code=new_code).first():
            new_code = f"{base_code}_{counter}"
            counter += 1
        
        # Crea il nuovo corso
        new_course = Course(
            code=new_code,
            name=data.get('name', course_data['name']),
            description=data.get('description', course_data.get('description', '')),
            total_hours=course_data['total_hours'],
            theory_hours=course_data['theory_hours'],
            practice_hours=course_data['practice_hours'],
            num_lessons=course_data.get('num_lessons', 0)
        )
        db.session.add(new_course)
        db.session.flush()
        
        # Crea le lezioni
        for lesson_data in lessons_data:
            new_lesson = Lesson(
                course_id=new_course.id,
                title=lesson_data['title'],
                description=lesson_data.get('description', ''),
                lesson_type=lesson_data.get('lesson_type', 'theory'),
                duration_hours=lesson_data.get('duration_hours', 2.0),
                order=lesson_data.get('order', len(lessons_data)),
                content=lesson_data.get('content', ''),
                objectives=json.dumps(lesson_data.get('objectives', [])),
                materials=json.dumps(lesson_data.get('materials', [])),
                exercises=json.dumps(lesson_data.get('exercises', []))
            )
            db.session.add(new_lesson)
        
        db.session.commit()
        return jsonify({
            'id': new_course.id,
            'code': new_course.code,
            'message': 'Corso creato dal template con successo'
        }), 201
    except Exception as e:
        db.session.rollback()
        logger.error(f"Errore nella creazione corso da template: {e}")
        return jsonify({'error': f'Errore nella creazione corso da template: {str(e)}'}), 500

@app.route('/api/templates/<int:template_id>/export', methods=['GET'])
def export_template(template_id):
    """Esporta un template come JSON"""
    try:
        template = CourseTemplate.query.get_or_404(template_id)
        export_data = {
            'name': template.name,
            'description': template.description,
            'code_prefix': template.code_prefix,
            'course_data': json.loads(template.course_data),
            'lessons_data': json.loads(template.lessons_data) if template.lessons_data else [],
            'exported_at': datetime.utcnow().isoformat()
        }
        return jsonify(export_data), 200
    except Exception as e:
        logger.error(f"Errore nell'esportazione template: {e}")
        return jsonify({'error': f'Errore nell\'esportazione template: {str(e)}'}), 500

@app.route('/api/templates/import', methods=['POST'])
def import_template():
    """Importa un template da JSON"""
    try:
        data = request.json
        
        if not data.get('course_data'):
            return jsonify({'error': 'Dati template non validi'}), 400
        
        template = CourseTemplate(
            name=data.get('name', 'Template Importato'),
            description=data.get('description', ''),
            code_prefix=data.get('code_prefix', 'IMPORTED_'),
            course_data=json.dumps(data['course_data']),
            lessons_data=json.dumps(data.get('lessons_data', [])),
            is_predefined=False
        )
        db.session.add(template)
        db.session.commit()
        
        return jsonify({
            'id': template.id,
            'message': 'Template importato con successo'
        }), 201
    except Exception as e:
        db.session.rollback()
        logger.error(f"Errore nell'importazione template: {e}")
        return jsonify({'error': f'Errore nell\'importazione template: {str(e)}'}), 500

@app.route('/api/dashboard/usage-stats', methods=['GET'])
def get_usage_stats():
    """Ottiene le statistiche di utilizzo (lezioni più modificate, corsi più esportati)"""
    try:
        # Parametri paginazione per lezioni modificate
        page = int(request.args.get('page', 1))
        per_page = int(request.args.get('per_page', 5))
        
        # Conta totale lezioni modificate
        total_modified = db.session.query(Lesson).filter(
            Lesson.modification_count > 0
        ).count()
        
        # Calcola offset
        offset = (page - 1) * per_page
        
        # Lezioni più modificate con paginazione
        most_modified_lessons = db.session.query(Lesson, Course).join(Course).filter(
            Lesson.modification_count > 0
        ).order_by(
            Lesson.modification_count.desc(),
            Lesson.updated_at.desc()
        ).offset(offset).limit(per_page).all()
        
        most_modified_data = []
        for lesson, course in most_modified_lessons:
            most_modified_data.append({
                'id': lesson.id,
                'title': lesson.title,
                'course_name': course.name,
                'course_code': course.code,
                'modification_count': lesson.modification_count or 0,
                'last_modified': lesson.updated_at.isoformat() if lesson.updated_at else None
            })
        
        # Parametri paginazione per corsi esportati
        exported_page = int(request.args.get('exported_page', 1))
        exported_per_page = int(request.args.get('exported_per_page', 5))
        
        # Conta totale corsi esportati
        total_exported = Course.query.filter(
            Course.export_count > 0
        ).count()
        
        # Calcola offset
        exported_offset = (exported_page - 1) * exported_per_page
        
        # Corsi più esportati con paginazione
        most_exported_courses = Course.query.filter(
            Course.export_count > 0
        ).order_by(
            Course.export_count.desc()
        ).offset(exported_offset).limit(exported_per_page).all()
        
        most_exported_data = [{
            'id': c.id,
            'name': c.name,
            'code': c.code,
            'export_count': c.export_count or 0,
            'total_hours': c.total_hours
        } for c in most_exported_courses]
        
        # Calcola totale pagine per corsi esportati
        exported_total_pages = (total_exported + exported_per_page - 1) // exported_per_page if total_exported > 0 else 1
        
        # Calcola totale pagine
        total_pages = (total_modified + per_page - 1) // per_page if total_modified > 0 else 1
        
        return jsonify({
            'most_modified_lessons': most_modified_data,
            'modified_pagination': {
                'page': page,
                'per_page': per_page,
                'total': total_modified,
                'total_pages': total_pages,
                'has_prev': page > 1,
                'has_next': page < total_pages
            },
            'most_exported_courses': most_exported_data,
            'exported_pagination': {
                'page': exported_page,
                'per_page': exported_per_page,
                'total': total_exported,
                'total_pages': exported_total_pages,
                'has_prev': exported_page > 1,
                'has_next': exported_page < exported_total_pages
            }
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Errore nel recupero delle statistiche di utilizzo: {str(e)}'}), 500

if __name__ == '__main__':
    # Per sviluppo locale, il database è già inizializzato da init_database()
    app.run(host='0.0.0.0', port=5000, debug=True, use_reloader=False, use_debugger=True)

